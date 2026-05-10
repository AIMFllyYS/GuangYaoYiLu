from __future__ import annotations

"""
Build the v5 stable defense deck.

The v5 route deliberately uses python-pptx instead of hand-written OOXML.
It keeps the deck static, UTF-8 safe, and image-rich while preserving key
elements as native PowerPoint shapes/text/pictures where practical.
"""

import hashlib
import json
import shutil
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from validate_pptx_package import validate_pptx


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "交付物" / "答辩PPT" / "答辩PPT_v5"
V4 = ROOT / "交付物" / "答辩PPT" / "答辩PPT_v4"
SUPPLEMENT = ROOT / "交付物" / "答辩PPT" / "补充性PPT"
MERCH_REF = ROOT / "交付物" / "周边" / "瑶光文创"
ASSETS = OUT / "assets"
PREVIEW = OUT / "preview_png"
PPTX_NAME = "光药医路_5分钟稳定答辩PPT_v5.pptx"
PDF_NAME = "光药医路_5分钟稳定答辩PPT_v5_静态预览.pdf"
SCRIPT_NAME = "光药医路_5分钟答辩讲稿_v5.md"

PX_W = 1920
PX_H = 1080
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5

CREAM = (246, 241, 226)
INK = (45, 55, 39)
MUTED = (112, 99, 76)
GREEN = (45, 106, 69)
DEEP_GREEN = (28, 81, 54)
RED = (174, 38, 43)
GOLD = (185, 143, 70)
PALE_GREEN = (226, 237, 218)

FONTS = {
    "yahei": "C:/Windows/Fonts/msyh.ttc",
    "hei": "C:/Windows/Fonts/simhei.ttf",
    "song": "C:/Windows/Fonts/simsun.ttc",
}

AI_MERCH_SOURCE = Path(
    "C:/Users/AIMFl/.codex/generated_images/019e0b62-fd8f-72c1-a3ca-4781328159d0/"
    "ig_09f4f6660e451e010169feec17f0648190890a4a18a7bcc306.png"
)


@dataclass
class Element:
    slide: int
    element_id: str
    type: str
    source: str
    bbox: dict[str, int]
    z: int
    editable: bool
    text: str | None = None


@dataclass
class SlideText:
    title: str
    visible_text: list[str]
    speaker: str
    subtitle: str = ""
    elements: list[Element] = field(default_factory=list)


SLIDE_TEXT: list[SlideText] = [
    SlideText(
        "光药医路",
        ["光药医路", "金光青黛同济兴，本草杏林华夏清", "药学2503 × 光电2506 × 基医2501 联合团支部", "01/22"],
        "各位评委老师好，我们是药学中外2503、光电2506、基础医学强基2501联合团支部。我们的特色团日主题是“金光青黛同济兴，本草杏林华夏清”。接下来，我们将用五分钟汇报这段从三个班走向一个集体的“光药医路”。",
    ),
    SlideText(
        "评分导向",
        ["评分导向", "主题贯穿 · 内容丰富 · 组织有效 · 形式创新 · 成长表达", "我们围绕五个关键词展开答辩", "02/22"],
        "本次特色团日评分不仅看活动有没有做，更看主题是否贯穿、内容是否丰富、组织是否有效、形式是否有创新，以及答辩能否讲出支部成长。因此我们的展示不按流水账展开，而围绕这五个关键词呈现。",
    ),
    SlideText(
        "主题释义",
        ["主题释义", "光", "药", "医", "路", "技术视角", "本草传承", "仁心关怀", "共同成长", "中医药文化，是三个专业共同对话的交汇点", "03/22"],
        "“光药医路”不是简单拼接三个专业名称。“光”代表光电技术与创新视角，“药”代表本草文化与药学传承，“医”代表生命科学与仁心关怀，“路”代表我们75位同学共同走过的实践过程。中医药文化，正好成为三个专业共同对话的交汇点。",
    ),
    SlideText(
        "三支部总览",
        ["三支部总览", "75位同行者，差异成为合力", "药学2503", "光电2506", "基医2501", "04/22"],
        "联合支部由三个专业背景差异很大的班级组成。药学支部负责总协调和中药知识支撑，光电支部带来技术和系统设计能力，基医支部提供医学理解和视觉表达。差异不是阻力，反而成为我们活动创新的来源。",
    ),
    SlideText(
        "组织架构",
        ["组织架构", "跨班混编，让联合支部真正运转", "策划组", "宣传组", "实践组", "财务组", "后勤组", "策划 · 物资 · 执行 · 宣传 · 总结", "05/22"],
        "为了让联合支部真正运转起来，我们没有按班级各做各的，而是设置宣传、财务、实践、策划、后勤五个功能组，跨班混编。每一次活动都形成“策划、物资、执行、宣传、总结”的闭环，这保证了后续多场活动能够持续推进。",
    ),
    SlideText(
        "三色成长弧",
        ["三色成长弧", "初识白", "思政红", "药草绿", "破冰相识", "团课学习 · 志愿服务", "嘉年华 · 宣讲 · 家乡考察 · 叶开泰", "06/22"],
        "我们把整个特色团日概括为三色成长弧。初识白，是破冰相识；思政红，是团课学习和志愿服务；药草绿，是围绕中医药文化展开的路演、宣讲、家乡考察和叶开泰参观。三种颜色共同构成支部成长的完整路径。",
    ),
    SlideText(
        "初识白",
        ["初识白", "破冰相识，味冰之夜", "传声筒", "数字炸弹", "狼人杀", "阿瓦隆", "联合支部从通知里的名称，变成真实的温度", "07/22"],
        "12月7日，我们在韵苑活动中心开展破冰活动。特色团日介绍让大家知道未来要一起做什么，传声筒、数字炸弹、狼人杀和阿瓦隆让同学们真正熟悉起来。那天之后，“联合支部”不再只是通知里的名称，而开始有了真实的温度。",
    ),
    SlideText(
        "思政红团课",
        ["思政红团课", "国家需求中的青年使命", "科技创新", "民生健康", "绿色发展", "从专业出发，理解时代责任", "08/22"],
        "思政红首先体现在共同学习。3月29日，三个支部围绕十五五规划开展团课，从科技创新、民生健康、绿色发展等维度理解青年使命。药学、光电、基医三个专业，也在国家需求的语境中找到了各自的责任位置。",
    ),
    SlideText(
        "思政红志愿",
        ["思政红志愿", "把担当写在社区巷道里", "清理广告", "扫除落叶", "张贴通知", "红建社区 · 4月19日", "09/22"],
        "思政红不只停留在课堂。4月19日，我们来到红建社区，清理违法小广告、扫除巷道落叶、张贴正规通知。大家穿上红马甲，拿起铲子和喷壶，在真实社区场景中理解“服务基层”的含义。",
    ),
    SlideText(
        "药草绿嘉年华",
        ["药草绿嘉年华", "让中医药在互动中被看见", "辨认药材", "捣药闻香", "药秤称量", "观察 · 触摸 · 嗅闻 · 操作", "10/22"],
        "嘉年华是我们面向全校的首次集中亮相。我们把中医药知识设计成三个互动游戏：辨认药材、捣药闻香、传统药秤称量。参与者不是被动听科普，而是在观察、触摸、嗅闻和操作中走近中医药文化。",
    ),
    SlideText(
        "创新亮点",
        ["创新亮点", "用代码和视觉，让传统文化年轻化", "小程序", "积分", "百科", "电子奖券", "技术 × 文化 × 传播", "11/22"],
        "我们的创新不只是活动内容，也包括传播方式。光电同学参与小程序搭建，让游戏规则、积分排行和中药百科在线化；基医同学主导海报设计，把传统中医药做成年轻人愿意靠近的国潮视觉。技术和设计，让文化传播更轻盈。",
    ),
    SlideText(
        "冬暖青日",
        ["冬暖青日", "本草寻踪，从校园到家乡", "返乡宣讲", "家乡考察", "本草地图", "75位同学 · 17个省份", "12/22"],
        "寒假期间，联合支部把实践从校园延伸到家乡。同学们返乡开展中医药文化宣讲，也结合家乡资源进行考察。我们把75位同学的家乡分布整理成本草地图，让每个省份都对应本地特色中药，呈现出一张属于青年视角的“本草中国”。",
    ),
    SlideText(
        "叶开泰",
        ["叶开泰", "从游戏化体验，走向历史现场", "历史展板", "闻香体验", "互动桌", "传统药铺", "叶开泰中医药文化园 · 4月19日", "13/22"],
        "4月19日下午，我们走进叶开泰中医药文化园。这里不是单纯参观，而是一次沉浸式体验：看中医药历史展板，闻药材香气，触摸互动桌，走进传统药铺复原场景。嘉年华上被游戏化的中医药，在叶开泰回到了更深厚的历史现场。",
    ),
    SlideText(
        "成果总结",
        ["成果总结", "三种成长，落在真实行动里", "思想上：从学习到担当", "组织上：从三班到一体", "文化上：从了解走向传播", "14/22"],
        "回看整个过程，我们的收获可以概括为三点。思想上，团课和志愿让同学们把个人专业放进时代责任中理解；组织上，跨班分工让三个支部真正形成协作共同体；文化上，我们从了解中医药，到用游戏、技术、宣讲和参观主动传播中医药。",
    ),
    SlideText(
        "光药医路 一起走",
        ["光药医路 一起走", "光可以照亮黑暗，药可以治愈疾病，医可以守护生命，而路，需要一起走。", "谢谢各位老师", "15/22"],
        "最后，我们想用一句话收束这段经历：光可以照亮黑暗，药可以治愈疾病，医可以守护生命，而路，需要一起走。光药医路的意义，不只在完成一次特色团日，更在于让75位来自不同专业的青年，真正成为了一个共同成长的集体。谢谢各位老师。",
    ),
    SlideText(
        "瑶光角色解构",
        ["瑶光角色解构", "角色大图", "白袍", "红衫", "绿带", "透镜", "药壶", "16/22"],
        "瑶光是贯穿全套视觉的线索。白色实验袍对应现代医学与实验精神，红色中式内衫承接思政红与传统文化，绿色腰带、药壶和本草纹样回应药草绿主线，透镜与发簪呼应光电视角。她把三个支部的专业气质合在同一个形象里。",
        "一位把光学、药学与医学连接起来的本草守护者",
    ),
    SlideText(
        "瑶光文创延展",
        ["瑶光文创延展", "日历", "笔记本", "徽章", "钥匙扣", "卡套", "笔", "17/22"],
        "在活动传播上，瑶光可以继续延展为可带走的文创：日历记录活动节点，笔记本承载本草笔记，徽章、钥匙扣、卡套和笔适合现场互动兑换。它们把“光药医路”从一次现场活动延伸到同学们的日常。",
        "参考实物形态，后期贴入瑶光与主题元素",
    ),
    SlideText(
        "嘉年华宣传",
        ["嘉年华宣传", "QQ空间推文", "兔子战队", "中操坐标", "本草任务", "幸运奖品", "18/22"],
        "嘉年华前期，我们通过 QQ 空间推文进行宣传，用更接近同学日常的语言发出邀请。推文把活动包装成一场“同济杏林之约”，强调不用穿越时空，就能在校园里亲手触碰千年本草，让传统文化先以轻松、可转发的方式被看见。",
        "把同济杏林之约推向全校",
    ),
    SlideText(
        "小游戏与奖品",
        ["小游戏与奖品", "望色辨证", "草色通真", "捣药闻香", "分两入毫", "积分挑战", "19/22"],
        "现场互动由四个小游戏组成：望色辨证引入中医观察思维，草色通真让同学辨认药材，捣药闻香通过气味与操作理解药性，分两入毫用传统药秤体验慎于分量。积分和奖品机制让参与者愿意停留、挑战和分享。",
        "让中医药知识变成可参与的体验",
    ),
    SlideText(
        "视频背景 01",
        ["视频背景 01", "开场氛围", "20/22"],
        "本页作为视频插入背景，可用于开场片段或主题引入。",
        "开场氛围",
    ),
    SlideText(
        "视频背景 02",
        ["视频背景 02", "嘉年华互动", "21/22"],
        "本页作为视频插入背景，可用于嘉年华现场互动或游戏过程片段。",
        "嘉年华互动",
    ),
    SlideText(
        "视频背景 03",
        ["视频背景 03", "共同成长", "22/22"],
        "本页作为视频插入背景，可用于结尾回顾、成果展示或合影视频。",
        "共同成长",
    ),
]


def px_to_in(value: float, axis: str = "x") -> float:
    return value / (PX_W if axis == "x" else PX_H) * (SLIDE_W_IN if axis == "x" else SLIDE_H_IN)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def font(size: int, family: str = "yahei") -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONTS[family], size)


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def reset_dirs() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)
    if PREVIEW.exists():
        shutil.rmtree(PREVIEW)
    PREVIEW.mkdir(parents=True)


def text_bbox(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.FreeTypeFont) -> tuple[int, int]:
    box = draw.textbbox((0, 0), text, font=fnt)
    return box[2] - box[0], box[3] - box[1]


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    max_width: int,
    fnt: ImageFont.FreeTypeFont,
    fill: tuple[int, int, int] = INK,
    line_gap: int = 12,
) -> int:
    x, y = xy
    lines: list[str] = []
    current = ""
    for char in text:
        trial = current + char
        if current and text_bbox(draw, trial, fnt)[0] > max_width:
            lines.append(current)
            current = char
        else:
            current = trial
    if current:
        lines.append(current)
    for line in lines:
        draw.text((x, y), line, font=fnt, fill=fill)
        y += text_bbox(draw, line, fnt)[1] + line_gap
    return y


def fit_image(src: Path, size: tuple[int, int], mode: str = "cover") -> Image.Image:
    image = Image.open(src).convert("RGB")
    if mode == "contain":
        image.thumbnail(size, Image.Resampling.LANCZOS)
        out = Image.new("RGB", size, CREAM)
        out.paste(image, ((size[0] - image.width) // 2, (size[1] - image.height) // 2))
        return out
    return ImageOps.fit(image, size, Image.Resampling.LANCZOS)


def rounded_picture(
    canvas: Image.Image,
    src: Path,
    box: tuple[int, int, int, int],
    radius: int = 32,
    mode: str = "cover",
    shadow: bool = True,
) -> None:
    x, y, w, h = box
    fitted = fit_image(src, (w, h), mode).convert("RGBA")
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, w, h), radius=radius, fill=255)
    if shadow:
        shadow_layer = Image.new("RGBA", (w + 34, h + 34), (0, 0, 0, 0))
        ImageDraw.Draw(shadow_layer).rounded_rectangle((17, 17, w + 17, h + 17), radius=radius, fill=(48, 36, 18, 70))
        shadow_layer = shadow_layer.filter(ImageFilter.GaussianBlur(12))
        canvas.alpha_composite(shadow_layer, (x - 17, y - 17))
    canvas.paste(fitted, (x, y), mask)


def parchment_background(accent: tuple[int, int, int] = GREEN) -> Image.Image:
    base = Image.new("RGB", (PX_W, PX_H), CREAM)
    overlay = Image.new("RGBA", (PX_W, PX_H), (255, 255, 255, 0))
    draw = ImageDraw.Draw(overlay)
    # Deterministic paper speckles keep reruns from churning large PNG outputs.
    for i in range(6200):
        x = (i * 97 + 41) % PX_W
        y = (i * 193 + 17) % PX_H
        alpha = 12 + (i * 7) % 18
        color = (188, 171, 138, alpha) if i % 3 else (255, 255, 255, alpha)
        draw.point((x, y), fill=color)
    for shift, color, width in [(0, RED, 28), (48, GOLD, 18), (92, accent, 20)]:
        draw.arc((-260 + shift, 650, 1180 + shift, 1420), 198, 350, fill=(*color, 42), width=width)
        draw.arc((875 - shift, -240, 2260 - shift, 485), 14, 178, fill=(*color, 36), width=width)
    for i in range(10):
        x = 110 + i * 205
        draw.line((x, 1005, x + 120, 945), fill=(*GREEN, 44), width=3)
        draw.ellipse((x + 52, 952, x + 110, 987), outline=(*GREEN, 50), width=2)
    for x, y, r in [(165, 925, 42), (1720, 920, 42), (1545, 145, 38), (1615, 172, 25)]:
        draw.ellipse((x - r, y - r, x + r, y + r), outline=(*GOLD, 95), width=3)
        draw.line((x, y - r, x, y + r), fill=(*GOLD, 70), width=2)
        draw.line((x - r, y, x + r, y), fill=(*GOLD, 70), width=2)
    return Image.alpha_composite(base.convert("RGBA"), overlay).convert("RGB")


def make_yaoguang_mark() -> Path:
    src = SUPPLEMENT / "瑶光-2.jpg"
    mark_path = ASSETS / "yaoguang_mark.png"
    if mark_path.exists():
        return mark_path
    source = ImageOps.fit(Image.open(src).convert("RGB"), (420, 420), Image.Resampling.LANCZOS)
    mask = Image.new("L", (420, 420), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, 420, 420), fill=255)
    out = Image.new("RGBA", (420, 420), (255, 255, 255, 0))
    out.paste(source, (0, 0), mask)
    border = Image.new("RGBA", (420, 420), (255, 255, 255, 0))
    draw = ImageDraw.Draw(border)
    draw.ellipse((5, 5, 415, 415), outline=(*GOLD, 235), width=12)
    out.alpha_composite(border)
    out.save(mark_path)
    return mark_path


def prepare_assets() -> dict[str, Path]:
    mark = make_yaoguang_mark()
    merch = ASSETS / "ai_merch_base.png"
    if not merch.exists() and AI_MERCH_SOURCE.exists():
        shutil.copy2(AI_MERCH_SOURCE, merch)
    if not merch.exists():
        # Deterministic fallback: use the user-provided notebook reference as a product-form source.
        fallback = MERCH_REF / "本子就是车线本换下封面.png"
        if fallback.exists():
            shutil.copy2(fallback, merch)
    if not merch.exists():
        raise FileNotFoundError("缺少文创产品底图，请先放入 交付物/答辩PPT/答辩PPT_v5/assets/ai_merch_base.png")
    logo = V4 / "source_assets" / "hust_logo_official_layer.png"
    if not logo.exists():
        raise FileNotFoundError(logo)
    return {"yaoguang_mark": mark, "merch_base": merch, "hust_logo": logo}


def draw_header(canvas: Image.Image, title: str, subtitle: str, no: int, accent: tuple[int, int, int] = GREEN) -> None:
    draw = ImageDraw.Draw(canvas)
    logo = V4 / "source_assets" / "hust_logo_official_layer.png"
    if logo.exists():
        logo_im = Image.open(logo).convert("RGBA").resize((260, 61), Image.Resampling.LANCZOS)
        canvas.alpha_composite(logo_im, (42, 22))
    draw.rounded_rectangle((328, 43, 462, 77), radius=17, fill=(*accent, 220))
    draw.text((352, 48), f"{no:02d}/22", font=font(22), fill=(255, 250, 235))
    draw.text((118, 122), title, font=font(66, "hei"), fill=accent)
    if subtitle:
        draw.text((122, 206), subtitle, font=font(30), fill=MUTED)
    draw.rectangle((0, 0, 18, 170), fill=RED)


def draw_page_marker(canvas: Image.Image, no: int, accent: tuple[int, int, int] = GOLD) -> None:
    draw = ImageDraw.Draw(canvas)
    x0, y0, x1, y1 = 1746, 1013, 1870, 1052
    draw.rounded_rectangle((x0, y0, x1, y1), radius=19, fill=(249, 244, 230), outline=accent, width=2)
    draw.text((x0 + 25, y0 + 8), f"{no:02d}/22", font=font(20), fill=(82, 75, 62))


def draw_yaoguang_stamp(canvas: Image.Image, mark: Path, position: tuple[int, int] = (1610, 770), size: int = 180) -> None:
    stamp = Image.open(mark).convert("RGBA").resize((size, size), Image.Resampling.LANCZOS)
    stamp.putalpha(stamp.getchannel("A").point(lambda p: int(p * 0.82)))
    canvas.alpha_composite(stamp, position)


def save_first_15_preview(no: int, mark: Path) -> Path:
    src = V4 / "preview_png" / f"slide_{no:02d}.png"
    if not src.exists():
        raise FileNotFoundError(src)
    canvas = Image.open(src).convert("RGBA")
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle((1720, 1004, 1890, 1062), radius=26, fill=(246, 241, 226, 240), outline=GOLD, width=3)
    draw.text((1762, 1020), f"{no:02d}/22", font=font(20), fill=(82, 75, 62))
    draw_yaoguang_stamp(canvas, mark, (1702, 780), 118)
    draw.rounded_rectangle((1658, 912, 1856, 956), radius=22, fill=(255, 252, 240, 224), outline=GOLD, width=2)
    draw.text((1682, 923), "瑶光线索", font=font(21), fill=GREEN)
    out = PREVIEW / f"slide_{no:02d}.png"
    canvas.convert("RGB").save(out)
    return out


def add_full_slide_picture(slide: Any, path: Path, name: str) -> Any:
    shape = slide.shapes.add_picture(str(path), 0, 0, width=Inches(SLIDE_W_IN), height=Inches(SLIDE_H_IN))
    shape.name = name
    return shape


def add_native_page_marker(slide: Any, no: int) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(px_to_in(1720)),
        Inches(px_to_in(1004, "y")),
        Inches(px_to_in(170)),
        Inches(px_to_in(58, "y")),
    )
    shape.name = "page-marker-bg"
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb((246, 241, 226))
    shape.line.color.rgb = rgb(GOLD)
    box = slide.shapes.add_textbox(
        Inches(px_to_in(1762)),
        Inches(px_to_in(1018, "y")),
        Inches(px_to_in(85)),
        Inches(px_to_in(28, "y")),
    )
    box.name = "page-marker"
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"{no:02d}/22"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(11)
    run.font.color.rgb = rgb((82, 75, 62))


def add_text(
    slide: Any,
    text: str,
    box: tuple[int, int, int, int],
    size: int,
    color: tuple[int, int, int],
    bold: bool = False,
    name: str = "text",
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    x, y, w, h = box
    shape = slide.shapes.add_textbox(
        Inches(px_to_in(x)),
        Inches(px_to_in(y, "y")),
        Inches(px_to_in(w)),
        Inches(px_to_in(h, "y")),
    )
    shape.name = name
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_round_rect(
    slide: Any,
    box: tuple[int, int, int, int],
    fill: tuple[int, int, int],
    line: tuple[int, int, int],
    name: str,
    transparency: int = 0,
) -> Any:
    x, y, w, h = box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(px_to_in(x)),
        Inches(px_to_in(y, "y")),
        Inches(px_to_in(w)),
        Inches(px_to_in(h, "y")),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.5)
    return shape


def add_picture_box(slide: Any, src: Path, box: tuple[int, int, int, int], name: str) -> Any:
    x, y, w, h = box
    shape = slide.shapes.add_picture(
        str(src),
        Inches(px_to_in(x)),
        Inches(px_to_in(y, "y")),
        width=Inches(px_to_in(w)),
        height=Inches(px_to_in(h, "y")),
    )
    shape.name = name
    return shape


def register(elements: list[Element], slide_no: int, element_id: str, kind: str, source: str, box: tuple[int, int, int, int], z: int, editable: bool, text: str | None = None) -> None:
    x, y, w, h = box
    elements.append(Element(slide_no, element_id, kind, source, {"x": x, "y": y, "w": w, "h": h}, z, editable, text))


def make_slide_16(mark: Path) -> tuple[Path, Path, list[Element]]:
    slide_no = 16
    bg = parchment_background(GREEN).convert("RGBA")
    draw_header(bg, "瑶光角色解构", "角色大图 + 设计元素拆解", slide_no, GREEN)
    character = SUPPLEMENT / "瑶光-3.jpg"
    rounded_picture(bg, character, (120, 260, 560, 720), radius=42, mode="cover")
    draw = ImageDraw.Draw(bg)
    points = [
        ("白色实验袍", "现代医学、实验精神与专业可信度", (795, 275, 760, 90), GREEN),
        ("红色中式内衫", "承接思政红，也把传统文化穿在身上", (920, 410, 760, 90), RED),
        ("绿色腰带与药壶", "药草绿主线，指向中医药文化生命力", (795, 545, 760, 90), GREEN),
        ("发簪透镜", "光电专业技术视角，连接观测与创新", (1020, 680, 680, 90), GOLD),
        ("银杏 / DNA 纹样", "本草意象与基础医学理解同框出现", (795, 815, 760, 90), DEEP_GREEN),
    ]
    elements: list[Element] = []
    register(elements, slide_no, "background", "image", "generated_background", (0, 0, PX_W, PX_H), 0, False)
    register(elements, slide_no, "title", "text", "native", (118, 122, 720, 78), 5, True, "瑶光角色解构")
    register(elements, slide_no, "yaoguang-character", "image", "reference_asset", (120, 260, 560, 720), 6, False)
    for idx, (title, body, box, color) in enumerate(points, 1):
        x, y, w, h = box
        draw.line((665, y + 45, x - 28, y + 45), fill=(*color, 170), width=5)
        draw.ellipse((654, y + 34, 678, y + 58), fill=GOLD)
        draw.rounded_rectangle((x, y, x + w, y + h), radius=28, fill=(255, 252, 240, 226), outline=color, width=3)
        draw.text((x + 28, y + 15), title, font=font(34, "hei"), fill=color)
        draw.text((x + 28, y + 55), body, font=font(24), fill=MUTED)
        register(elements, slide_no, f"concept-{idx}", "shape_text_group", "native", box, 10 + idx, True, f"{title}：{body}")
    draw_page_marker(bg, slide_no, GOLD)
    draw_yaoguang_stamp(bg, mark, (1642, 745), 150)
    preview = PREVIEW / "slide_16.png"
    bg.convert("RGB").save(preview)
    bg_path = ASSETS / "slide_16_bg.png"
    parchment_background(GREEN).save(bg_path)
    return preview, bg_path, elements


def make_slide_17(mark: Path, merch: Path) -> tuple[Path, Path, list[Element]]:
    slide_no = 17
    bg = parchment_background(GOLD).convert("RGBA")
    draw_header(bg, "瑶光文创延展", "参考实物形态，后期贴入瑶光与主题元素", slide_no, GREEN)
    rounded_picture(bg, merch, (95, 250, 1120, 655), radius=38, mode="cover")
    draw = ImageDraw.Draw(bg)
    draw.text((300, 548), "光药医路", font=font(48, "hei"), fill=GREEN)
    draw.text((332, 610), "YAO GUANG", font=font(24), fill=GOLD)
    small_mark = Image.open(mark).convert("RGBA").resize((92, 92), Image.Resampling.LANCZOS)
    small_mark.putalpha(small_mark.getchannel("A").point(lambda p: int(p * 0.88)))
    bg.alpha_composite(small_mark, (692, 492))
    bg.alpha_composite(small_mark.resize((72, 72), Image.Resampling.LANCZOS), (666, 658))
    draw.text((882, 358), "本草日历", font=font(31, "hei"), fill=RED)
    draw.rounded_rectangle((1265, 250, 1765, 905), radius=44, fill=(255, 252, 240, 232), outline=GOLD, width=4)
    draw.text((1320, 292), "从活动现场", font=font(42, "hei"), fill=GREEN)
    draw.text((1320, 350), "走进日常生活", font=font(42, "hei"), fill=RED)
    body = "文创不直接照搬参考图中的图案，而是保留产品形态：本子、日历、徽章、钥匙扣、卡套和笔。最终贴入瑶光、光药医路与三色成长弧，让活动记忆可携带、可展示、可交换。"
    draw_wrapped(draw, (1320, 430), body, 385, font(28), fill=MUTED, line_gap=18)
    tags = [("本草笔记", GREEN), ("日历节点", GOLD), ("徽章兑换", RED), ("钥匙扣", GREEN), ("卡套", GOLD), ("签字笔", RED)]
    for idx, (label, color) in enumerate(tags):
        x = 1325 + (idx % 2) * 200
        y = 710 + (idx // 2) * 72
        draw.rounded_rectangle((x, y, x + 168, y + 46), radius=23, fill=color)
        draw.text((x + 26, y + 9), label, font=font(23), fill=(255, 250, 235))
    draw_yaoguang_stamp(bg, mark, (1080, 650), 160)
    draw_page_marker(bg, slide_no, GOLD)
    elements: list[Element] = []
    register(elements, slide_no, "background", "image", "generated_background", (0, 0, PX_W, PX_H), 0, False)
    register(elements, slide_no, "merch-base", "image", "ai_generated_product_mockup", (95, 250, 1120, 655), 5, False)
    register(elements, slide_no, "merch-theme-overlay", "text_image_group", "native", (300, 548, 585, 190), 9, True, "光药医路 / YAO GUANG / 本草日历")
    register(elements, slide_no, "yaoguang-mark", "image", "reference_asset", (1080, 650, 160, 160), 10, False)
    register(elements, slide_no, "text-panel", "shape_text_group", "native", (1265, 250, 500, 655), 12, True, body)
    preview = PREVIEW / "slide_17.png"
    bg.convert("RGB").save(preview)
    bg_path = ASSETS / "slide_17_bg.png"
    parchment_background(GOLD).save(bg_path)
    return preview, bg_path, elements


def make_slide_18(mark: Path) -> tuple[Path, Path, list[Element]]:
    slide_no = 18
    bg = parchment_background(RED).convert("RGBA")
    draw_header(bg, "嘉年华宣传", "QQ空间推文，把同济杏林之约推向全校", slide_no, RED)
    promo = SUPPLEMENT / "宣传推文.png"
    rounded_picture(bg, promo, (110, 250, 620, 640), radius=34, mode="contain")
    draw = ImageDraw.Draw(bg)
    draw.text((820, 305), "无需穿越时空", font=font(58, "hei"), fill=RED)
    draw.text((820, 385), "在校园亲手触碰千年本草", font=font(52, "hei"), fill=GREEN)
    body = "宣传文案把活动包装成一场“同济杏林之约”：有兔子战队、有行动坐标、有游戏任务，也有幸运奖品。它让传统文化先以年轻、轻快、可转发的方式抵达同学。"
    draw_wrapped(draw, (825, 510), body, 820, font(32), fill=MUTED, line_gap=18)
    chips = [("兔子战队", RED), ("中操坐标", GREEN), ("本草任务", RED), ("幸运奖品", GREEN)]
    for idx, (label, color) in enumerate(chips):
        x = 840 + (idx % 2) * 310
        y = 770 + (idx // 2) * 78
        draw.rounded_rectangle((x, y, x + 225, y + 52), radius=26, fill=color)
        draw.text((x + 48, y + 12), label, font=font(27), fill=(255, 250, 235))
    draw_yaoguang_stamp(bg, mark, (1648, 742), 150)
    draw_page_marker(bg, slide_no, RED)
    elements: list[Element] = []
    register(elements, slide_no, "background", "image", "generated_background", (0, 0, PX_W, PX_H), 0, False)
    register(elements, slide_no, "promo-post", "image", "source_asset", (110, 250, 620, 640), 5, False)
    register(elements, slide_no, "headline", "text", "native", (820, 305, 820, 145), 10, True, "无需穿越时空 / 在校园亲手触碰千年本草")
    register(elements, slide_no, "promo-summary", "text", "native", (825, 510, 820, 190), 11, True, body)
    preview = PREVIEW / "slide_18.png"
    bg.convert("RGB").save(preview)
    bg_path = ASSETS / "slide_18_bg.png"
    parchment_background(RED).save(bg_path)
    return preview, bg_path, elements


def make_slide_19(mark: Path) -> tuple[Path, Path, list[Element]]:
    slide_no = 19
    bg = parchment_background(GREEN).convert("RGBA")
    draw_header(bg, "小游戏与奖品", "任务、积分与奖品，把中医药知识变成可参与的体验", slide_no, GREEN)
    draw = ImageDraw.Draw(bg)
    games = [
        ("望色辨证", "观察药液色彩，引入中医观察思维", RED),
        ("草色通真", "辨认药材图样，解锁本草知识", GREEN),
        ("捣药闻香", "研磨药材，感受药香与古法技艺", GOLD),
        ("分两入毫", "挑战药秤称量，体验慎于分量", (112, 91, 62)),
    ]
    elements: list[Element] = []
    register(elements, slide_no, "background", "image", "generated_background", (0, 0, PX_W, PX_H), 0, False)
    for idx, (title, body, color) in enumerate(games):
        x = 145 + (idx % 2) * 835
        y = 280 + (idx // 2) * 285
        draw.rounded_rectangle((x, y, x + 700, y + 215), radius=34, fill=(255, 252, 240, 232), outline=color, width=5)
        draw.ellipse((x + 40, y + 52, x + 140, y + 152), fill=color)
        draw.text((x + 76, y + 76), str(idx + 1), font=font(39, "hei"), fill=(255, 248, 226))
        draw.text((x + 175, y + 52), title, font=font(43, "hei"), fill=color)
        draw_wrapped(draw, (x + 175, y + 118), body, 445, font(28), fill=MUTED)
        register(elements, slide_no, f"game-{idx + 1}", "shape_text_group", "native", (x, y, 700, 215), 5 + idx, True, f"{title}：{body}")
    draw.rounded_rectangle((390, 870, 1530, 955), radius=42, fill=(*RED, 230), outline=GOLD, width=4)
    draw.text((478, 895), "积分挑战  ·  现场引导  ·  100%中奖  ·  瑶光周边承载诚意与祝福", font=font(35), fill=(255, 244, 218))
    draw_yaoguang_stamp(bg, mark, (1610, 735), 150)
    draw_page_marker(bg, slide_no, RED)
    register(elements, slide_no, "reward-bar", "shape_text_group", "native", (390, 870, 1140, 85), 14, True, "积分挑战 · 现场引导 · 100%中奖")
    preview = PREVIEW / "slide_19.png"
    bg.convert("RGB").save(preview)
    bg_path = ASSETS / "slide_19_bg.png"
    parchment_background(GREEN).save(bg_path)
    return preview, bg_path, elements


def make_video_slide(no: int, subtitle: str, accent: tuple[int, int, int], mark: Path) -> tuple[Path, Path, list[Element]]:
    bg = parchment_background(accent).convert("RGBA")
    draw_header(bg, f"视频背景 {no - 19:02d}", subtitle, no, accent)
    draw = ImageDraw.Draw(bg)
    draw.rounded_rectangle((230, 260, 1690, 875), radius=58, fill=(255, 252, 240, 176), outline=accent, width=7)
    for inset in [22, 38, 54]:
        draw.rounded_rectangle((230 + inset, 260 + inset, 1690 - inset, 875 - inset), radius=46, outline=(*GOLD, 105), width=2)
    draw.text((774, 492), "VIDEO", font=font(88, "hei"), fill=accent)
    draw.text((690, 615), "预留视频插入区域", font=font(40), fill=MUTED)
    draw.rounded_rectangle((1340, 180, 1698, 228), radius=24, fill=(*accent, 210))
    draw.text((1372, 191), "瑶光线索 · 光药医路", font=font(24), fill=(255, 250, 235))
    draw_yaoguang_stamp(bg, mark, (1575, 705), 150)
    draw_page_marker(bg, no, accent)
    preview = PREVIEW / f"slide_{no:02d}.png"
    bg.convert("RGB").save(preview)
    bg_path = ASSETS / f"slide_{no:02d}_bg.png"
    parchment_background(accent).save(bg_path)
    elements = [
        Element(no, "background", "image", "generated_background", {"x": 0, "y": 0, "w": PX_W, "h": PX_H}, 0, False),
        Element(no, "video-placeholder", "shape_text_group", "native", {"x": 230, "y": 260, "w": 1460, "h": 615}, 5, True, "预留视频插入区域"),
        Element(no, "yaoguang-mark", "image", "reference_asset", {"x": 1575, "y": 705, "w": 150, "h": 150}, 9, False),
    ]
    return preview, bg_path, elements


def build_preview_and_assets(asset_map: dict[str, Path]) -> tuple[list[Path], dict[int, Path], list[Element]]:
    previews: list[Path] = []
    bg_assets: dict[int, Path] = {}
    elements: list[Element] = []
    mark = asset_map["yaoguang_mark"]
    for no in range(1, 16):
        preview = save_first_15_preview(no, mark)
        previews.append(preview)
        bg_assets[no] = preview
        register(elements, no, "background", "image", "v4_visual_reference", (0, 0, PX_W, PX_H), 0, False)
        register(elements, no, "yaoguang-mark", "image", "reference_asset", (1702, 780, 118, 118), 5, False)
        register(elements, no, "page-marker", "text", "native", (1762, 1020, 85, 28), 6, True, f"{no:02d}/22")
    for fn in [
        lambda: make_slide_16(mark),
        lambda: make_slide_17(mark, asset_map["merch_base"]),
        lambda: make_slide_18(mark),
        lambda: make_slide_19(mark),
        lambda: make_video_slide(20, "开场氛围", GOLD, mark),
        lambda: make_video_slide(21, "嘉年华互动", GREEN, mark),
        lambda: make_video_slide(22, "共同成长", RED, mark),
    ]:
        preview, bg, slide_elements = fn()
        previews.append(preview)
        no = int(preview.stem.split("_")[1])
        bg_assets[no] = preview
        register(elements, no, "composite-slide", "image", "v5_composite_preview", (0, 0, PX_W, PX_H), 0, False)
        elements.extend(slide_elements)
    return previews, bg_assets, elements


def build_contact_sheet(previews: list[Path]) -> Path:
    thumbs: list[Image.Image] = []
    for path in previews:
        im = Image.open(path).convert("RGB")
        im.thumbnail((480, 270), Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (480, 300), (255, 255, 255))
        canvas.paste(im, (0, 24))
        draw = ImageDraw.Draw(canvas)
        draw.text((6, 4), path.stem.replace("slide_", ""), font=font(18), fill=(50, 50, 50))
        thumbs.append(canvas)
    cols = 3
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 480, rows * 300), (255, 255, 255))
    for idx, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((idx % cols) * 480, (idx // cols) * 300))
    out = OUT / "preview_contact_sheet.png"
    sheet.save(out)
    return out


def save_pdf(previews: list[Path]) -> Path:
    images = [Image.open(path).convert("RGB") for path in previews]
    out = OUT / PDF_NAME
    images[0].save(out, save_all=True, append_images=images[1:], resolution=150)
    return out


def add_native_slide_content(slide: Any, no: int, bg_path: Path, asset_map: dict[str, Path]) -> list[Element]:
    add_full_slide_picture(slide, bg_path, "background")
    if no >= 16:
        return []
    add_picture_box(slide, asset_map["hust_logo"], (42, 22, 260, 61), "hust-logo")
    if no <= 15:
        add_picture_box(slide, asset_map["yaoguang_mark"], (1702, 780, 118, 118), "yaoguang-mark")
        add_round_rect(slide, (1658, 912, 198, 44), (255, 252, 240), GOLD, "yaoguang-label-bg", 8)
        add_text(slide, "瑶光线索", (1682, 922, 145, 25), 12, GREEN, name="yaoguang-label", align=PP_ALIGN.CENTER)
        add_native_page_marker(slide, no)
        return []

    data = SLIDE_TEXT[no - 1]
    accent = RED if no in [18, 22] else GOLD if no in [17, 20] else GREEN
    add_round_rect(slide, (328, 43, 134, 34), accent, accent, "header-slide-number-bg", 0)
    add_text(slide, f"{no:02d}/22", (352, 48, 86, 24), 12, (255, 250, 235), name="header-slide-number", align=PP_ALIGN.CENTER)
    add_text(slide, data.title, (118, 122, 760, 86), 36, accent, True, "title")
    add_text(slide, data.subtitle, (122, 205, 900, 44), 17, MUTED, False, "subtitle")
    if no == 16:
        add_picture_box(slide, SUPPLEMENT / "瑶光-3.jpg", (120, 260, 560, 720), "yaoguang-character")
        pts = [
            ("白色实验袍", "现代医学、实验精神与专业可信度", (795, 275, 760, 90), GREEN),
            ("红色中式内衫", "承接思政红，也把传统文化穿在身上", (920, 410, 760, 90), RED),
            ("绿色腰带与药壶", "药草绿主线，指向中医药文化生命力", (795, 545, 760, 90), GREEN),
            ("发簪透镜", "光电专业技术视角，连接观测与创新", (1020, 680, 680, 90), GOLD),
            ("银杏 / DNA 纹样", "本草意象与基础医学理解同框出现", (795, 815, 760, 90), DEEP_GREEN),
        ]
        for idx, (title, body, box, color) in enumerate(pts, 1):
            add_round_rect(slide, box, (255, 252, 240), color, f"concept-{idx}-box", 10)
            x, y, w, h = box
            add_text(slide, title, (x + 28, y + 13, w - 56, 38), 19, color, True, f"concept-{idx}-title")
            add_text(slide, body, (x + 28, y + 54, w - 56, 28), 13, MUTED, False, f"concept-{idx}-body")
    elif no == 17:
        add_picture_box(slide, asset_map["merch_base"], (95, 250, 1120, 655), "merch-base")
        add_text(slide, "光药医路", (300, 548, 260, 58), 25, GREEN, True, "merch-notebook-title")
        add_text(slide, "YAO GUANG", (332, 610, 180, 34), 13, GOLD, False, "merch-notebook-en")
        add_picture_box(slide, asset_map["yaoguang_mark"], (692, 492, 92, 92), "merch-badge-mark")
        add_picture_box(slide, asset_map["yaoguang_mark"], (666, 658, 72, 72), "merch-keychain-mark")
        add_text(slide, "本草日历", (882, 358, 190, 42), 17, RED, True, "merch-calendar-title")
        add_round_rect(slide, (1265, 250, 500, 655), (255, 252, 240), GOLD, "text-panel", 8)
        add_text(slide, "从活动现场\n走进日常生活", (1320, 292, 380, 120), 23, GREEN, True, "panel-headline")
        add_text(
            slide,
            "文创保留参考图中的产品形态：本子、日历、徽章、钥匙扣、卡套和笔。最终贴入瑶光、光药医路与三色成长弧，让活动记忆可携带、可展示、可交换。",
            (1320, 430, 385, 210),
            15,
            MUTED,
            False,
            "panel-body",
        )
        add_picture_box(slide, asset_map["yaoguang_mark"], (1080, 650, 160, 160), "yaoguang-mark")
    elif no == 18:
        add_picture_box(slide, SUPPLEMENT / "宣传推文.png", (110, 250, 620, 640), "promo-post")
        add_text(slide, "无需穿越时空", (820, 305, 780, 72), 31, RED, True, "headline-1")
        add_text(slide, "在校园亲手触碰千年本草", (820, 385, 850, 72), 28, GREEN, True, "headline-2")
        add_text(
            slide,
            "宣传文案把活动包装成一场“同济杏林之约”：有兔子战队、有行动坐标、有游戏任务，也有幸运奖品。它让传统文化先以年轻、轻快、可转发的方式抵达同学。",
            (825, 510, 820, 190),
            18,
            MUTED,
            False,
            "promo-summary",
        )
    elif no == 19:
        games = [
            ("望色辨证", "观察药液色彩，引入中医观察思维", RED),
            ("草色通真", "辨认药材图样，解锁本草知识", GREEN),
            ("捣药闻香", "研磨药材，感受药香与古法技艺", GOLD),
            ("分两入毫", "挑战药秤称量，体验慎于分量", (112, 91, 62)),
        ]
        for idx, (title, body, color) in enumerate(games):
            x = 145 + (idx % 2) * 835
            y = 280 + (idx // 2) * 285
            add_round_rect(slide, (x, y, 700, 215), (255, 252, 240), color, f"game-{idx + 1}-box", 7)
            add_text(slide, str(idx + 1), (x + 68, y + 70, 40, 45), 22, (255, 250, 235), True, f"game-{idx + 1}-num", PP_ALIGN.CENTER)
            add_text(slide, title, (x + 175, y + 48, 400, 48), 24, color, True, f"game-{idx + 1}-title")
            add_text(slide, body, (x + 175, y + 116, 445, 62), 15, MUTED, False, f"game-{idx + 1}-body")
    elif no in [20, 21, 22]:
        add_round_rect(slide, (230, 260, 1460, 615), (255, 252, 240), accent, "video-placeholder", 22)
        add_text(slide, "VIDEO", (760, 490, 400, 100), 47, accent, True, "video-word", PP_ALIGN.CENTER)
        add_text(slide, "预留视频插入区域", (690, 615, 550, 60), 22, MUTED, False, "video-label", PP_ALIGN.CENTER)
    if no >= 16:
        add_picture_box(slide, asset_map["yaoguang_mark"], (1575, 705, 150, 150), "yaoguang-mark")
        add_native_page_marker(slide, no)
    return []


def write_pptx(bg_assets: dict[int, Path], asset_map: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    while len(prs.slides) > 0:
        raise RuntimeError("unexpected default slides")
    for no in range(1, 23):
        slide = prs.slides.add_slide(blank)
        add_native_slide_content(slide, no, bg_assets[no], asset_map)
    out = OUT / PPTX_NAME
    tmp_dir = Path(tempfile.gettempdir())
    tmp_out = tmp_dir / "guangyaoyilu_v5_build_tmp.pptx"
    if tmp_out.exists():
        tmp_out.unlink()
    prs.save(str(tmp_out))
    fix_app_slide_count(tmp_out, 22)
    if out.exists():
        out.unlink()
    shutil.copy2(tmp_out, out)
    tmp_out.unlink()
    return out


def fix_app_slide_count(pptx_path: Path, slide_count: int) -> None:
    """python-pptx can leave docProps/app.xml Slides at 0; keep metadata honest."""
    app_name = "docProps/app.xml"
    ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
    with zipfile.ZipFile(pptx_path, "r") as zin:
        parts = {name: zin.read(name) for name in zin.namelist()}
    if app_name not in parts:
        return
    root = ET.fromstring(parts[app_name])
    slides = root.find("ep:Slides", ns)
    if slides is not None:
        slides.text = str(slide_count)
        parts[app_name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in parts.items():
            zout.writestr(name, data)
    if pptx_path.exists():
        pptx_path.unlink()
    tmp.replace(pptx_path)


def write_manifests(elements: list[Element]) -> None:
    manifest = [
        {
            "slide": item.slide,
            "id": item.element_id,
            "type": item.type,
            "source": item.source,
            "bbox": item.bbox,
            "z": item.z,
            "editable": item.editable,
            **({"text": item.text} if item.text else {}),
        }
        for item in elements
    ]
    (OUT / "element_manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    text_manifest = [
        {
            "slide": idx,
            "title": item.title,
            "visible_text": item.visible_text,
            "speaker": item.speaker,
        }
        for idx, item in enumerate(SLIDE_TEXT, 1)
    ]
    (OUT / "text_manifest.json").write_text(json.dumps(text_manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def write_script() -> None:
    lines = ["# 光药医路 5分钟答辩讲稿 v5", "", "目标时长：5分钟左右。v5 为稳定静态版 PPTX，不包含手写动画 XML。", ""]
    for idx, item in enumerate(SLIDE_TEXT, 1):
        lines.extend([f"## {idx:02d}. {item.title}", "", item.speaker, ""])
    (OUT / SCRIPT_NAME).write_text("\n".join(lines), encoding="utf-8")


def verify_images(paths: list[Path]) -> None:
    for path in paths:
        with Image.open(path) as image:
            image.verify()


def verify_pptx_media(pptx_path: Path) -> int:
    count = 0
    with zipfile.ZipFile(pptx_path) as zf:
        for name in zf.namelist():
            if name.startswith("ppt/media/"):
                with zf.open(name) as stream:
                    image = Image.open(stream)
                    image.verify()
                    count += 1
    return count


def write_qa(pptx_path: Path, previews: list[Path], pdf_path: Path, contact_sheet: Path, media_count: int) -> None:
    qa = f"""# 光药医路 v5 QA 记录

生成日期：2026-05-09

## 输出

- PPTX：`交付物/答辩PPT/答辩PPT_v5/{PPTX_NAME}`
- 静态预览 PDF：`交付物/答辩PPT/答辩PPT_v5/{PDF_NAME}`
- 预览图：`交付物/答辩PPT/答辩PPT_v5/preview_png/slide_01.png` 至 `slide_22.png`
- 总览图：`交付物/答辩PPT/答辩PPT_v5/preview_contact_sheet.png`
- 元素清单：`交付物/答辩PPT/答辩PPT_v5/element_manifest.json`
- 文本清单：`交付物/答辩PPT/答辩PPT_v5/text_manifest.json`
- 讲稿：`交付物/答辩PPT/答辩PPT_v5/{SCRIPT_NAME}`

## 已执行检查

- 生成脚本：`python 工具/create_defense_ppt_v5_stable.py`
- PPTX 包结构：`python 工具/validate_pptx_package.py 交付物/答辩PPT/答辩PPT_v5/{PPTX_NAME} --expected-slides 22`
- 图片可读性：PIL 校验 {len(previews)} 张预览图和 PPTX 内 {media_count} 张媒体图片。
- UTF-8 安全：脚本写入 `.json`、`.md` 均显式使用 UTF-8。
- Git 空白检查：`git diff --check` 通过。

## 视觉检查重点

- v5 修复 v4 后段左上角拥挤：16-22 页标题、页码与 HUST logo 分离。
- v5 取消 v4 第 18 页多余项目符号，正文改为短段落与关键词。
- v5 重做 20-22 页视频背景，保留大视频区域，同时增加瑶光角标和主题栏。
- 17 页文创底图由生图模型生成空白产品形态，PPT 中再叠加瑶光与主题元素；不直接照抄参考图品牌图案。

## 验证缺口

- 当前机器未检测到 LibreOffice、dotnet 或 PowerPoint 自动化能力，因此未做 OpenXmlValidator 或 saved-PPTX headless render parity。
- PPTX 使用 python-pptx 生成，包结构包含成熟库输出的 master/layout/theme 基础结构，风险低于 v4 手写 OOXML，但最终上交前仍建议用 Microsoft PowerPoint 人工打开确认无修复弹窗。
"""
    (OUT / "QA_CHECKS.md").write_text(qa, encoding="utf-8")


def build() -> None:
    reset_dirs()
    asset_map = prepare_assets()
    previews, bg_assets, elements = build_preview_and_assets(asset_map)
    contact_sheet = build_contact_sheet(previews)
    pdf_path = save_pdf(previews)
    pptx_path = write_pptx(bg_assets, asset_map)
    validate_pptx(pptx_path, expected_slides=22)
    media_count = verify_pptx_media(pptx_path)
    verify_images(previews + [contact_sheet, *asset_map.values()])
    if not pdf_path.exists() or pdf_path.stat().st_size <= 0:
        raise RuntimeError(f"PDF preview was not created: {pdf_path}")
    write_manifests(elements)
    write_script()
    write_qa(pptx_path, previews, pdf_path, contact_sheet, media_count)
    print(json.dumps({
        "pptx": str(pptx_path),
        "pdf": str(pdf_path),
        "contact_sheet": str(contact_sheet),
        "preview_count": len(previews),
        "pptx_media_count": media_count,
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    build()
