"""
光药医路 · 交付物/快速答辩 快速交付构建脚本 v2

映射规则（已通过照片数量与照片框数量验证）：
  layers/slide_01-16  →  full-all-slides/new08-new22  (offset +7)
  layers/slide_17-22  →  full-all-slides/new01-new06  (仅logo)
  new07 缺失，跳过

画布：1672×941 px，EMU = px * 9525 (96 DPI)

Logo：等比例放置，仅指定 width，height 由 python-pptx 自动按图片原始比例计算
  hust_logo.png 原始尺寸 425×320，比例 1.328:1
  LOGO_W = 80px → 自动 h ≈ 60px（之前为 235×57 且被拉伸，现缩小约2倍并修正比例）
"""

import os
from pptx import Presentation
from pptx.util import Emu

BASE        = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR  = os.path.join(BASE, 'full-all-slides')
LAYERS_DIR  = os.path.join(BASE, 'layers')

# ── Canvas ─────────────────────────────────────────────────────────────────────
CANVAS_W_PX = 1672
CANVAS_H_PX = 941
EMU_PER_PX  = 9525   # 96 DPI: 1 inch = 914400 EMU, 1 px = 9525 EMU

def px(n):
    return Emu(int(n * EMU_PER_PX))

# ── HUST Logo ─────────────────────────────────────────────────────────────────
# 等比例放置：只指定 width；height 由 python-pptx 按图片原始长宽比自动计算
# 原图 425×320 = 1.328:1，width=80px → height≈60px（缩小约2倍，正确比例）
LOGO_W = 80    # px
LOGO_X = 15    # px，距左边缘
LOGO_Y = 12    # px，距顶边缘

# ── Photo placements (px, 1672×941 canvas) ────────────────────────────────────
# key = layers/slide_XX 文件夹编号
# 检测方法：在对应 new_XX.png 中寻找最亮矩形区域（高置信度 score > 0.87）
PHOTO_PLACEMENTS = {
    # layer_04 → new11 (三支部总览：3等宽框)  scores: 0.98/0.98/0.98
    '04': [
        {'file': 'photo_02.png', 'x':  80, 'y': 260, 'w': 458, 'h': 338},
        {'file': 'photo_01.png', 'x': 590, 'y': 310, 'w': 458, 'h': 338},
        {'file': 'photo_03.png', 'x':1110, 'y': 340, 'w': 458, 'h': 338},
    ],
    # layer_07 → new14 (初识白·破冰：1大+3小)  scores: 0.97/0.96/0.91/0.90
    '07': [
        {'file': 'photo_01.png', 'x': 140, 'y':  80, 'w': 848, 'h': 493},
        {'file': 'photo_04.png', 'x':1240, 'y':   0, 'w': 303, 'h': 228},
        {'file': 'photo_02.png', 'x': 990, 'y': 300, 'w': 303, 'h': 228},
        {'file': 'photo_03.png', 'x':1290, 'y': 290, 'w': 303, 'h': 228},
    ],
    # layer_08 → new15 (思政红·团课：2框上下)  scores: 0.93/0.46
    '08': [
        {'file': 'photo_01.png', 'x': 150, 'y': 240, 'w': 708, 'h': 413},
        {'file': 'photo_02.png', 'x': 860, 'y': 120, 'w': 708, 'h': 238},
    ],
    # layer_09 → new16 (思政红·志愿：1大+3小)  scores: 0.94/0.93/0.82/0.78
    '09': [
        {'file': 'photo_01.png', 'x': 260, 'y': 240, 'w': 763, 'h': 443},
        {'file': 'photo_02.png', 'x':1030, 'y': 270, 'w': 273, 'h': 343},
        {'file': 'photo_03.png', 'x':1300, 'y': 320, 'w': 273, 'h': 343},
        {'file': 'photo_04.png', 'x':   0, 'y':   0, 'w': 273, 'h': 343},
    ],
    # layer_10 → new17 (药草绿·嘉年华：1框)  score: 0.88
    '10': [
        {'file': 'photo_01.png', 'x':  30, 'y':  60, 'w': 843, 'h': 573},
    ],
    # layer_11 → new18  scores: 0.93/0.91/0.94
    '11': [
        {'file': 'photo_01.png', 'x': 170, 'y':  20, 'w': 358, 'h': 663},
        {'file': 'photo_02.png', 'x':1120, 'y': 170, 'w': 418, 'h': 418},
        {'file': 'photo_03.png', 'x':   0, 'y':   0, 'w': 170, 'h': 170},
    ],
    # layer_12 → new19  scores: 0.95/0.94
    '12': [
        {'file': 'photo_01.png', 'x':   0, 'y': 130, 'w': 313, 'h': 323},
        {'file': 'photo_02.png', 'x': 340, 'y': 130, 'w': 313, 'h': 323},
    ],
    # layer_13 → new20  scores: 0.94/0.87/0.86/0.70
    '13': [
        {'file': 'photo_01.png', 'x': 220, 'y': 220, 'w': 788, 'h': 428},
        {'file': 'photo_02.png', 'x':1050, 'y': 290, 'w': 248, 'h': 283},
        {'file': 'photo_03.png', 'x':1330, 'y': 280, 'w': 248, 'h': 283},
        {'file': 'photo_04.png', 'x': 130, 'y':   0, 'w': 248, 'h': 283},
    ],
    # layer_14 → new21  scores: 0.97/0.96/0.93
    '14': [
        {'file': 'photo_02.png', 'x': 100, 'y': 280, 'w': 498, 'h': 308},
        {'file': 'photo_01.png', 'x': 600, 'y': 270, 'w': 498, 'h': 308},
        {'file': 'photo_03.png', 'x':1100, 'y': 320, 'w': 498, 'h': 308},
    ],
}

# ── Slide sequence ─────────────────────────────────────────────────────────────
# (bg_filename, layer_folder_or_None, photo_placement_key_or_None)
# logo 取自 layer_folder/hust_logo.png；layer_folder=None 时用 slide_01 的 logo

def logo_path(layer_folder):
    return os.path.join(LAYERS_DIR, f'slide_{layer_folder}', 'hust_logo.png')

FALLBACK_LOGO = logo_path('01')

SLIDE_SEQUENCE = [
    # new01-06: 对应 layer_17-22（仅 logo）
    ('new01.png', '17', None),
    ('new02.png', '18', None),
    ('new03.png', '19', None),
    ('new04.png', '20', None),
    ('new05.png', '21', None),
    ('new06.png', '22', None),
    # new07 缺失，跳过
    # new08-22: 对应 layer_01-15（部分有照片）
    ('new08.png', '01', None),
    ('new09.png', '02', None),
    ('new10.png', '03', None),
    ('new11.png', '04', '04'),   # 三支部总览
    # 补充：幕后制作（插在 new11 和 new12 之间）
    ('补充-幕后制作-11和12间插入.png', None, None),
    ('new12.png', '05', None),
    ('new13.png', '06', None),
    ('new14.png', '07', '07'),   # 初识白·破冰
    ('new15.png', '08', '08'),   # 思政红·团课
    ('new16.png', '09', '09'),   # 思政红·志愿
    ('new17.png', '10', '10'),   # 药草绿·嘉年华
    # 补充：嘉年华（插在 new17 之后）
    ('嘉年华补充01-补充在17页后.png', None, None),
    ('嘉年华补充02-补充在17页后.png', None, None),
    ('new18.png', '11', '11'),
    ('new19.png', '12', '12'),
    ('new20.png', '13', '13'),
    ('new21.png', '14', '14'),
    ('new22.png', '15', None),
]

# ── Build ──────────────────────────────────────────────────────────────────────

def add_slide(prs, bg_path, layer_folder, photo_key):
    slide_layout = prs.slide_layouts[6]   # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 背景图：铺满全画布
    slide.shapes.add_picture(
        bg_path,
        px(0), px(0),
        px(CANVAS_W_PX), px(CANVAS_H_PX)
    )

    # HUST Logo：仅指定 width，height 自动按原图比例计算（等比，不拉伸）
    lp = logo_path(layer_folder) if layer_folder else FALLBACK_LOGO
    if not os.path.exists(lp):
        lp = FALLBACK_LOGO
    if os.path.exists(lp):
        slide.shapes.add_picture(
            lp,
            px(LOGO_X), px(LOGO_Y),
            width=px(LOGO_W)     # height=None → 自动等比计算
        )
    else:
        print(f'  [WARN] logo not found: {lp}')

    # 照片
    if photo_key and photo_key in PHOTO_PLACEMENTS:
        for p in PHOTO_PLACEMENTS[photo_key]:
            photo_dir = os.path.join(LAYERS_DIR, f'slide_{photo_key}')
            photo_path = os.path.join(photo_dir, p['file'])
            if os.path.exists(photo_path):
                slide.shapes.add_picture(
                    photo_path,
                    px(p['x']), px(p['y']),
                    px(p['w']), px(p['h'])
                )
            else:
                print(f'  [WARN] photo not found: {photo_path}')


def build():
    prs = Presentation()
    prs.slide_width  = Emu(CANVAS_W_PX * EMU_PER_PX)
    prs.slide_height = Emu(CANVAS_H_PX * EMU_PER_PX)

    for i, (bg_name, layer_folder, photo_key) in enumerate(SLIDE_SEQUENCE, 1):
        bg_path = os.path.join(SLIDES_DIR, bg_name)
        if not os.path.exists(bg_path):
            print(f'[SKIP] slide {i:02d}: {bg_name} not found')
            continue
        layer_label = f'layer_{layer_folder}' if layer_folder else 'fallback_logo'
        photo_label = f'+ {len(PHOTO_PLACEMENTS[photo_key])} photos' if photo_key else ''
        print(f'[{i:02d}/{len(SLIDE_SEQUENCE)}] {bg_name}  ({layer_label}{photo_label})')
        add_slide(prs, bg_path, layer_folder, photo_key)

    out = os.path.join(BASE, 'GuangYaoYiLu_quick.pptx')
    prs.save(out)
    print(f'\nSaved: {out}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    build()
