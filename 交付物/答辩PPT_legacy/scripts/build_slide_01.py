from __future__ import annotations

import json
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt


ROOT = Path(__file__).resolve().parents[1]
SLIDE_NO = 1
CANVAS_W = 1920
CANVAS_H = 1080
PPT_W = 9144000
PPT_H = 5143500

LAYER_DIR = ROOT / "layers" / "slide_01"
MANIFEST_PATH = LAYER_DIR / "manifest.json"
SINGLE_DIR = ROOT / "single_slides"
PREVIEW_DIR = ROOT / "preview"
PPTX_OUT = SINGLE_DIR / "slide_01.pptx"
PREVIEW_OUT = PREVIEW_DIR / "slide_01_preview.png"

FONT_DIR = Path("C:/Windows/Fonts")
FONT_FALLBACKS = {
    "\u5fae\u8f6f\u96c5\u9ed1": ["NotoSansSC-VF.ttf", "msyh.ttc", "simhei.ttf"],
    "\u5b8b\u4f53": ["NotoSerifSC-VF.ttf", "simsun.ttc", "msyh.ttc"],
    "\u534e\u6587\u884c\u6977": ["STXINGKA.TTF", "STKAITI.TTF", "simkai.ttf", "FZSTK.TTF"],
}


def px_to_emu_x(px: int | float) -> int:
    return int(px * PPT_W / CANVAS_W)


def px_to_emu_y(px: int | float) -> int:
    return int(px * PPT_H / CANVAS_H)


def parse_rgb(value: str) -> tuple[int, int, int]:
    raw = value.strip()
    if not raw.startswith("RGB(") or not raw.endswith(")"):
        raise ValueError(f"Unsupported color format: {value}")
    return tuple(int(part.strip()) for part in raw[4:-1].split(","))  # type: ignore[return-value]


def align_value(value: str) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(value.lower(), PP_ALIGN.LEFT)


def ensure_east_asian_font(run, font_name: str) -> None:
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()  # noqa: SLF001 - python-pptx exposes no public East Asian font setter.
    r_pr.set(qn("a:ea"), font_name)


def local_font_path(font_name: str) -> str | None:
    for candidate in FONT_FALLBACKS.get(font_name, []):
        path = FONT_DIR / candidate
        if path.exists():
            return str(path)
    for candidates in FONT_FALLBACKS.values():
        for candidate in candidates:
            path = FONT_DIR / candidate
            if path.exists():
                return str(path)
    return None


def preview_font(font_name: str, size_px: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    path = local_font_path(font_name)
    if path:
        return ImageFont.truetype(path, size_px)
    return ImageFont.load_default()


def add_manifest_images(slide, manifest: dict) -> None:
    for layer in sorted(manifest["layers"], key=lambda item: item["z"]):
        if layer["type"] != "image":
            continue
        image_path = LAYER_DIR / layer["file"]
        if not image_path.exists():
            raise FileNotFoundError(image_path)
        bbox = layer["bbox"]
        slide.shapes.add_picture(
            str(image_path),
            left=px_to_emu_x(bbox["x"]),
            top=px_to_emu_y(bbox["y"]),
            width=px_to_emu_x(bbox["w"]),
            height=px_to_emu_y(bbox["h"]),
        )


def add_manifest_text(slide, manifest: dict) -> None:
    for textbox in sorted(manifest["textboxes"], key=lambda item: item["z"]):
        bbox = textbox["bbox"]
        box = slide.shapes.add_textbox(
            px_to_emu_x(bbox["x"]),
            px_to_emu_y(bbox["y"]),
            px_to_emu_x(bbox["w"]),
            px_to_emu_y(bbox["h"]),
        )
        frame = box.text_frame
        frame.clear()
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align_value(textbox["align"])
        run = paragraph.add_run()
        run.text = textbox["content"]
        ensure_east_asian_font(run, textbox["font"])
        run.font.size = Pt(textbox["font_size"])
        run.font.color.rgb = RGBColor(*parse_rgb(textbox["color"]))


def build_preview(manifest: dict) -> None:
    preview = Image.new("RGBA", (CANVAS_W, CANVAS_H), (255, 255, 255, 255))
    for layer in sorted(manifest["layers"], key=lambda item: item["z"]):
        if layer["type"] != "image":
            continue
        image_path = LAYER_DIR / layer["file"]
        bbox = layer["bbox"]
        image = Image.open(image_path).convert("RGBA").resize((bbox["w"], bbox["h"]), Image.Resampling.LANCZOS)
        preview.alpha_composite(image, (bbox["x"], bbox["y"]))

    draw = ImageDraw.Draw(preview)
    for textbox in sorted(manifest["textboxes"], key=lambda item: item["z"]):
        bbox = textbox["bbox"]
        font = preview_font(textbox["font"], int(textbox["font_size"] * 1.38))
        color = parse_rgb(textbox["color"])
        text = textbox["content"]
        text_bbox = draw.textbbox((0, 0), text, font=font)
        text_w = text_bbox[2] - text_bbox[0]
        text_h = text_bbox[3] - text_bbox[1]
        if textbox["align"] == "center":
            x = bbox["x"] + (bbox["w"] - text_w) / 2
        elif textbox["align"] == "right":
            x = bbox["x"] + bbox["w"] - text_w
        else:
            x = bbox["x"]
        y = bbox["y"] + (bbox["h"] - text_h) / 2 - 2
        draw.text((x, y), text, font=font, fill=color)

    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    preview.convert("RGB").save(PREVIEW_OUT)


def validate_pptx(manifest: dict) -> dict[str, int | bool]:
    if not PPTX_OUT.exists():
        raise FileNotFoundError(PPTX_OUT)
    with zipfile.ZipFile(PPTX_OUT) as package:
        slide_xml = package.read("ppt/slides/slide1.xml")
    root = ET.fromstring(slide_xml)
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    pic_count = len(root.findall(".//p:pic", ns))
    sp_count = len(root.findall(".//p:sp", ns))
    all_text = "".join(node.text or "" for node in root.findall(".//a:t", ns))
    forbidden = ["01/22", "\u9875\u7801", "\u6821\u5fbd\u4f4d\u7f6e\u9884\u7559\u533a"]
    if pic_count <= 1:
        raise AssertionError("PPTX must contain multiple picture layers")
    if sp_count <= 0:
        raise AssertionError("PPTX must contain editable text boxes")
    if not any(layer["id"] == "hust_logo" for layer in manifest["layers"]):
        raise AssertionError("manifest must contain hust_logo")
    if any(item in all_text for item in forbidden):
        raise AssertionError("PPTX contains forbidden placeholder or page-number text")
    return {"pic_count": pic_count, "sp_count": sp_count, "forbidden_text_found": False}


def main() -> None:
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    SINGLE_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Emu(PPT_W)
    prs.slide_height = Emu(PPT_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_manifest_images(slide, manifest)
    add_manifest_text(slide, manifest)
    prs.save(PPTX_OUT)
    build_preview(manifest)
    result = validate_pptx(manifest)
    print(json.dumps({"pptx": str(PPTX_OUT), "preview": str(PREVIEW_OUT), **result}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
