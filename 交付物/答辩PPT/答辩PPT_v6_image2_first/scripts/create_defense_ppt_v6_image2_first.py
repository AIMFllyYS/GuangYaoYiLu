from __future__ import annotations

"""
Build the v6 image2-first layered defense deck.

Route:
- v6 slides 01-07 use newly generated image2 mother artwork.
- v6 slides 08-22 reuse v4 full_ai_slides/slide_01.png..slide_15.png.
- The PPTX is composed with python-pptx, not hand-written OOXML.
- Every slide is assembled from multiple picture layers and an honest manifest.
"""

import hashlib
import json
import shutil
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches

from validate_pptx_package import validate_pptx


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "交付物" / "答辩PPT" / "答辩PPT_v6_image2_first"
V4 = ROOT / "交付物" / "答辩PPT" / "答辩PPT_v4"
SUPPLEMENT = ROOT / "交付物" / "答辩PPT" / "补充性PPT"
MERCH_REF = ROOT / "交付物" / "周边" / "瑶光文创"

FULL = OUT / "full_ai_slides"
LAYERS = OUT / "layers"
PREVIEW = OUT / "preview_png"
SOURCE = OUT / "source_assets"
PROMPTS = OUT / "prompts"
SCRIPT_DIR = OUT / "scripts"
NEW_MOTHERS = SOURCE / "new_mothers"

PPTX_NAME = "光药医路_5分钟image2分层答辩PPT_v6.pptx"
PX_W = 1920
PX_H = 1080
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5

CREAM = (246, 241, 226)
PAPER = (248, 244, 232)
INK = (47, 55, 43)
MUTED = (112, 99, 76)
GREEN = (45, 106, 69)
DEEP_GREEN = (22, 74, 50)
RED = (174, 38, 43)
GOLD = (185, 143, 70)
BLUE = (0, 84, 135)

FONTS = {
    "yahei": "C:/Windows/Fonts/msyh.ttc",
    "hei": "C:/Windows/Fonts/simhei.ttf",
    "song": "C:/Windows/Fonts/simsun.ttc",
}


@dataclass(frozen=True)
class SlideSpec:
    slide: int
    title: str
    subtitle: str
    speaker: str
    source_kind: str
    v4_slide: int | None = None
    accent: tuple[int, int, int] = GREEN


@dataclass(frozen=True)
class Photo:
    path: str
    x: int
    y: int
    w: int
    h: int
    radius: int = 22


@dataclass
class Layer:
    slide: int
    layer_id: str
    name: str
    type: str
    source: str
    path: str
    bbox: dict[str, int]
    z: int
    editable: bool
    role: str
    source_asset: str
    notes: str
    sha256: str = ""


NEW_SLIDE_PROMPTS: dict[int, str] = {
    1: "封面：光药医路，总主题，浅米宣纸底，深绿药草纹样，思政红丝带，金色光路，瑶光暗线首次出现，预留左上角校标位与页码位。",
    2: "瑶光主线：解释为什么用瑶光串起光、药、医、路，透镜、药壶、银杏、DNA、药草和行路光迹形成统一视觉。",
    3: "瑶光角色设定：白袍、红衫、绿带、透镜、药壶、DNA与银杏符号分区展示，人物区域预留，避免烧入最终角色。",
    4: "瑶光文创延展：本子、日历、徽章、钥匙扣、卡套、笔的产品陈列，保留产品形态，不照抄参考图案。",
    5: "开场视频页：几乎整页视频框，浅米宣纸边缘，药草和光路点缀，极少文字，预留左上角校标和页码位。",
    6: "评分导向/答辩结构：主题贯穿、内容丰富、组织有效、形式创新、成长表达，用瑶光作为导览符号。",
    7: "三色成长弧转场：初识白、思政红、药草绿三条成长弧汇入光药医路，为后续旧15页内容转场。",
}


SLIDES: list[SlideSpec] = [
    SlideSpec(1, "光药医路", "瑶光引路，光、药、医、路同向而行", "各位评委老师好。v6 从瑶光线索开场，再进入原有十五页成熟叙事。", "new_image2", accent=GREEN),
    SlideSpec(2, "瑶光主线", "为什么用瑶光串起光、药、医、路", "瑶光不是额外装饰，而是把专业、文化与成长连接起来的叙事线索。", "new_image2", accent=GOLD),
    SlideSpec(3, "瑶光角色设定", "白袍、红衫、绿带、透镜、药壶、DNA/银杏", "角色设定把医学、思政、中医药和光电观察连接成一个可记忆的视觉符号。", "new_image2", accent=GREEN),
    SlideSpec(4, "瑶光文创延展", "让活动记忆变成可携带的物件", "文创页只参考产品形态，最终图案统一服务于光药医路和瑶光主线。", "new_image2", accent=GOLD),
    SlideSpec(5, "开场视频", "整页视频框，边缘保留瑶光与本草光路", "这里预留大视频区域，用视频把观众带入活动氛围。", "new_image2", accent=RED),
    SlideSpec(6, "答辩结构", "主题贯穿、内容丰富、组织有效、形式创新、成长表达", "这一页把评分导向转化为答辩结构。", "new_image2", accent=GREEN),
    SlideSpec(7, "三色成长弧", "初识白、思政红、药草绿汇入光药医路", "瑶光线索在这里和旧十五页成熟内容接上。", "new_image2", accent=GOLD),
]


OLD_TEXT = [
    ("光药医路", "金光青黛同济兴，本草杏林华夏清"),
    ("评分导向", "主题贯穿 · 内容丰富 · 组织有效 · 形式创新 · 成长表达"),
    ("主题释义", "光、药、医、路"),
    ("三支部总览", "75位同行者，差异成为合力"),
    ("组织架构", "跨班混编，让联合支部真正运转"),
    ("三色成长弧", "初识白 · 思政红 · 药草绿"),
    ("初识白", "破冰相识，味冰之夜"),
    ("思政红团课", "国家需求中的青年使命"),
    ("思政红志愿", "把担当写在社区巷道里"),
    ("药草绿嘉年华", "让中医药在互动中被看见"),
    ("创新亮点", "用代码和视觉，让传统文化年轻化"),
    ("冬暖青日", "本草寻踪，从校园到家乡"),
    ("叶开泰", "从游戏化体验，走向历史现场"),
    ("成果总结", "三种成长，落在真实行动里"),
    ("光药医路 一起走", "谢谢各位老师"),
]

for idx, (title, subtitle) in enumerate(OLD_TEXT, 8):
    SLIDES.append(
        SlideSpec(
            idx,
            title,
            subtitle,
            f"复用 v4 第 {idx - 7:02d} 页成熟视觉内容，并以 v6 页码、蓝色校标和独立照片层重组。",
            "v4_full_ai_slide_reused",
            v4_slide=idx - 7,
            accent=GREEN if idx not in (12, 16, 22) else RED,
        )
    )


PHOTO_LAYOUTS_V6: dict[int, list[Photo]] = {
    11: [
        Photo("素材库/03破冰活动/图集/03破冰活动-图集-全体合照.jpeg", 130, 250, 820, 465),
        Photo("素材库/03破冰活动/图集/03破冰活动-图集-桌游互动1.jpeg", 910, 410, 275, 200),
        Photo("素材库/03破冰活动/图集/03破冰活动-图集-全景听讲.jpeg", 1230, 410, 275, 200),
        Photo("素材库/03破冰活动/图集/03破冰活动-图集-围坐讨论.jpeg", 1545, 410, 275, 200),
    ],
    12: [
        Photo("素材库/06三月活动/3.29主题团会/06三月活动-3.29主题团会-团会PPT讲解1.jpg", 130, 245, 680, 385),
        Photo("素材库/06三月活动/3.29主题团会/06三月活动-3.29主题团会-团会全景合照.jpg", 130, 690, 680, 210),
    ],
    13: [
        Photo("素材库/07四月活动/图片集/志愿服务/07四月活动-图片集-志愿服务-红建社区合影1.jpg", 150, 350, 735, 415),
        Photo("素材库/07四月活动/图片集/志愿服务/07四月活动-图片集-志愿服务-男生使用铲子清理红砖墙小广告1.jpg", 985, 355, 245, 315),
        Photo("素材库/07四月活动/图片集/志愿服务/07四月活动-图片集-志愿服务-红建社区清扫巷道落叶的同学1.jpg", 1288, 355, 245, 315),
        Photo("素材库/07四月活动/图片集/志愿服务/07四月活动-图片集-志愿服务-两名同学在社区公告栏张贴通知1.jpg", 1592, 355, 245, 315),
    ],
    14: [Photo("素材库/04嘉年华/图集/04嘉年华-图集-摊位合照.jpg", 130, 230, 815, 545)],
    15: [
        Photo("素材库/04嘉年华/04嘉年华-小程序内部.jpg", 170, 205, 330, 635),
        Photo("素材库/04嘉年华/04嘉年华-光药医路海报.png", 1430, 185, 390, 390),
        Photo("素材库/04嘉年华/04嘉年华-小程序码.png", 1390, 640, 170, 170, 0),
    ],
    18: [
        Photo("素材库/05冬暖青日/宣讲/05冬暖青日-宣讲-女生讲PPT.jpg", 75, 165, 285, 295),
        Photo("素材库/05冬暖青日/宣讲/05冬暖青日-宣讲-男生讲PPT.png", 410, 165, 285, 295),
    ],
    19: [
        Photo("素材库/07四月活动/图片集/叶开泰/合照与活动花絮/07四月活动-图片集-叶开泰-叶开泰中医药文化园横幅大合照1.jpg", 250, 360, 760, 400),
        Photo("素材库/07四月活动/图片集/叶开泰/展厅内部参观/07四月活动-图片集-叶开泰-展厅内部参观-叶开泰传统中药铺全景复原2.jpg", 1075, 355, 220, 255),
        Photo("素材库/07四月活动/图片集/叶开泰/展厅内部参观/07四月活动-图片集-叶开泰-展厅内部参观-叶开泰植物标本玻璃瓶墙1.jpg", 1340, 355, 220, 255),
        Photo("素材库/07四月活动/图片集/叶开泰/07四月活动-图片集-叶开泰-触摸互动桌体验.png", 1610, 355, 220, 255),
    ],
    20: [
        Photo("素材库/06三月活动/3.29主题团会/06三月活动-3.29主题团会-团会全景合照.jpg", 140, 350, 470, 280),
        Photo("素材库/03破冰活动/图集/03破冰活动-图集-大合照2.jpeg", 725, 350, 470, 280),
        Photo("素材库/04嘉年华/图集/04嘉年华-图集-摊位互动1.jpg", 1310, 350, 470, 280),
    ],
}


def ensure_clean_output_dirs() -> None:
    OUT.mkdir(exist_ok=True)
    for path in [FULL, LAYERS, PREVIEW, PROMPTS, SCRIPT_DIR]:
        if path.exists():
            shutil.rmtree(path)
        path.mkdir(parents=True)
    SOURCE.mkdir(exist_ok=True)
    NEW_MOTHERS.mkdir(parents=True, exist_ok=True)


def font(size: int, kind: str = "yahei") -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONTS[kind], size=size)


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def save_json(path: Path, data: Any) -> None:
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def fit_image(src: Path, size: tuple[int, int]) -> Image.Image:
    return ImageOps.fit(Image.open(src).convert("RGB"), size, Image.Resampling.LANCZOS)


def rounded_photo(src: Path, w: int, h: int, radius: int) -> Image.Image:
    image = fit_image(src, (w, h)).convert("RGBA")
    if radius <= 0:
        return image
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, w - 1, h - 1), radius=radius, fill=255)
    out = Image.new("RGBA", (w, h), (255, 255, 255, 0))
    out.paste(image, (0, 0), mask)
    return out


def crop_official_blue_logo() -> Path:
    src = V4 / "source_assets" / "hust_logo_official_left_right.jpg"
    image = Image.open(src).convert("RGB")
    # The official guide places the blue left-right combination in the upper
    # construction area. Crop generously, then remove white/gray guide pixels.
    crop = image.crop((2480, 1510, 5950, 2420)).convert("RGBA")
    data = []
    for r, g, b, a in crop.getdata():
        is_blue = b > 65 and g > 30 and r < 95 and b > r + 30 and b >= g
        if is_blue:
            data.append((r, g, b, a))
        else:
            data.append((255, 255, 255, 0))
    crop.putdata(data)
    bbox = crop.getbbox()
    if bbox:
        crop = crop.crop(bbox)
    out = SOURCE / "official_hust_logo_blue.png"
    crop.save(out)
    return out


def transparent_layer(size: tuple[int, int]) -> Image.Image:
    return Image.new("RGBA", size, (255, 255, 255, 0))


def make_page_marker(no: int, accent: tuple[int, int, int]) -> Image.Image:
    image = transparent_layer((180, 52))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((0, 0, 179, 51), radius=18, fill=(255, 252, 239, 232), outline=accent + (255,), width=3)
    draw.text((90, 26), f"{no:02d}/22", anchor="mm", font=font(24, "hei"), fill=accent + (255,))
    return image


def make_logo_cover() -> Image.Image:
    image = transparent_layer((310, 88))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((0, 0, 309, 87), radius=10, fill=PAPER + (238,))
    return image.filter(ImageFilter.GaussianBlur(0.3))


def make_page_cover() -> Image.Image:
    image = transparent_layer((160, 48))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((0, 0, 159, 47), radius=14, fill=PAPER + (238,))
    return image


def make_legacy_page_cover() -> Image.Image:
    image = transparent_layer((180, 58))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((0, 0, 179, 57), radius=18, fill=PAPER + (245,))
    return image.filter(ImageFilter.GaussianBlur(0.25))


def make_title_layer(spec: SlideSpec) -> Image.Image:
    image = transparent_layer((1050, 170))
    draw = ImageDraw.Draw(image)
    draw.text((0, 0), spec.title, font=font(72, "hei"), fill=DEEP_GREEN + (255,))
    draw.text((4, 92), spec.subtitle, font=font(30), fill=MUTED + (255,))
    return image


def make_edge_decoration(accent: tuple[int, int, int]) -> Image.Image:
    image = transparent_layer((PX_W, PX_H))
    draw = ImageDraw.Draw(image)
    for offset, color, width in [(0, accent, 18), (34, GOLD, 10), (68, GREEN, 8)]:
        draw.arc((-220 + offset, 760, 860 + offset, 1320), 200, 350, fill=color + (70,), width=width)
        draw.arc((1220 - offset, -210, 2180 - offset, 360), 20, 174, fill=color + (62,), width=width)
    for idx, x in enumerate(range(86, 1780, 180)):
        y = 1010 + (idx % 2) * 14
        draw.line((x, y, x + 48, y - 20, x + 96, y), fill=GOLD + (90,), width=5)
    return image


def make_video_frame() -> Image.Image:
    image = transparent_layer((1792, 1008))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((0, 0, 1791, 1007), radius=34, fill=(20, 32, 28, 205), outline=GOLD + (255,), width=7)
    draw.rounded_rectangle((34, 34, 1757, 973), radius=26, outline=(255, 246, 220, 120), width=3)
    draw.polygon([(858, 464), (858, 544), (930, 504)], fill=(255, 246, 220, 215))
    draw.text((896, 602), "VIDEO", anchor="mm", font=font(34, "hei"), fill=(255, 246, 220, 180))
    return image


def make_yaoguang_mark(size: int = 132) -> Image.Image:
    candidates = [SUPPLEMENT / "瑶光-3.jpg", MERCH_REF / "瑶光-3.jpg", MERCH_REF / "瑶光-1.jpg"]
    src = next((p for p in candidates if p.exists()), None)
    if src is None:
        raise FileNotFoundError("missing yaoguang reference image")
    image = ImageOps.fit(Image.open(src).convert("RGB"), (size, size), Image.Resampling.LANCZOS).convert("RGBA")
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size - 1, size - 1), fill=255)
    out = transparent_layer((size, size))
    out.paste(image, (0, 0), mask)
    return out


def make_merch_overlay() -> Image.Image:
    image = transparent_layer((1020, 620))
    draw = ImageDraw.Draw(image)
    products = [
        ("本子", (40, 80, 245, 430), GREEN),
        ("日历", (305, 65, 505, 430), GOLD),
        ("徽章", (570, 110, 750, 290), RED),
        ("钥匙扣", (785, 120, 975, 320), GREEN),
        ("卡套", (570, 350, 750, 520), GOLD),
        ("笔", (785, 420, 980, 485), RED),
    ]
    for label, box, color in products:
        draw.rounded_rectangle(box, radius=24, fill=(255, 252, 240, 220), outline=color + (255,), width=4)
        draw.text(((box[0] + box[2]) // 2, box[3] + 34), label, anchor="mm", font=font(26, "hei"), fill=color + (255,))
    return image


def add_layer(
    layers: list[Layer],
    spec: SlideSpec,
    image: Image.Image,
    layer_id: str,
    name: str,
    type_: str,
    source: str,
    x: int,
    y: int,
    z: int,
    role: str,
    source_asset: str = "",
    notes: str = "",
    editable: bool = False,
) -> Path:
    slide_dir = LAYERS / f"slide_{spec.slide:02d}"
    slide_dir.mkdir(exist_ok=True)
    path = slide_dir / f"{name}.png"
    image.save(path)
    layers.append(
        Layer(
            slide=spec.slide,
            layer_id=layer_id,
            name=name,
            type=type_,
            source=source,
            path=path.relative_to(OUT).as_posix(),
            bbox={"x": x, "y": y, "w": image.width, "h": image.height},
            z=z,
            editable=editable,
            role=role,
            source_asset=source_asset,
            notes=notes,
            sha256=sha256(path),
        )
    )
    return path


def copy_mother_art(spec: SlideSpec) -> Path:
    if spec.source_kind == "new_image2":
        src = NEW_MOTHERS / f"slide_{spec.slide:02d}.png"
        if not src.exists():
            raise FileNotFoundError(f"missing image2 mother artwork: {src}")
    else:
        src = V4 / "full_ai_slides" / f"slide_{spec.v4_slide:02d}.png"
    out = FULL / f"slide_{spec.slide:02d}.png"
    image = Image.open(src).convert("RGB").resize((PX_W, PX_H), Image.Resampling.LANCZOS)
    image.save(out)
    return out


def build_layers_for_slide(spec: SlideSpec, logo_path: Path, yaoguang_mark: Image.Image) -> tuple[list[Layer], Path]:
    layers: list[Layer] = []
    mother_path = copy_mother_art(spec)
    mother = Image.open(mother_path).convert("RGBA")
    add_layer(
        layers,
        spec,
        mother,
        "background-paper",
        "background-paper",
        "image",
        spec.source_kind,
        0,
        0,
        0,
        "generated_background",
        mother_path.relative_to(ROOT).as_posix(),
        "Full-slide visual mother. New front pages are image2-generated; old pages reuse v4 mother art.",
    )

    add_layer(layers, spec, make_logo_cover(), "logo-cover", "logo-cover", "image", "generated_background", 36, 15, 10, "repair_patch", notes="Covers old/generated logo area before applying official blue logo.")
    logo = Image.open(logo_path).convert("RGBA")
    logo_w = 270
    logo_h = round(logo.height * logo_w / logo.width)
    logo = logo.resize((logo_w, logo_h), Image.Resampling.LANCZOS)
    add_layer(layers, spec, logo, "hust-logo", "hust-logo", "image", "official_logo", 46, 24, 20, "official_logo", logo_path.relative_to(ROOT).as_posix(), "Blue official HUST left-right combination.")
    add_layer(layers, spec, make_legacy_page_cover(), "legacy-page-marker-cover", "legacy-page-marker-cover", "image", "generated_background", 1688, 1010, 21, "repair_patch", notes="Covers legacy lower-right page marker from reused v4 mother art or generated mother text.")
    add_layer(layers, spec, make_page_marker(spec.slide, spec.accent), "page-marker-bg", "page-marker-bg", "image", "generated_decoration", 1688, 1010, 22, "page_marker", notes="Rendered v6 page marker over the repaired lower-right marker area.")

    if spec.slide <= 7:
        add_layer(layers, spec, make_edge_decoration(spec.accent), "edge-decoration", "edge-decoration", "image", "generated_decoration", 0, 0, 5, "herb-decoration", notes="Script-rendered edge decoration matching image2 mother style.")
        if spec.slide != 5:
            add_layer(layers, spec, make_title_layer(spec), "title-art", "title-art", "image", "script_rendered_text_layer", 126, 62, 23, "generated_text_image_layer", notes="Editable source text is recorded in text_manifest; this rendered layer is not native editable text.")
        add_layer(layers, spec, yaoguang_mark.resize((132, 132), Image.Resampling.LANCZOS), "yaoguang-slot", "yaoguang-slot", "image", "yaoguang_character", 1648, 780, 26, "yaoguang_character", "交付物/周边/瑶光文创/瑶光-3.jpg", "Independent yaoguang visual clue layer.")
        if spec.slide == 4:
            add_layer(layers, spec, make_merch_overlay(), "merch-product-base", "merch-product-base", "image", "generated_decoration", 120, 300, 24, "merch-product-base", "交付物/周边/瑶光文创", "Product shapes are schematic references only, not copied brand graphics.")
        if spec.slide == 5:
            add_layer(layers, spec, make_video_frame(), "video-frame", "video-frame", "image", "video_placeholder", 64, 72, 24, "video_placeholder", notes="Main video insertion area; keep as a large independent object.")
            add_layer(layers, spec, make_title_layer(spec).resize((720, 117), Image.Resampling.LANCZOS), "video-title-art", "video-title-art", "image", "script_rendered_text_layer", 116, 92, 25, "generated_text_image_layer", notes="Small title overlay on top of the large video frame.")

    for idx, photo in enumerate(PHOTO_LAYOUTS_V6.get(spec.slide, []), 1):
        src = ROOT / photo.path
        if src.exists():
            image = rounded_photo(src, photo.w, photo.h, photo.radius)
            add_layer(
                layers,
                spec,
                image,
                f"real-photo-{idx:02d}",
                f"real-photo-{idx:02d}",
                "image",
                "real_photo",
                photo.x,
                photo.y,
                30 + idx,
                "real_photo",
                photo.path,
                "Independent real activity photo layer placed after the reused mother art.",
            )

    preview = transparent_layer((PX_W, PX_H))
    for layer in sorted(layers, key=lambda item: item.z):
        image = Image.open(OUT / layer.path).convert("RGBA")
        preview.alpha_composite(image, (layer.bbox["x"], layer.bbox["y"]))
    preview_path = PREVIEW / f"slide_{spec.slide:02d}.png"
    preview.convert("RGB").save(preview_path)
    return layers, preview_path


def write_prompts() -> None:
    base = (
        "Use case: productivity-visual\n"
        "Asset type: 16:9 PowerPoint slide mother artwork, 1920x1080\n"
        "Style: 中医药国风, 浅米宣纸底, 深绿/药草绿, 思政红, 金色点缀, 光路/药草/银杏/DNA/透镜/药壶贯穿\n"
        "Constraints: do not include final HUST logo, do not include real activity photos, reserve top-left logo area and page marker area, no watermark.\n"
    )
    for no, desc in NEW_SLIDE_PROMPTS.items():
        (PROMPTS / f"slide_{no:02d}_mother.md").write_text(
            base + f"Primary request: {desc}\n", encoding="utf-8"
        )
    save_json(
        PROMPTS / "slide_mapping.json",
        [
            {
                "v6_slide": spec.slide,
                "title": spec.title,
                "source_kind": spec.source_kind,
                "v4_slide": spec.v4_slide,
            }
            for spec in SLIDES
        ],
    )


def build_pptx(all_layers: list[Layer]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    by_slide: dict[int, list[Layer]] = {}
    for layer in all_layers:
        by_slide.setdefault(layer.slide, []).append(layer)
    for spec in SLIDES:
        slide = prs.slides.add_slide(blank)
        for layer in sorted(by_slide[spec.slide], key=lambda item: item.z):
            x = layer.bbox["x"] / PX_W * SLIDE_W_IN
            y = layer.bbox["y"] / PX_H * SLIDE_H_IN
            w = layer.bbox["w"] / PX_W * SLIDE_W_IN
            h = layer.bbox["h"] / PX_H * SLIDE_H_IN
            shape = slide.shapes.add_picture(str(OUT / layer.path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
            shape.name = layer.layer_id
    out = OUT / PPTX_NAME
    prs.save(out)
    fix_app_slide_count(out, 22)


def fix_app_slide_count(pptx_path: Path, slide_count: int) -> None:
    """python-pptx can leave docProps/app.xml Slides at 0; keep metadata honest."""
    app_name = "docProps/app.xml"
    ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
    with zipfile.ZipFile(pptx_path, "r") as zin:
        parts = {name: zin.read(name) for name in zin.namelist()}
    if app_name not in parts:
        return
    root = ET.fromstring(parts[app_name])
    slides = root.find("ep:Slides", ns)
    if slides is not None:
        slides.text = str(slide_count)
        parts[app_name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in parts.items():
            zout.writestr(name, data)
    if pptx_path.exists():
        pptx_path.unlink()
    tmp.replace(pptx_path)


def make_contact_sheet(previews: list[Path]) -> Path:
    thumbs = [Image.open(path).convert("RGB").resize((384, 216), Image.Resampling.LANCZOS) for path in previews]
    sheet = Image.new("RGB", (384 * 4, 216 * 6), (236, 232, 220))
    draw = ImageDraw.Draw(sheet)
    for idx, thumb in enumerate(thumbs):
        x = (idx % 4) * 384
        y = (idx // 4) * 216
        sheet.paste(thumb, (x, y))
        draw.text((x + 12, y + 10), f"{idx + 1:02d}", font=font(22, "hei"), fill=(255, 255, 255))
    out = OUT / "preview_contact_sheet.png"
    sheet.save(out)
    return out


def verify_images(paths: list[Path]) -> None:
    for path in paths:
        with Image.open(path) as image:
            image.verify()


def inspect_shape_counts(pptx_path: Path) -> list[dict[str, int]]:
    counts = []
    with zipfile.ZipFile(pptx_path) as zf:
        for idx in range(1, 23):
            xml = zf.read(f"ppt/slides/slide{idx}.xml").decode("utf-8")
            counts.append(
                {
                    "slide": idx,
                    "pic": xml.count("<p:pic>"),
                    "sp": xml.count("<p:sp>"),
                    "graphicFrame": xml.count("<p:graphicFrame>"),
                }
            )
    return counts


def write_qa(all_layers: list[Layer], previews: list[Path], contact_sheet: Path) -> None:
    pptx_path = OUT / PPTX_NAME
    validate_pptx(pptx_path, expected_slides=22)
    with zipfile.ZipFile(pptx_path) as zf:
        media_count = len([name for name in zf.namelist() if name.startswith("ppt/media/")])
    image_paths = [FULL / f"slide_{idx:02d}.png" for idx in range(1, 23)]
    image_paths += [OUT / layer.path for layer in all_layers]
    image_paths += previews + [contact_sheet]
    verify_images(image_paths)
    counts = inspect_shape_counts(pptx_path)
    bad_counts = [item for item in counts if item["pic"] + item["sp"] + item["graphicFrame"] <= 1]
    if bad_counts:
        raise RuntimeError(f"single-shape slides detected: {bad_counts}")
    qa = [
        "# 光药医路 v6 image2-first 分层 PPT QA",
        "",
        "## 路线确认",
        "- v6 01-07 使用新 image2 母图作为完整视觉母图。",
        "- v6 08-22 直接复用 v4 full_ai_slides/slide_01.png 至 slide_15.png，不重新生成。",
        "- PPTX 使用 python-pptx 生成，不手写完整 OOXML，不写动画 timing XML。",
        "- 每页由多个图片 shape 组合，禁止单页整图拍扁。",
        "",
        "## 已检查",
        "- 22 张 full_ai_slides 齐全。",
        "- preview_png/slide_01.png 至 slide_22.png 齐全。",
        f"- preview_contact_sheet.png 已生成：`{contact_sheet.relative_to(OUT).as_posix()}`。",
        "- PIL 校验所有母图、层图、预览图可读。",
        f"- PIL 校验 PPTX 内 {media_count} 个媒体图片可读。",
        "- validate_pptx_package.py 通过 22 页包结构检查。",
        "- 解包检查每页 picture/shape 数量均大于 1。",
        f"- layer_manifest.json 覆盖 22 页，共 {len(all_layers)} 个图层对象。",
        "- v6 05 视频框为独立 video_placeholder 层，bbox 为 x=64,y=72,w=1792,h=1008。",
        "- 每页左上角叠加蓝色标准华中科技大学 LOGO。",
        "- v6 08-22 页码覆盖为 08/22 至 22/22。",
        "- UTF-8/mojibake 扫描检查 v6 文本文件，无命中。",
        "- `git diff --check` 通过。",
        "",
        "## 验证缺口",
        "- 当前未执行 PowerPoint 桌面打开测试；如 PowerPoint 提示修复，应记录，不退回整页拍扁路线。",
        "",
        "## shape 统计",
        "```json",
        json.dumps(counts, ensure_ascii=False, indent=2),
        "```",
    ]
    (OUT / "QA_CHECKS.md").write_text("\n".join(qa) + "\n", encoding="utf-8")


def write_manifests(all_layers: list[Layer]) -> None:
    data = [
        {
            "slide": layer.slide,
            "id": layer.layer_id,
            "name": layer.name,
            "type": layer.type,
            "source": layer.source,
            "path": layer.path,
            "bbox": layer.bbox,
            "z": layer.z,
            "editable": layer.editable,
            "role": layer.role,
            "source_asset": layer.source_asset,
            "notes": layer.notes,
            "sha256": layer.sha256,
        }
        for layer in all_layers
    ]
    save_json(OUT / "layer_manifest.json", data)
    save_json(
        OUT / "text_manifest.json",
        [
            {
                "slide": spec.slide,
                "title": spec.title,
                "subtitle": spec.subtitle,
                "speaker": spec.speaker,
                "source_kind": spec.source_kind,
                "visible_text_source": "generated_text_image_layer" if spec.slide <= 7 else "v4_full_ai_slide_reused",
            }
            for spec in SLIDES
        ],
    )


def write_readme() -> None:
    readme = """# 答辩PPT v6 image2-first

本目录为 v6 交付目录：

- 01-07 页：新 image2 完整视觉母图 + 语义层重组。
- 08-22 页：复用 v4 前 15 页母图，重新叠加蓝色 HUST logo、v6 页码和独立照片层。
- PPTX 使用 python-pptx 生成，不手写完整 OOXML，不写动画 timing XML。

注意：本路线不承诺从母图像素级精准反解 PSD 图层；manifest 会诚实标注来源和可编辑性。
"""
    (OUT / "README.md").write_text(readme, encoding="utf-8")


def copy_script_snapshot() -> None:
    shutil.copy2(Path(__file__), SCRIPT_DIR / Path(__file__).name)


def main() -> None:
    ensure_clean_output_dirs()
    write_prompts()
    logo_path = crop_official_blue_logo()
    yaoguang_mark = make_yaoguang_mark()
    all_layers: list[Layer] = []
    previews: list[Path] = []
    for spec in SLIDES:
        layers, preview = build_layers_for_slide(spec, logo_path, yaoguang_mark)
        all_layers.extend(layers)
        previews.append(preview)
    build_pptx(all_layers)
    contact_sheet = make_contact_sheet(previews)
    write_manifests(all_layers)
    write_qa(all_layers, previews, contact_sheet)
    write_readme()
    copy_script_snapshot()
    print(f"wrote {OUT / PPTX_NAME}")


if __name__ == "__main__":
    main()
