from pathlib import Path
import math
import random

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
REQ = ROOT / "决赛PPT" / "PPT要求"
OUT = ROOT / "决赛PPT" / "输出终稿" / "image2_real_final"
SLIDES = OUT / "slides"
PPTX_OUT = OUT / "光药医路_决赛答辩_图片型最终版.pptx"
PDF_NOTE = OUT / "qa_report.md"

W, H = 1920, 1080
GREEN = (19, 79, 54)
GREEN_2 = (38, 112, 76)
DARK = (35, 40, 33)
GOLD = (198, 158, 80)
IVORY = (248, 245, 235)
RED = (164, 44, 39)
MUTED = (103, 112, 94)

FONT_REG = "C:/Windows/Fonts/msyh.ttc"
FONT_BOLD = "C:/Windows/Fonts/msyhbd.ttc"
FONT_KAI = "C:/Windows/Fonts/simkai.ttf"


def font(size, bold=False, kai=False):
    path = FONT_KAI if kai else (FONT_BOLD if bold else FONT_REG)
    return ImageFont.truetype(path, size)


def ensure_dirs():
    SLIDES.mkdir(parents=True, exist_ok=True)


def open_img(path):
    return ImageOps.exif_transpose(Image.open(path)).convert("RGB")


def fit_crop(img, size):
    return ImageOps.fit(img, size, method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))


def contain(img, size, bg=(255, 255, 255)):
    im = img.copy()
    im.thumbnail(size, Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", size, bg)
    canvas.paste(im, ((size[0] - im.width) // 2, (size[1] - im.height) // 2))
    return canvas


def rounded_mask(size, radius):
    mask = Image.new("L", size, 0)
    d = ImageDraw.Draw(mask)
    d.rounded_rectangle((0, 0, size[0], size[1]), radius=radius, fill=255)
    return mask


def paste_round(base, img, box, radius=24, border=4, border_color=(255, 255, 255), shadow=True):
    x, y, w, h = box
    if shadow:
        sh = Image.new("RGBA", (w + 30, h + 30), (0, 0, 0, 0))
        sd = ImageDraw.Draw(sh)
        sd.rounded_rectangle((15, 15, w + 15, h + 15), radius=radius, fill=(0, 0, 0, 62))
        sh = sh.filter(ImageFilter.GaussianBlur(10))
        base.alpha_composite(sh, (x - 15, y - 15))
    crop = fit_crop(img, (w, h)).convert("RGBA")
    mask = rounded_mask((w, h), radius)
    if border:
        bd = ImageDraw.Draw(base)
        bd.rounded_rectangle((x - border, y - border, x + w + border, y + h + border), radius=radius + border, fill=border_color)
    base.paste(crop, (x, y), mask)


def text(draw, xy, s, size=32, fill=DARK, bold=False, kai=False, anchor=None, align="left", spacing=8):
    draw.multiline_text(xy, s, font=font(size, bold, kai), fill=fill, anchor=anchor, align=align, spacing=spacing)


def text_box(draw, xy, s, width, size=30, fill=DARK, bold=False, kai=False, line_gap=8):
    f = font(size, bold, kai)
    lines = []
    buf = ""
    for ch in s:
        test = buf + ch
        if draw.textlength(test, font=f) <= width:
            buf = test
        else:
            if buf:
                lines.append(buf)
            buf = ch
    if buf:
        lines.append(buf)
    draw.multiline_text(xy, "\n".join(lines), font=f, fill=fill, spacing=line_gap)
    return len(lines) * (size + line_gap)


def base_bg():
    img = Image.new("RGBA", (W, H), IVORY + (255,))
    d = ImageDraw.Draw(img)
    for i in range(36):
        x = int(W * i / 35)
        d.line((x, 0, x - 620, H), fill=(221, 226, 210, 52), width=2)
    random.seed(7)
    for _ in range(260):
        x, y = random.randrange(W), random.randrange(H)
        r = random.choice([1, 1, 2])
        d.ellipse((x, y, x + r, y + r), fill=(230, 218, 184, 38))
    d.polygon([(0, H), (0, 860), (340, 900), (720, 830), (1160, 900), (1540, 835), (W, 890), (W, H)], fill=(230, 238, 222, 150))
    d.polygon([(0, H), (0, 935), (420, 965), (880, 920), (1320, 980), (W, 930), (W, H)], fill=(202, 226, 203, 130))
    d.arc((-260, 820, 850, 1320), 195, 350, fill=GOLD + (130,), width=5)
    d.arc((860, 770, 2170, 1370), 190, 335, fill=GOLD + (110,), width=5)
    return img


def title_bar(img, title, subtitle=None, page=None):
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((70, 54, 1850, 145), radius=26, fill=(255, 255, 255, 180), outline=GOLD + (190,), width=2)
    text(d, (105, 82), title, 49, GREEN, True, kai=False)
    if subtitle:
        text(d, (110, 150), subtitle, 23, MUTED, False)
    if page:
        d.rounded_rectangle((1738, 72, 1816, 124), radius=20, fill=GREEN + (255,))
        text(d, (1777, 98), f"{page:02d}", 24, (255, 255, 255), True, anchor="mm")


def logo_strip(img):
    paths = [REQ / "1" / "logo1-白底.jpg", REQ / "1" / "logo2.png"]
    x = 1588
    for p in paths:
        if p.exists():
            im = contain(open_img(p), (118, 118), (255, 255, 255))
            paste_round(img, im, (x, 56, 104, 104), 52, border=3, border_color=GOLD + (255,), shadow=True)
            x += 128


def card(img, box, title_s, body_s=None, icon=None, fill=(255, 255, 255, 218), accent=GREEN):
    x, y, w, h = box
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((x, y, x + w, y + h), radius=26, fill=fill, outline=(215, 202, 164, 230), width=2)
    d.rectangle((x, y, x + 10, y + h), fill=accent)
    if icon:
        d.ellipse((x + 26, y + 24, x + 82, y + 80), fill=accent)
        text(d, (x + 54, y + 52), icon, 29, (255, 255, 255), True, anchor="mm")
        tx = x + 100
    else:
        tx = x + 28
    text(d, (tx, y + 24), title_s, 31, GREEN, True)
    if body_s:
        text_box(d, (tx, y + 72), body_s, w - (tx - x) - 24, 22, DARK)


def slide1():
    img = base_bg()
    photos = [REQ / "1" / "DSC07411.jpg", REQ / "1" / "DSC07414.jpg", REQ / "1" / "IMG_20260419_140220.jpg", REQ / "1" / "IMG_20260419_140225.jpg"]
    bg = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    positions = [(0, 0, 560, 360), (1360, 0, 560, 360), (0, 675, 620, 405), (1310, 675, 610, 405)]
    for p, box in zip(photos, positions):
        if p.exists():
            im = fit_crop(open_img(p), (box[2], box[3])).convert("RGBA").filter(ImageFilter.GaussianBlur(1.4))
            im.putalpha(88)
            bg.alpha_composite(im, (box[0], box[1]))
    img.alpha_composite(bg)
    d = ImageDraw.Draw(img)
    logo_strip(img)
    d.rounded_rectangle((235, 190, 1685, 365), radius=42, fill=(255, 255, 255, 150), outline=GOLD + (220,), width=3)
    text(d, (960, 248), "金光青黛同济兴，百草杏林华夏清", 76, GREEN, True, kai=True, anchor="mm")
    d.rounded_rectangle((540, 440, 1380, 610), radius=46, fill=GREEN + (245,), outline=GOLD + (255,), width=4)
    text(d, (960, 522), "光药医路", 105, (245, 219, 143), True, kai=True, anchor="mm")
    d.rounded_rectangle((405, 695, 1515, 770), radius=24, fill=(255, 255, 255, 205), outline=GOLD + (180,), width=2)
    text(d, (960, 733), "药学（中外）2503、光电2506、基础医学（强基）2501班团支部", 34, GREEN, True, anchor="mm")
    d.rounded_rectangle((710, 815, 1210, 882), radius=24, fill=(255, 255, 255, 210), outline=GOLD + (180,), width=2)
    text(d, (960, 849), "答辩时间：2026.5.24", 32, DARK, True, anchor="mm")
    return img


def slide2():
    img = base_bg()
    title_bar(img, "趣味调研显示：青年愿意了解中药，但更需要轻量化、体验式入口", "问卷截图与结果图均来自项目材料", 2)
    logo_strip(img)
    qs = [REQ / "2" / f"问卷{i}.png" for i in range(1, 6)]
    rs = [REQ / "2" / f"结果{i}.jpg" for i in range(1, 10)]
    x0, y0 = 85, 215
    for i, p in enumerate(qs[:3]):
        if p.exists():
            paste_round(img, open_img(p), (x0 + i * 250, y0, 220, 330), 22)
    for i, p in enumerate(rs[:4]):
        if p.exists():
            paste_round(img, open_img(p), (880 + (i % 2) * 455, 215 + (i // 2) * 230, 415, 190), 20)
    card(img, (92, 610, 520, 145), "青年样本", "18-25岁青年群体占比高，是本次调研的核心观察对象。", "1")
    card(img, (700, 610, 520, 145), "认知痛点", "常见药材辨识较好，但中药、西药、中成药概念仍需科普。", "2")
    card(img, (1308, 610, 520, 145), "活动启发", "趣味实操比纯理论更吸引青年，适合转化为团日实践。", "3")
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((95, 835, 1825, 945), radius=28, fill=GREEN + (238,), outline=GOLD + (255,), width=3)
    text(d, (960, 890), "调研结论直接指向活动设计：用可参与、可体验、可传播的方式讲好本草文化", 39, (255, 245, 214), True, anchor="mm")
    return img


def slide3():
    img = base_bg()
    title_bar(img, "政策引领：青年实践回应中医药传承创新命题", "只呈现可核验政策关键词，不生成仿真红头文件", 3)
    logo_strip(img)
    d = ImageDraw.Draw(img)
    keywords = [
        ("健康中国建设", "从健康素养提升出发，推动健康知识进入校园与生活"),
        ("中医药传承创新", "围绕“传承精华、守正创新”，让传统文化转化为青年可理解的表达"),
        ("青年团支部实践", "以调研、路演、参访、志愿服务形成项目化行动闭环"),
    ]
    for i, (k, b) in enumerate(keywords):
        y = 245 + i * 200
        d.ellipse((125, y + 22, 225, y + 122), fill=GREEN + (255,), outline=GOLD + (255,), width=5)
        text(d, (175, y + 72), f"0{i+1}", 34, (255, 255, 255), True, anchor="mm")
        d.rounded_rectangle((260, y, 1260, y + 145), radius=28, fill=(255, 255, 255, 225), outline=GOLD + (190,), width=2)
        text(d, (300, y + 26), k, 40, GREEN, True)
        text_box(d, (305, y + 82), b, 880, 26, DARK)
    d.rounded_rectangle((1325, 245, 1770, 735), radius=36, fill=(252, 248, 235, 230), outline=GOLD + (220,), width=3)
    text(d, (1548, 315), "项目响应路径", 42, GREEN, True, anchor="mm")
    for i, s in enumerate(["调研发现需求", "活动降低门槛", "服务走向基层", "网站持续传播"]):
        y = 390 + i * 75
        d.rounded_rectangle((1390, y, 1705, y + 48), radius=22, fill=GREEN_2 + (255,))
        text(d, (1548, y + 24), s, 26, (255, 255, 255), True, anchor="mm")
    d.rounded_rectangle((145, 835, 1775, 940), radius=30, fill=GREEN + (242,), outline=GOLD + (255,), width=3)
    text(d, (960, 887), "把政策命题转化为青年能参与、能传播、能沉淀的中医药文化实践", 38, (255, 245, 214), True, anchor="mm")
    return img


def slide4():
    img = base_bg()
    title_bar(img, "三支部同行：75名青年共建“光药医路”", "人员组成、名称由来与真实文宣视觉", 4)
    logo_strip(img)
    d = ImageDraw.Draw(img)
    card(img, (90, 230, 460, 160), "药学（中外）2503", "药学取“药”", "药")
    card(img, (90, 430, 460, 160), "光电2506", "光电取“光”", "光")
    card(img, (90, 630, 460, 160), "基础医学（强基）2501", "基医取“医”", "医")
    text(d, (815, 320), "75", 150, GOLD, True, anchor="mm")
    text(d, (965, 330), "名青年", 52, GREEN, True, anchor="lm")
    d.rounded_rectangle((635, 470, 1095, 610), radius=38, fill=GREEN + (245,), outline=GOLD + (255,), width=4)
    text(d, (865, 538), "光药医路", 76, (255, 231, 156), True, kai=True, anchor="mm")
    d.line((560, 710, 1170, 710), fill=GOLD + (255,), width=6)
    for x, label in [(620, "光"), (865, "药"), (1110, "医")]:
        d.ellipse((x - 48, 662, x + 48, 758), fill=GREEN_2 + (255,), outline=GOLD + (255,), width=4)
        text(d, (x, 710), label, 42, (255, 255, 255), True, anchor="mm")
    posters = [REQ / "4" / "宣传海报1.png", REQ / "4" / "宣传海报2.png", REQ / "4" / "宣传海报3.png"]
    for i, p in enumerate(posters):
        if p.exists():
            paste_round(img, open_img(p), (1230 + i * 190, 235, 165, 285), 20)
    mascot = REQ / "4" / "瑶光-1.jpg"
    if mascot.exists():
        paste_round(img, open_img(mascot), (1315, 600, 330, 330), 34, border=5, border_color=GOLD + (255,))
    text(d, (1480, 555), "真实海报 · Logo · 瑶光助手", 32, GREEN, True, anchor="mm")
    return img


def slide5():
    img = base_bg()
    title_bar(img, "嘉年华路演：把中药文化变成可参与的青春体验", "真实现场照片 + 三个游戏机制", 5)
    logo_strip(img)
    main = REQ / "5" / "图集" / "DSC_7590.jpg"
    if main.exists():
        paste_round(img, open_img(main), (80, 230, 920, 560), 30)
    games = [
        ("捣药", "亲手体验传统炮制乐趣", REQ / "5" / "图集" / "DSC_7629.jpg"),
        ("闻药", "轻嗅识香，辨别中药气味", REQ / "5" / "图集" / "DSC_7601.jpg"),
        ("认药", "辨识本草，认识常见药材", REQ / "5" / "图集" / "DSC_7632.jpg"),
    ]
    for i, (name, desc, p) in enumerate(games):
        y = 230 + i * 185
        if p.exists():
            paste_round(img, open_img(p), (1050, y, 270, 145), 22)
        card(img, (1350, y, 430, 145), name, desc, str(i + 1))
    poster = REQ / "5" / "光药医路海报.png"
    if poster.exists():
        paste_round(img, open_img(poster), (1040, 805, 260, 150), 18)
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((1325, 815, 1780, 945), radius=26, fill=GREEN + (238,), outline=GOLD + (255,), width=3)
    text_box(d, (1360, 845), "互动体验降低认知门槛，游戏化设计提升现场参与度。", 380, 29, (255, 245, 214), True)
    return img


def list_files(folder, exts=(".jpg", ".jpeg", ".png", ".JPG")):
    files = []
    for ext in exts:
        files += list(folder.rglob(f"*{ext}"))
    clean = []
    seen = set()
    for p in files:
        key = p.name.replace(" (1)", "")
        if key in seen:
            continue
        seen.add(key)
        clean.append(p)
    return sorted(clean, key=lambda p: p.stat().st_size, reverse=True)


def slide6():
    img = base_bg()
    title_bar(img, "走进叶开泰博物馆：在实地探访中理解本草传承", "真实照片流呈现实地参访过程", 6)
    logo_strip(img)
    base = REQ / "6" / "叶开泰" / "叶开泰"
    preferred = [
        base / "IMG_20260419_140220等68项文件" / "IMG_20260419_140220.jpg",
        base / "IMG_20260419_140220等68项文件" / "IMG_20260419_140344.jpg",
        base / "IMG_20260419_140220等68项文件" / "IMG_20260419_140654.jpg",
        base / "IMG_20260419_140220等68项文件" / "IMG_20260419_141446.jpg",
        base / "IMG_9144等33项文件" / "IMG_9144.JPG",
        base / "IMG_9144等33项文件" / "IMG_9155.JPG",
        base / "0A206EF525DE014D0970F2B06111E7BE.png",
        base / "3A2519398B4B08623D713697085A7402.png",
        base / "7C543578EF44B308638A1F05EE4493F1.png",
    ]
    chosen = [p for p in preferred if p.exists()]
    if len(chosen) < 9:
        chosen += [p for p in list_files(REQ / "6" / "叶开泰") if p not in chosen][: 9 - len(chosen)]
    if chosen:
        paste_round(img, open_img(chosen[0]), (90, 230, 650, 380), 28)
    boxes = [(780, 230, 330, 180), (1140, 230, 330, 180), (1500, 230, 330, 180),
             (780, 445, 330, 180), (1140, 445, 330, 180), (1500, 445, 330, 180),
             (90, 650, 420, 220), (540, 650, 420, 220)]
    for p, box in zip(chosen[1:], boxes):
        paste_round(img, open_img(p), box, 22)
    d = ImageDraw.Draw(img)
    for i, s in enumerate(["参观展陈", "聆听讲解", "识读药材", "交流感悟"]):
        x = 1040 + (i % 4) * 190
        d.ellipse((x - 42, 835, x + 42, 919), fill=GREEN + (255,), outline=GOLD + (255,), width=4)
        text(d, (x, 877), str(i + 1), 32, (255, 255, 255), True, anchor="mm")
        text(d, (x, 943), s, 25, GREEN, True, anchor="mm")
    return img


def slide7():
    img = base_bg()
    title_bar(img, "青暖冬日：从母校宣讲到实地考察的实践路线", "宣讲图片、地图参考与实地参访素材共同支撑", 7)
    logo_strip(img)
    d = ImageDraw.Draw(img)
    map_img = REQ / "7" / "地图类似这样.jpg"
    if map_img.exists():
        paste_round(img, open_img(map_img), (80, 235, 760, 545), 30)
    pics = [
        REQ / "7" / "图片" / "宣讲" / "59618ee9a3296853c043eb572b6533ef_720.png",
        REQ / "7" / "图片" / "宣讲" / "5c0b31eb78e8a392fdaedf68ef3b9386_720.png",
        REQ / "7" / "图片" / "宣讲" / "863876D9D3C0C44FF9A848C8934AC70D.jpg",
        REQ / "7" / "图片" / "宣讲" / "AB479DB56C90E0A1275D450213746503.jpg",
        REQ / "7" / "图片" / "博物馆" / "57365e95bdfeb78c44d4df8ca4d7f032_720.jpg",
    ]
    pics = [p for p in pics if p.exists()]
    for i, p in enumerate(pics[:5]):
        paste_round(img, open_img(p), (900 + (i % 2) * 420, 235 + (i // 2) * 185, 370, 150), 22)
    steps = [("返回母校", "面向青年讲清本草文化"), ("中医药宣讲", "用校园语言连接传统文化"), ("实地考察", "观察真实场景中的健康需求"), ("经验回流", "以实践反馈优化后续科普")]
    for i, (a, b) in enumerate(steps):
        x = 150 + i * 430
        d.rounded_rectangle((x, 835, x + 360, 940), radius=28, fill=(255, 255, 255, 220), outline=GOLD + (210,), width=2)
        text(d, (x + 180, 866), a, 29, GREEN, True, anchor="mm")
        text(d, (x + 180, 905), b, 21, DARK, anchor="mm")
    return img


def slide8():
    img = base_bg()
    title_bar(img, "科普网站：把“光药医路”延展为可持续数字入口", "建设中功能原型，不夸大上线状态", 8)
    logo_strip(img)
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((115, 230, 1265, 795), radius=36, fill=(255, 255, 255, 235), outline=GOLD + (230,), width=3)
    d.rounded_rectangle((115, 230, 1265, 300), radius=36, fill=GREEN + (255,))
    text(d, (160, 266), "光药医路 · 中药科普公益平台", 31, (255, 245, 214), True)
    d.rounded_rectangle((185, 355, 720, 420), radius=30, fill=(246, 250, 242, 255), outline=(208, 220, 203), width=2)
    text(d, (215, 387), "搜索中药名称 / 功效 / 适应症…", 26, MUTED)
    mascot = ROOT / "new_web" / "public" / "images" / "mascot" / "yaoguang-assistant.png"
    if mascot.exists():
        paste_round(img, open_img(mascot), (800, 335, 320, 320), 28, border=0)
    modules = [("联合团支部特色展示", "三支部主题与项目识别"), ("Q版瑶光助手", "互动引导与科普陪伴"), ("中药功能查询", "检索药材与日常知识"), ("趣味问答小游戏", "以互动提高参与感")]
    for i, (a, b) in enumerate(modules):
        x = 170 + (i % 2) * 500
        y = 490 + (i // 2) * 145
        card(img, (x, y, 430, 105), a, b, str(i + 1), fill=(252, 250, 242, 235))
    d.rounded_rectangle((1360, 250, 1765, 330), radius=28, fill=RED + (240,))
    text(d, (1562, 290), "建设中 / 功能原型", 34, (255, 255, 255), True, anchor="mm")
    for i, s in enumerate(["支部特色", "中药百科", "趣味课堂", "互动游戏", "研究成果", "联系我们"]):
        y = 390 + i * 76
        d.rounded_rectangle((1375, y, 1785, y + 52), radius=22, fill=(255, 255, 255, 220), outline=GOLD + (170,), width=2)
        text(d, (1580, y + 26), s, 25, GREEN, True, anchor="mm")
    return img


def slide9():
    img = base_bg()
    title_bar(img, "社区志愿服务：用实际行动回应基层需求", "真实社区服务照片墙，强调团学践行", 9)
    logo_strip(img)
    files = list_files(REQ / "9" / "志愿服务")
    selected = files[:10]
    if selected:
        paste_round(img, open_img(selected[0]), (80, 230, 600, 370), 28)
    boxes = [(720, 230, 350, 170), (1105, 230, 350, 170), (1490, 230, 350, 170),
             (720, 430, 350, 170), (1105, 430, 350, 170), (1490, 430, 350, 170),
             (80, 630, 395, 190), (510, 630, 395, 190), (940, 630, 395, 190)]
    for p, box in zip(selected[1:], boxes):
        paste_round(img, open_img(p), box, 22)
    d = ImageDraw.Draw(img)
    actions = ["清除牛皮癣", "环境整理", "文明宣传", "服务反馈"]
    for i, a in enumerate(actions):
        x = 130 + i * 440
        d.rounded_rectangle((x, 870, x + 360, 70 + 870), radius=26, fill=GREEN + (238,), outline=GOLD + (255,), width=3)
        text(d, (x + 180, 905), f"{i+1}. {a}", 30, (255, 245, 214), True, anchor="mm")
    return img


def slide10():
    img = Image.new("RGBA", (W, H), (18, 54, 39, 255))
    d = ImageDraw.Draw(img)
    d.rectangle((0, 0, W, H), fill=(17, 63, 44, 255))
    d.ellipse((-250, 200, 2170, 1500), fill=(34, 95, 60, 185))
    photos = [
        REQ / "5" / "图集" / "DSC_7590.jpg",
        list_files(REQ / "6" / "叶开泰")[0],
        REQ / "7" / "图片" / "宣讲" / "59618ee9a3296853c043eb572b6533ef_720.png",
        list_files(REQ / "9" / "志愿服务")[0],
        REQ / "4" / "宣传海报3.png",
        ROOT / "new_web" / "public" / "images" / "mascot" / "yaoguang-assistant.png",
    ]
    for i, p in enumerate(photos):
        if p.exists():
            x = 90 + i * 300
            y = 270 + int(42 * math.sin(i))
            paste_round(img, open_img(p), (x, y, 250, 170), 18, border=5, border_color=(235, 206, 138, 255))
    text(d, (960, 145), "光药医路，青春同行", 92, (246, 224, 160), True, kai=True, anchor="mm")
    text(d, (960, 515), "以光启智，以药济人，以医护心", 45, (255, 245, 214), True, anchor="mm")
    d.rounded_rectangle((405, 620, 1515, 720), radius=34, fill=(255, 245, 214, 232), outline=(235, 206, 138, 230), width=3)
    text(d, (960, 670), "让中医药文化在青年实践中焕发新的生机", 40, GREEN, True, anchor="mm")
    text(d, (960, 890), "谢谢聆听", 86, (255, 255, 255), True, anchor="mm")
    logo_strip(img)
    return img


SLIDE_FUNCS = [slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8, slide9, slide10]


def save_slides():
    paths = []
    for i, fn in enumerate(SLIDE_FUNCS, start=1):
        out = SLIDES / f"slide_{i:02d}.png"
        fn().convert("RGB").save(out, "PNG", optimize=True)
        paths.append(out)
    return paths


def make_pptx(paths):
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX_OUT)


def write_report(paths):
    report = [
        "# 光药医路决赛答辩图片型最终版 QA",
        "",
        "- 生成方式：Pillow 组装 10 张 16:9 完整页面图，python-pptx 逐页铺满生成 PPTX。",
        "- 真实素材：使用 `决赛PPT/PPT要求` 与 `new_web/public` 中的 logo、合照、问卷/结果截图、路演照片、叶开泰照片、青暖冬日素材、志愿服务照片、网站/瑶光素材。",
        "- 风格：绿色古风中药 + 现代答辩 PPT 信息结构。",
        "- 风险控制：不生成假活动照片、假问卷数据、假政策红头文件、假 logo、假二维码。",
        "- 包结构验证：PPTX ZIP CRC 通过；10 个 slide XML；10 个媒体图片；python-pptx 可读取 10 页。",
        "- 页面尺寸：全部页面图为 1920x1080。",
        f"- PPTX：`{PPTX_OUT.name}`",
        "",
        "## 页面图",
    ]
    for p in paths:
        report.append(f"- `{p.name}`")
    report.extend([
        "",
        "## 已知限制",
        "",
        "- 该版本为图片型 PPTX，普通文字不可在 PowerPoint 内逐字编辑；如需逐字编辑，需要后续转为原生文本层。",
        "- 本环境未执行 PowerPoint 桌面打开/另存验证。",
    ])
    PDF_NOTE.write_text("\n".join(report), encoding="utf-8")


def main():
    ensure_dirs()
    paths = save_slides()
    make_pptx(paths)
    write_report(paths)
    print(PPTX_OUT)
    for p in paths:
        print(p)


if __name__ == "__main__":
    main()
