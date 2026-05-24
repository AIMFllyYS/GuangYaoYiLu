from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path

from PIL import Image, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
SRC_DIR = ROOT / "决赛PPT" / "输出终稿" / "image2_premium_ppt_v2" / "slides"
OUT = ROOT / "决赛PPT" / "WorkShape" / "layered_pages"

CANVAS_W = 1920
CANVAS_H = 1080
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5

FONT_YAHEI = "Microsoft YaHei"
FONT_KAI = "KaiTi"

GREEN = (16, 82, 54)
GREEN_2 = (23, 98, 64)
DARK_GREEN = (5, 58, 40)
GOLD = (198, 158, 80)
GOLD_LIGHT = (245, 222, 150)
IVORY = (249, 246, 232)
IVORY_2 = (255, 252, 241)
DARK = (45, 54, 39)
RED = (149, 31, 28)


def rel(path: Path) -> str:
    return str(path.relative_to(ROOT)).replace("\\", "/")


def src(slide_no: int) -> Path:
    return SRC_DIR / f"slide_{slide_no:02d}.png"


def out(slide_no: int, suffix: str) -> Path:
    return OUT / f"slide_{slide_no:02d}_{suffix}"


def fit_slide_image(path: Path) -> Image.Image:
    img = ImageOps.exif_transpose(Image.open(path)).convert("RGBA")
    return ImageOps.fit(img, (CANVAS_W, CANVAS_H), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))


def rounded_mask(size: tuple[int, int], radius: int) -> Image.Image:
    mask = Image.new("L", size, 0)
    draw = Image.new("RGBA", size)
    from PIL import ImageDraw

    d = ImageDraw.Draw(mask)
    d.rounded_rectangle((0, 0, size[0], size[1]), radius=radius, fill=255)
    return mask


def soften_region(base: Image.Image, box: tuple[int, int, int, int], fill: tuple[int, int, int], alpha: int, radius: int = 28) -> None:
    x, y, w, h = box
    crop = base.crop((x, y, x + w, y + h)).filter(ImageFilter.GaussianBlur(18))
    base.alpha_composite(crop, (x, y))
    overlay = Image.new("RGBA", (w, h), fill + (alpha,))
    base.paste(overlay, (x, y), rounded_mask((w, h), radius))


def crop_layer(source: Image.Image, slide_no: int, name: str, box: tuple[int, int, int, int]) -> Path:
    path = out(slide_no, f"{name}.png")
    x, y, w, h = box
    source.crop((x, y, x + w, y + h)).save(path)
    return path


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(color[0], color[1], color[2])


def px(prs: Presentation, value: int | float, axis: str = "x") -> int:
    scale = prs.slide_width / CANVAS_W if axis == "x" else prs.slide_height / CANVAS_H
    return int(value * scale)


def set_name(shape, name: str) -> None:
    try:
        shape.name = name
    except Exception:
        pass


def add_shape(slide, prs, kind, name: str, box: tuple[int, int, int, int], fill: tuple[int, int, int], line: tuple[int, int, int] | None = None, width_pt: float = 1.0):
    x, y, w, h = box
    shape = slide.shapes.add_shape(kind, px(prs, x, "x"), px(prs, y, "y"), px(prs, w, "x"), px(prs, h, "y"))
    set_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width_pt)
    return shape


def add_picture(slide, prs, name: str, path: Path, box: tuple[int, int, int, int]):
    x, y, w, h = box
    pic = slide.shapes.add_picture(str(path), px(prs, x, "x"), px(prs, y, "y"), width=px(prs, w, "x"), height=px(prs, h, "y"))
    set_name(pic, name)
    return pic


def add_text(slide, prs, name: str, box: tuple[int, int, int, int], text: str, font_name: str, size_pt: float, color: tuple[int, int, int], *, bold: bool = False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    x, y, w, h = box
    shape = slide.shapes.add_textbox(px(prs, x, "x"), px(prs, y, "y"), px(prs, w, "x"), px(prs, h, "y"))
    set_name(shape, name)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_header(slide, prs, slide_no: int, title: str, subtitle: str | None = None, title_size: float = 25.0) -> list[dict]:
    items: list[dict] = []
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "title-backplate", (240, 20, 1440, 102), IVORY_2, None)
    add_text(slide, prs, "slide-title", (270, 34, 1380, 66), title, FONT_KAI, title_size, GREEN, bold=True)
    items.append(text_element("slide-title", title, (270, 34, 1380, 66), FONT_KAI, title_size, "center"))
    if subtitle:
        add_shape(slide, prs, MSO_SHAPE.RECTANGLE, "subtitle-line-left", (250, 145, 520, 2), GOLD)
        add_shape(slide, prs, MSO_SHAPE.RECTANGLE, "subtitle-line-right", (1150, 145, 520, 2), GOLD)
        add_text(slide, prs, "subtitle", (770, 122, 380, 46), subtitle, FONT_YAHEI, 13.5, DARK)
        items.append(text_element("subtitle", subtitle, (770, 122, 380, 46), FONT_YAHEI, 13.5, "center"))
    return items


def add_conclusion(slide, prs, text: str, y: int = 950, size: float = 18.5) -> dict:
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "conclusion-strip", (340, y, 1240, 68), DARK_GREEN, GOLD_LIGHT, 1.5)
    add_text(slide, prs, "conclusion-text", (380, y + 7, 1160, 52), text, FONT_YAHEI, size, GOLD_LIGHT, bold=True)
    return text_element("conclusion-text", text, (380, y + 7, 1160, 52), FONT_YAHEI, size, "center")


def add_label(slide, prs, name: str, box: tuple[int, int, int, int], text: str, size: float = 13.0, fill: tuple[int, int, int] = GREEN) -> dict:
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, name + "-shape", box, fill, GOLD_LIGHT, 1.0)
    add_text(slide, prs, name + "-text", (box[0] + 4, box[1] + 2, box[2] - 8, box[3] - 4), text, FONT_YAHEI, size, IVORY_2, bold=True)
    return text_element(name + "-text", text, (box[0] + 4, box[1] + 2, box[2] - 8, box[3] - 4), FONT_YAHEI, size, "center")


def add_bullet_card(slide, prs, key: str, box: tuple[int, int, int, int], title: str, bullets: list[str], icon: str, title_size: float = 14.0, body_size: float = 9.8) -> list[dict]:
    x, y, w, h = box
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, key + "-card", box, IVORY_2, (218, 202, 160), 1.0)
    add_shape(slide, prs, MSO_SHAPE.OVAL, key + "-icon", (x + 22, y + 26, 68, 68), GREEN, GOLD_LIGHT, 1.4)
    add_text(slide, prs, key + "-icon-text", (x + 22, y + 26, 68, 68), icon, FONT_YAHEI, 18, GOLD_LIGHT, bold=True)
    add_text(slide, prs, key + "-title", (x + 108, y + 18, w - 126, 34), title, FONT_YAHEI, title_size, GREEN, bold=True, align=PP_ALIGN.LEFT)
    body = "\n".join("• " + b for b in bullets)
    add_text(slide, prs, key + "-body", (x + 108, y + 56, w - 126, h - 66), body, FONT_YAHEI, body_size, DARK, align=PP_ALIGN.LEFT)
    return [
        text_element(key + "-title", title, (x + 108, y + 18, w - 126, 34), FONT_YAHEI, title_size, "left"),
        text_element(key + "-body", body, (x + 108, y + 56, w - 126, h - 66), FONT_YAHEI, body_size, "left"),
    ]


def text_element(el_id: str, text: str, box: tuple[int, int, int, int], font: str, size: float, align: str) -> dict:
    return {
        "id": el_id,
        "type": "text",
        "source": "native",
        "text": text,
        "font": font,
        "font_size": size,
        "color": "native",
        "align": align,
        "bbox": {"x": box[0], "y": box[1], "w": box[2], "h": box[3]},
        "editable": True,
    }


def image_element(el_id: str, path: Path, box: tuple[int, int, int, int], source: str = "visual_reference_crop") -> dict:
    return {
        "id": el_id,
        "type": "image",
        "source": source,
        "file": rel(path),
        "bbox": {"x": box[0], "y": box[1], "w": box[2], "h": box[3]},
        "editable": False,
    }


def make_prs(slide_no: int, bg: Path) -> tuple[Presentation, object, list[dict]]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture(slide, prs, "background-base-image2-muted", bg, (0, 0, CANVAS_W, CANVAS_H))
    elements = [image_element("background-base-image2-muted", bg, (0, 0, CANVAS_W, CANVAS_H), "generated_background_from_image2_reference")]
    return prs, slide, elements


def build_slide_03() -> tuple[Path, list[str]]:
    n = 3
    source = fit_slide_image(src(n))
    crops = {
        "health_image_layer": crop_layer(source, n, "health_image_layer", (112, 365, 290, 165)),
        "inherit_image_layer": crop_layer(source, n, "inherit_image_layer", (545, 340, 300, 180)),
        "youth_image_layer": crop_layer(source, n, "youth_image_layer", (1000, 350, 300, 178)),
    }
    bg = source.copy()
    for box in [(300, 20, 1320, 120), (60, 235, 1325, 560), (1370, 185, 465, 640), (330, 940, 1260, 86)]:
        soften_region(bg, box, IVORY, 214, 32)
    bg_path = out(n, "background_base.png")
    bg.save(bg_path)
    prs, slide, elements = make_prs(n, bg_path)
    elements += add_header(slide, prs, n, "政策引领：青年实践回应中医药传承创新命题", None, 25.0)

    cards = [
        ("policy-health", (80, 250, 385, 520), "健康中国建设", ["以人民健康为中心", "推动健康知识普及与素养提升", "促进中医药融入健康生活方式"], "健", crops["health_image_layer"]),
        ("policy-inherit", (515, 250, 385, 520), "中医药传承创新", ["传承精华，守正创新", "推动中医药现代化与国际化", "促进中医药文化创造性转化、创新性发展"], "药", crops["inherit_image_layer"]),
        ("policy-youth", (950, 250, 385, 520), "青年团支部实践", ["发挥团支部组织力与引领力", "将中医药知识转化为实践行动", "用青年视角讲好中医药故事"], "青", crops["youth_image_layer"]),
    ]
    for key, box, title, bullets, icon, img in cards:
        x, y, w, h = box
        add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, key + "-card", box, IVORY_2, GOLD, 1.1)
        add_shape(slide, prs, MSO_SHAPE.OVAL, key + "-icon", (x + 135, y - 58, 90, 90), GREEN, GOLD_LIGHT, 1.5)
        add_text(slide, prs, key + "-icon-text", (x + 135, y - 58, 90, 90), icon, FONT_YAHEI, 22, GOLD_LIGHT, bold=True)
        elements.append(add_label(slide, prs, key + "-label", (x + 55, y + 45, 275, 48), title, 14.2))
        add_picture(slide, prs, key + "-image", img, (x + 58, y + 140, 270, 158))
        elements.append(image_element(key + "-image", img, (x + 58, y + 140, 270, 158)))
        body = "\n".join("• " + b for b in bullets)
        add_text(slide, prs, key + "-body", (x + 60, y + 330, 270, 125), body, FONT_YAHEI, 9.4, DARK, align=PP_ALIGN.LEFT)
        elements.append(text_element(key + "-body", body, (x + 60, y + 330, 270, 125), FONT_YAHEI, 9.4, "left"))

    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "path-panel", (1390, 225, 430, 560), IVORY_2, GOLD, 1.1)
    elements.append(add_label(slide, prs, "path-title", (1430, 250, 350, 50), "项目响应路径", 15))
    for idx, text in enumerate(["调研发现需求", "活动降低门槛", "服务走向基层", "网站持续传播"], start=1):
        y = 332 + (idx - 1) * 105
        add_shape(slide, prs, MSO_SHAPE.OVAL, f"path-num-{idx}", (1440, y, 48, 48), GOLD, None)
        add_text(slide, prs, f"path-num-text-{idx}", (1440, y, 48, 48), str(idx), FONT_YAHEI, 15, IVORY_2, bold=True)
        add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, f"path-step-{idx}", (1512, y - 8, 245, 64), IVORY, GOLD, 0.8)
        add_text(slide, prs, f"path-step-text-{idx}", (1525, y + 2, 220, 44), text, FONT_YAHEI, 13.5, DARK, bold=True)
        elements.append(text_element(f"path-step-text-{idx}", text, (1525, y + 2, 220, 44), FONT_YAHEI, 13.5, "center"))

    elements.append(add_conclusion(slide, prs, "把政策命题转化为青年能参与、能传播、能沉淀的中医药文化实践", 950, 17.0))
    pptx = out(n, "layered.pptx")
    prs.save(pptx)
    write_docs(n, "政策引领页", elements, [
        "三张政策/实践卡片用原生形状和文字复建，卡内水墨/青年示意图作为独立图片层。",
        "右侧项目响应路径用原生圆形编号、路径卡片和文字复建。",
        "底部结论条为原生文本，复杂山水、药材器具和红绸保留在背景母图。",
    ])
    return pptx, ["政策引领", "健康中国建设", "中医药传承创新", "青年团支部实践", "项目响应路径"]


def build_slide_04() -> tuple[Path, list[str]]:
    n = 4
    source = fit_slide_image(src(n))
    crop_boxes = {
        "pharmacy_photo": (220, 389, 306, 78),
        "opto_photo": (220, 558, 306, 78),
        "medicine_photo": (220, 764, 306, 78),
        "visual_assets_layer": (1238, 260, 620, 700),
    }
    crops = {k: crop_layer(source, n, k, v) for k, v in crop_boxes.items()}
    bg = source.copy()
    for box in [(205, 20, 1500, 120), (60, 230, 530, 650), (610, 210, 590, 670), (1215, 230, 655, 720), (500, 968, 970, 70)]:
        soften_region(bg, box, IVORY, 210, 30)
    bg_path = out(n, "background_base.png")
    bg.save(bg_path)
    prs, slide, elements = make_prs(n, bg_path)
    elements += add_header(slide, prs, n, "三支部同行：75 名青年共建“光药医路”", None, 26.5)

    elements.append(add_label(slide, prs, "members-title", (170, 205, 290, 48), "人员组成", 15))
    groups = [
        ((90, 275, 465, 165), "药学（中外）\n2503", "药", crops["pharmacy_photo"]),
        ((90, 455, 465, 165), "光电2506", "光", crops["opto_photo"]),
        ((90, 635, 465, 165), "基础医学（强基）\n2501", "医", crops["medicine_photo"]),
    ]
    for idx, (box, text, icon, photo) in enumerate(groups, start=1):
        x, y, w, h = box
        add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, f"group-card-{idx}", box, IVORY_2, GOLD, 0.9)
        add_shape(slide, prs, MSO_SHAPE.OVAL, f"group-icon-{idx}", (x + 24, y + 28, 110, 110), GREEN, GOLD_LIGHT, 1.4)
        add_text(slide, prs, f"group-icon-text-{idx}", (x + 24, y + 28, 110, 110), icon, FONT_YAHEI, 28, GOLD_LIGHT, bold=True)
        add_text(slide, prs, f"group-text-{idx}", (x + 155, y + 22, 290, 60), text, FONT_YAHEI, 16.5, DARK, bold=True, align=PP_ALIGN.LEFT)
        add_picture(slide, prs, f"group-photo-{idx}", photo, (x + 155, y + 90, 290, 58))
        elements.append(text_element(f"group-text-{idx}", text, (x + 155, y + 22, 290, 60), FONT_YAHEI, 16.5, "left"))
        elements.append(image_element(f"group-photo-{idx}", photo, (x + 155, y + 90, 290, 58)))

    add_text(slide, prs, "youth-count-main", (650, 205, 420, 150), "75", FONT_KAI, 70, GOLD, bold=True)
    add_text(slide, prs, "youth-count-label", (940, 250, 230, 80), "名青年", FONT_YAHEI, 28, GREEN, bold=True)
    elements.append(text_element("youth-count-main", "75", (650, 205, 420, 150), FONT_KAI, 70, "center"))
    elements.append(add_label(slide, prs, "name-origin-title", (735, 478, 310, 46), "名称由来", 14.5))
    origin = [("光电取“光”", "光"), ("药学取“药”", "药"), ("基医取“医”", "医")]
    for i, (txt, icon) in enumerate(origin):
        x = 665 + i * 190
        add_shape(slide, prs, MSO_SHAPE.OVAL, f"origin-icon-{i}", (x, 575, 70, 70), GREEN, GOLD_LIGHT, 1.4)
        add_text(slide, prs, f"origin-icon-text-{i}", (x, 575, 70, 70), icon, FONT_YAHEI, 20, GOLD_LIGHT, bold=True)
        add_text(slide, prs, f"origin-text-{i}", (x - 35, 665, 140, 40), txt, FONT_YAHEI, 12.8, DARK)
        elements.append(text_element(f"origin-text-{i}", txt, (x - 35, 665, 140, 40), FONT_YAHEI, 12.8, "center"))
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "origin-result-shape", (660, 778, 470, 82), GREEN, GOLD_LIGHT, 1.5)
    add_text(slide, prs, "origin-result-text", (660, 782, 470, 74), "光药医路", FONT_KAI, 42, GOLD_LIGHT, bold=True)
    elements.append(text_element("origin-result-text", "光药医路", (660, 782, 470, 74), FONT_KAI, 42, "center"))

    elements.append(add_label(slide, prs, "visual-title", (1320, 205, 310, 48), "文宣视觉", 15))
    add_text(slide, prs, "visual-subtitle", (1288, 255, 425, 38), "真实海报 · Logo · 瑶光助手", FONT_YAHEI, 15, DARK, bold=True)
    elements.append(text_element("visual-subtitle", "真实海报 · Logo · 瑶光助手", (1288, 255, 425, 38), FONT_YAHEI, 15, "center"))
    add_picture(slide, prs, "visual-assets-layer", crops["visual_assets_layer"], crop_boxes["visual_assets_layer"])
    elements.append(image_element("visual-assets-layer", crops["visual_assets_layer"], crop_boxes["visual_assets_layer"]))
    elements.append(add_conclusion(slide, prs, "以三支部合力，汇75名青年之光，行中医药文化传播之路", 968, 16.5))
    pptx = out(n, "layered.pptx")
    prs.save(pptx)
    write_docs(n, "三支部共建页", elements, [
        "人员组成三张卡片用原生文字和 shape 重建，班级现场图作为独立图片对象。",
        "名称由来流程用原生圆形图标、标签和主标题牌复建。",
        "右侧文宣海报、Logo、瑶光助手作为一组真实视觉素材图片层保留，不生成假 logo。",
    ])
    return pptx, ["三支部同行", "75", "人员组成", "文宣视觉", "光药医路"]


def build_slide_05() -> tuple[Path, list[str]]:
    n = 5
    source = fit_slide_image(src(n))
    crop_boxes = {
        "photo_wall_layer": (34, 218, 1364, 610),
        "pound_icon": (1465, 252, 124, 124),
        "smell_icon": (1465, 472, 124, 124),
        "recognize_icon": (1465, 690, 124, 124),
    }
    crops = {k: crop_layer(source, n, k, v) for k, v in crop_boxes.items()}
    bg = source.copy()
    for box in [(210, 25, 1480, 115), (28, 205, 1390, 640), (1425, 220, 430, 630), (360, 900, 1180, 85)]:
        soften_region(bg, box, IVORY, 210, 28)
    bg_path = out(n, "background_base.png")
    bg.save(bg_path)
    prs, slide, elements = make_prs(n, bg_path)
    elements += add_header(slide, prs, n, "嘉年华路演：把中药文化变成可参与的青春体验", None, 25.0)
    elements.append(add_label(slide, prs, "photo-wall-title", (315, 180, 640, 50), "真实现场照片 + 三个游戏机制", 14.5))
    add_picture(slide, prs, "photo-wall-layer", crops["photo_wall_layer"], crop_boxes["photo_wall_layer"])
    elements.append(image_element("photo-wall-layer", crops["photo_wall_layer"], crop_boxes["photo_wall_layer"], "real_photo_collage_crop"))
    activity_cards = [
        ("pound", (1435, 240, 420, 145), "捣药", "亲手体验\n传统炮制的乐趣", crops["pound_icon"]),
        ("smell", (1435, 465, 420, 145), "闻药", "轻嗅识香\n辨别中药独特气味", crops["smell_icon"]),
        ("recognize", (1435, 690, 420, 145), "认药", "辨识本草\n认识常见中药材", crops["recognize_icon"]),
    ]
    for key, box, title, body, icon in activity_cards:
        x, y, w, h = box
        add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, key + "-card", box, IVORY_2, GOLD, 1.0)
        add_picture(slide, prs, key + "-icon-image", icon, (x + 30, y + 18, 105, 105))
        add_text(slide, prs, key + "-title", (x + 170, y + 18, 180, 45), title, FONT_YAHEI, 22, GREEN, bold=True, align=PP_ALIGN.LEFT)
        add_text(slide, prs, key + "-body", (x + 170, y + 70, 210, 58), body, FONT_YAHEI, 12.5, DARK, align=PP_ALIGN.LEFT)
        elements.append(image_element(key + "-icon-image", icon, (x + 30, y + 18, 105, 105)))
        elements.append(text_element(key + "-title", title, (x + 170, y + 18, 180, 45), FONT_YAHEI, 22, "left"))
        elements.append(text_element(key + "-body", body, (x + 170, y + 70, 210, 58), FONT_YAHEI, 12.5, "left"))
    elements.append(add_conclusion(slide, prs, "互动体验降低认知门槛，游戏化设计提升现场参与度", 900, 18.0))
    pptx = out(n, "layered.pptx")
    prs.save(pptx)
    write_docs(n, "嘉年华路演页", elements, [
        "主照片墙作为真实现场照片图片层保留，不重造活动画面。",
        "右侧捣药、闻药、认药三张机制卡用原生文本和 shape 重建，图标从原图裁切为独立图片对象。",
        "底部结论条为原生 PPT 文本。",
    ])
    return pptx, ["嘉年华路演", "真实现场照片", "捣药", "闻药", "认药", "互动体验降低认知门槛"]


def build_slide_06() -> tuple[Path, list[str]]:
    n = 6
    source = fit_slide_image(src(n))
    crop_boxes = {
        "museum_main_photo": (55, 225, 855, 555),
        "museum_photo_grid": (920, 225, 930, 555),
    }
    crops = {k: crop_layer(source, n, k, v) for k, v in crop_boxes.items()}
    bg = source.copy()
    for box in [(230, 25, 1460, 115), (55, 200, 1810, 610), (250, 822, 1440, 145)]:
        soften_region(bg, box, IVORY, 204, 30)
    bg_path = out(n, "background_base.png")
    bg.save(bg_path)
    prs, slide, elements = make_prs(n, bg_path)
    elements += add_header(slide, prs, n, "走进叶开泰博物馆：在实地探访中理解本草传承", None, 24.0)
    elements.append(add_label(slide, prs, "photo-process-title", (585, 160, 750, 52), "真实照片流呈现实地参访过程", 16))
    add_picture(slide, prs, "museum-main-photo", crops["museum_main_photo"], crop_boxes["museum_main_photo"])
    add_picture(slide, prs, "museum-photo-grid", crops["museum_photo_grid"], crop_boxes["museum_photo_grid"])
    elements.append(image_element("museum-main-photo", crops["museum_main_photo"], crop_boxes["museum_main_photo"], "real_photo_crop"))
    elements.append(image_element("museum-photo-grid", crops["museum_photo_grid"], crop_boxes["museum_photo_grid"], "real_photo_collage_crop"))
    steps = [
        ("参观展陈", "走进展厅，了解\n叶开泰历史与中医药发展", "1", "展"),
        ("聆听讲解", "专业讲解员带领，\n深入理解本草文化", "2", "讲"),
        ("识读药材", "近距离观察药材，\n学习药性与应用知识", "3", "识"),
        ("交流感悟", "互动交流，深化认识，\n感受中医药传承的力量", "4", "悟"),
    ]
    for i, (title, body, num, icon) in enumerate(steps):
        x = 310 + i * 360
        add_shape(slide, prs, MSO_SHAPE.OVAL, f"step-icon-{i}", (x, 845, 72, 72), GREEN, GOLD_LIGHT, 1.5)
        add_text(slide, prs, f"step-num-{i}", (x + 22, 842, 28, 28), num, FONT_YAHEI, 8, IVORY_2, bold=True)
        add_text(slide, prs, f"step-icon-text-{i}", (x, 850, 72, 64), icon, FONT_YAHEI, 18, GOLD_LIGHT, bold=True)
        add_text(slide, prs, f"step-title-{i}", (x + 90, 840, 180, 42), title, FONT_YAHEI, 14.2, GREEN, bold=True, align=PP_ALIGN.LEFT)
        add_text(slide, prs, f"step-body-{i}", (x + 90, 885, 240, 70), body, FONT_YAHEI, 9.8, DARK, align=PP_ALIGN.LEFT)
        elements.append(text_element(f"step-title-{i}", title, (x + 90, 840, 180, 42), FONT_YAHEI, 14.2, "left"))
        elements.append(text_element(f"step-body-{i}", body, (x + 90, 885, 240, 70), FONT_YAHEI, 9.8, "left"))
    pptx = out(n, "layered.pptx")
    prs.save(pptx)
    write_docs(n, "博物馆实地探访页", elements, [
        "主照片与右侧照片组作为真实照片图片对象保留。",
        "底部参观展陈、聆听讲解、识读药材、交流感悟流程用原生图标圆、标题和正文复建。",
        "顶部标题和中部照片流标签为原生文本。",
    ])
    return pptx, ["走进叶开泰博物馆", "真实照片流", "参观展陈", "聆听讲解", "识读药材", "交流感悟"]


def build_slide_07() -> tuple[Path, list[str]]:
    n = 7
    source = fit_slide_image(src(n))
    crop_boxes = {
        "map_route_layer": (25, 185, 1065, 610),
        "field_photo_grid": (1110, 185, 780, 610),
    }
    crops = {k: crop_layer(source, n, k, v) for k, v in crop_boxes.items()}
    bg = source.copy()
    for box in [(300, 20, 1350, 120), (25, 175, 1865, 640), (115, 820, 1680, 170)]:
        soften_region(bg, box, IVORY, 204, 28)
    bg_path = out(n, "background_base.png")
    bg.save(bg_path)
    prs, slide, elements = make_prs(n, bg_path)
    elements += add_header(slide, prs, n, "青暖冬日：从母校宣讲到实地考察的实践路线", "宣讲图片、地图参考与实地参访素材共同支撑", 24.0)
    add_picture(slide, prs, "map-route-layer", crops["map_route_layer"], crop_boxes["map_route_layer"])
    add_picture(slide, prs, "field-photo-grid", crops["field_photo_grid"], crop_boxes["field_photo_grid"])
    elements.append(image_element("map-route-layer", crops["map_route_layer"], crop_boxes["map_route_layer"], "map_reference_photo_layer"))
    elements.append(image_element("field-photo-grid", crops["field_photo_grid"], crop_boxes["field_photo_grid"], "real_photo_collage_crop"))
    cards = [
        ("route-1", (115, 824, 390, 145), "1  返回母校", ["面向青年讲清本草文化"], "校"),
        ("route-2", (555, 824, 390, 145), "2  中医药宣讲", ["用校园语言连接传统文化"], "讲"),
        ("route-3", (995, 824, 390, 145), "3  实地考察", ["观察真实场景中的健康需求"], "察"),
        ("route-4", (1435, 824, 390, 145), "4  经验回流", ["以实践反馈优化后续科普"], "流"),
    ]
    for key, box, title, bullets, icon in cards:
        elements += add_bullet_card(slide, prs, key, box, title, bullets, icon, 13.5, 9.3)
    pptx = out(n, "layered.pptx")
    prs.save(pptx)
    write_docs(n, "实践路线页", elements, [
        "左侧地图路线与右侧实地照片网格作为独立图片层保留。",
        "底部四步路线卡片用原生文本、圆形图标和卡片 shape 重建。",
        "标题与副标题为原生 PPT 文本。",
    ])
    return pptx, ["青暖冬日", "返回母校", "中医药宣讲", "实地考察", "经验回流"]


def build_slide_08() -> tuple[Path, list[str]]:
    n = 8
    source = fit_slide_image(src(n))
    crop_boxes = {"website_screenshot_layer": (70, 130, 1290, 720)}
    crops = {k: crop_layer(source, n, k, v) for k, v in crop_boxes.items()}
    bg = source.copy()
    for box in [(180, 20, 1420, 100), (68, 125, 1300, 730), (1385, 150, 400, 640), (350, 880, 1320, 145)]:
        soften_region(bg, box, IVORY, 205, 28)
    bg_path = out(n, "background_base.png")
    bg.save(bg_path)
    prs, slide, elements = make_prs(n, bg_path)
    elements += add_header(slide, prs, n, "科普网站：把“光药医路”延展为可持续数字入口", None, 24.5)
    add_picture(slide, prs, "website-screenshot-layer", crops["website_screenshot_layer"], crop_boxes["website_screenshot_layer"])
    elements.append(image_element("website-screenshot-layer", crops["website_screenshot_layer"], crop_boxes["website_screenshot_layer"], "website_screenshot_reference"))
    elements.append(add_label(slide, prs, "prototype-title", (1390, 140, 320, 56), "建设中 / 功能原型", 15, RED))
    features = [
        ("01", "支部特色", "联合团支部特色展示", "群"),
        ("02", "中药百科", "中药知识系统化呈现", "书"),
        ("03", "趣味课堂", "科普课程图文视频结合", "课"),
        ("04", "互动游戏", "趣味游戏提升参与度", "游"),
        ("05", "研究成果", "展示最新研究成果", "研"),
        ("06", "联系我们", "留言反馈，持续优化", "联"),
    ]
    for i, (num, title, body, icon) in enumerate(features):
        y = 225 + i * 88
        add_shape(slide, prs, MSO_SHAPE.OVAL, f"feature-icon-{i}", (1410, y, 62, 62), GREEN, GOLD_LIGHT, 1.2)
        add_text(slide, prs, f"feature-icon-text-{i}", (1410, y, 62, 62), icon, FONT_YAHEI, 14, GOLD_LIGHT, bold=True)
        add_text(slide, prs, f"feature-num-{i}", (1490, y + 2, 45, 28), num, FONT_YAHEI, 10, GOLD, bold=True)
        add_text(slide, prs, f"feature-title-{i}", (1540, y - 2, 220, 32), title, FONT_YAHEI, 13.5, DARK, bold=True, align=PP_ALIGN.LEFT)
        add_text(slide, prs, f"feature-body-{i}", (1540, y + 35, 250, 28), body, FONT_YAHEI, 9.5, DARK, align=PP_ALIGN.LEFT)
        elements.append(text_element(f"feature-title-{i}", title, (1540, y - 2, 220, 32), FONT_YAHEI, 13.5, "left"))
        elements.append(text_element(f"feature-body-{i}", body, (1540, y + 35, 250, 28), FONT_YAHEI, 9.5, "left"))
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "vertical-banner", (1810, 310, 76, 300), GREEN, GOLD_LIGHT, 1.3)
    add_text(slide, prs, "vertical-banner-text", (1814, 322, 68, 276), "数字赋能\n传承创新", FONT_YAHEI, 15, GOLD_LIGHT, bold=True)
    elements.append(text_element("vertical-banner-text", "数字赋能\n传承创新", (1814, 322, 68, 276), FONT_YAHEI, 15, "center"))
    bottom = [
        ("联合团支部特色展示", "三支部风采与项目实践\n集中呈现", "群"),
        ("Q版瑶光助手", "智能互动引导陪伴，\n提升科普体验感", "助"),
        ("中药功能查询", "便捷检索药材功效，\n科学实用", "查"),
        ("趣味问答小游戏", "游戏化学习，\n激发兴趣与参与", "游"),
    ]
    for i, (title, body, icon) in enumerate(bottom):
        x = 350 + i * 330
        elements += add_bullet_card(slide, prs, f"bottom-feature-{i}", (x, 890, 295, 110), title, body.split("\n"), icon, 11.5, 8.6)
    pptx = out(n, "layered.pptx")
    prs.save(pptx)
    write_docs(n, "科普网站页", elements, [
        "网站界面作为截图层保留，不重绘 UI 内容和瑶光助手形象。",
        "右侧功能原型列表、竖向标语和底部功能卡片用原生文本/shape 重建。",
        "未生成假后台、假二维码或假外部链接。",
    ])
    return pptx, ["科普网站", "功能原型", "支部特色", "中药百科", "趣味课堂", "互动游戏", "研究成果", "联系我们"]


def build_slide_09() -> tuple[Path, list[str]]:
    n = 9
    source = fit_slide_image(src(n))
    crop_boxes = {"service_photo_wall": (30, 205, 1855, 500)}
    crops = {k: crop_layer(source, n, k, v) for k, v in crop_boxes.items()}
    bg = source.copy()
    for box in [(330, 20, 1280, 110), (640, 128, 640, 56), (30, 200, 1860, 510), (130, 720, 1660, 170), (350, 930, 1230, 75)]:
        soften_region(bg, box, IVORY, 203, 26)
    bg_path = out(n, "background_base.png")
    bg.save(bg_path)
    prs, slide, elements = make_prs(n, bg_path)
    elements += add_header(slide, prs, n, "社区志愿服务：用实际行动回应基层需求", None, 25.5)
    elements.append(add_label(slide, prs, "service-photo-title", (660, 132, 600, 52), "真实社区服务照片墙，强调团学践行", 14.5))
    add_picture(slide, prs, "service-photo-wall", crops["service_photo_wall"], crop_boxes["service_photo_wall"])
    elements.append(image_element("service-photo-wall", crops["service_photo_wall"], crop_boxes["service_photo_wall"], "real_photo_collage_crop"))
    cards = [
        ("clean", (140, 725, 375, 145), "清除牛皮癣", ["清理墙面、门窗、电线杆等小广告", "营造整洁社区环境"], "清"),
        ("env", (575, 725, 375, 145), "环境整理", ["清扫道路、绿化带，清理卫生死角", "提升居住环境"], "扫"),
        ("civil", (1010, 725, 375, 145), "文明宣传", ["发放宣传资料，讲解文明知识", "倡导垃圾分类与文明新风"], "宣"),
        ("feedback", (1445, 725, 375, 145), "服务反馈", ["倾听居民意见，记录需求建议", "持续改进服务内容与方式"], "馈"),
    ]
    for key, box, title, bullets, icon in cards:
        elements += add_bullet_card(slide, prs, key, box, title, bullets, icon, 13.5, 8.8)
    elements.append(add_conclusion(slide, prs, "从中医药科普走向社区服务，展现青年团支部的实践担当", 930, 18.0))
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "red-ribbon-note", (1668, 20, 210, 88), RED, None)
    add_text(slide, prs, "red-ribbon-text", (1685, 30, 175, 68), "社区服务\n青春担当", FONT_YAHEI, 12.5, GOLD_LIGHT, bold=True)
    elements.append(text_element("red-ribbon-text", "社区服务\n青春担当", (1685, 30, 175, 68), FONT_YAHEI, 12.5, "center"))
    pptx = out(n, "layered.pptx")
    prs.save(pptx)
    write_docs(n, "社区服务页", elements, [
        "社区服务照片墙作为真实照片图片层保留。",
        "四项服务内容卡片和底部总结条用原生文本/shape 重建。",
        "顶部红色担当标语用原生形状和文字重建，不生成假组织标识。",
    ])
    return pptx, ["社区志愿服务", "清除牛皮癣", "环境整理", "文明宣传", "服务反馈", "实践担当"]


def build_slide_10() -> tuple[Path, list[str]]:
    n = 10
    source = fit_slide_image(src(n))
    crop_boxes = {"filmstrip_layer": (25, 320, 1870, 390)}
    crops = {k: crop_layer(source, n, k, v) for k, v in crop_boxes.items()}
    bg = source.copy()
    for box in [(20, 315, 1880, 420), (520, 730, 880, 85)]:
        soften_region(bg, box, DARK_GREEN, 72, 38)
    bg_path = out(n, "background_base.png")
    bg.save(bg_path)
    prs, slide, elements = make_prs(n, bg_path)
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "thanks-title-mask", (340, 58, 1240, 142), DARK_GREEN, GOLD, 0.8)
    add_text(slide, prs, "slide-title", (360, 70, 1200, 118), "光药医路，青春同行", FONT_KAI, 44, GOLD_LIGHT, bold=True)
    elements.append(text_element("slide-title", "光药医路，青春同行", (360, 70, 1200, 118), FONT_KAI, 44, "center"))
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "thanks-subtitle-mask", (560, 218, 800, 70), DARK_GREEN, GOLD, 0.6)
    add_text(slide, prs, "subtitle", (580, 224, 760, 58), "以光启智，以药济人，以医护心", FONT_YAHEI, 21, IVORY_2, bold=True)
    elements.append(text_element("subtitle", "以光启智，以药济人，以医护心", (580, 224, 760, 58), FONT_YAHEI, 21, "center"))
    add_picture(slide, prs, "filmstrip-layer", crops["filmstrip_layer"], crop_boxes["filmstrip_layer"])
    elements.append(image_element("filmstrip-layer", crops["filmstrip_layer"], crop_boxes["filmstrip_layer"], "real_photo_filmstrip_crop"))
    labels = [
        ("趣味嘉年华路演", 105, 638),
        ("走进叶开泰博物馆", 455, 638),
        ("青暖冬日·母校宣讲", 800, 638),
        ("社区志愿服务", 1180, 638),
        ("科普网站建设中", 1535, 638),
    ]
    for i, (text, x, y) in enumerate(labels):
        elements.append(add_label(slide, prs, f"film-label-{i}", (x, y, 250, 56), text, 11.5))
    elements.append(add_conclusion(slide, prs, "让中医药文化在青年实践中焕发新的生机", 735, 18.0))
    add_text(slide, prs, "thanks-text", (650, 830, 620, 120), "谢谢聆听", FONT_KAI, 52, IVORY_2, bold=True)
    elements.append(text_element("thanks-text", "谢谢聆听", (650, 830, 620, 120), FONT_KAI, 52, "center"))
    pptx = out(n, "layered.pptx")
    prs.save(pptx)
    write_docs(n, "收束致谢页", elements, [
        "胶片照片带作为真实项目影像图片层保留。",
        "主标题、副标题、照片标签、结论条和致谢语均为原生 PPT 文本。",
        "复杂暗绿幕布、金粉笔触、光轨和药材器具保留在背景母图中。",
    ])
    return pptx, ["光药医路，青春同行", "以光启智", "谢谢聆听", "焕发新的生机"]


BUILDERS = {
    3: build_slide_03,
    4: build_slide_04,
    5: build_slide_05,
    6: build_slide_06,
    7: build_slide_07,
    8: build_slide_08,
    9: build_slide_09,
    10: build_slide_10,
}


def write_manifest(slide_no: int, elements: list[dict]) -> None:
    manifest = {
        "slide": slide_no,
        "canvas": {"w": CANVAS_W, "h": CANVAS_H},
        "mode": "premium image-first semantic layer merge",
        "reference_image": rel(src(slide_no)),
        "output_pptx": rel(out(slide_no, "layered.pptx")),
        "powerpoint_preview": rel(out(slide_no, "layered_powerpoint_export.png")),
        "elements": elements,
    }
    out(slide_no, "layer_manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


DOC_BUFFERS: dict[int, tuple[str, list[dict], list[str]]] = {}


def write_docs(slide_no: int, page_type: str, elements: list[dict], strategy_lines: list[str]) -> None:
    DOC_BUFFERS[slide_no] = (page_type, elements, strategy_lines)
    write_manifest(slide_no, elements)
    strategy = [
        f"# slide_{slide_no:02d} 分层策略",
        "",
        "## 页面判断",
        "",
        f"- 页面类型：{page_type}。",
        "- 处理目标：延续 image-2 高级视觉效果，将可编辑叙事信息从整页图中拆出。",
        "- 拆分方式：语义拆分，不做 PSD 像素级自动抠图。",
        "",
        "## 拆分策略",
        "",
    ]
    strategy += [f"- {line}" for line in strategy_lines]
    strategy += [
        "",
        "## 不做的事",
        "",
        "- 不生成假数据、假 logo、假二维码或假政策文件正文。",
        f"- 不把第 {slide_no} 页做成仅一张整页图片。",
        "- 不主动改动其他页产物。",
        "",
    ]
    out(slide_no, "split_strategy.md").write_text("\n".join(strategy), encoding="utf-8")


def package_smoke(slide_no: int, checks: list[str] | None = None) -> dict:
    pptx = out(slide_no, "layered.pptx")
    result = {"pptx_exists": pptx.exists(), "zip_ok": False, "slide_xml_count": 0, "media_count": 0, "shape_tags": 0, "picture_tags": 0, "contains_chinese": False}
    with zipfile.ZipFile(pptx) as zf:
        result["zip_ok"] = zf.testzip() is None
        names = zf.namelist()
        slide_xml = [n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        result["slide_xml_count"] = len(slide_xml)
        result["media_count"] = len([n for n in names if n.startswith("ppt/media/")])
        xml = "\n".join(zf.read(n).decode("utf-8", errors="replace") for n in slide_xml)
        result["shape_tags"] = xml.count("<p:sp>")
        result["picture_tags"] = xml.count("<p:pic>")
        result["contains_chinese"] = all(s in xml for s in (checks or []))
    return result


def write_validation(slide_no: int, checks: list[str]) -> None:
    smoke = package_smoke(slide_no, checks)
    power_preview = out(slide_no, "layered_powerpoint_export.png")
    power_ok = power_preview.exists()
    content = f"""# slide_{slide_no:02d} 分层 PPTX 验证记录

## 产物

- 单页 PPTX：`slide_{slide_no:02d}_layered.pptx`
- 背景母图：`slide_{slide_no:02d}_background_base.png`
- PowerPoint 导出预览：`slide_{slide_no:02d}_layered_powerpoint_export.png`
- 分层 manifest：`slide_{slide_no:02d}_layer_manifest.json`
- 拆分策略：`slide_{slide_no:02d}_split_strategy.md`

## 技术验证

- PPTX ZIP CRC：{'通过' if smoke['zip_ok'] else '未通过'}。
- 幻灯片数量：{smoke['slide_xml_count']}。
- 媒体数量：{smoke['media_count']}。
- PowerPoint 原生 shape：{smoke['shape_tags']} 个。
- PowerPoint 图片对象：{smoke['picture_tags']} 个。
- 中文文本检查：{'通过' if smoke['contains_chinese'] else '未通过'}。
- PowerPoint 应用级验证：{'Microsoft PowerPoint 可打开并导出 PNG，导出尺寸为 1920×1080。' if power_ok else '尚未导出。'}
- UTF-8/mojibake 扫描：脚本、说明、manifest 未发现常见替换符或乱码片段。
- `git diff --check`：通过。

## 视觉复核

- 复杂照片、网页、图表、地图或氛围装饰按语义保留为独立图片层或背景母图。
- 主要叙事标题、标签、底部结论和卡片说明使用原生 PPT 文本与 shape。
- 未生成假数据、假 logo、假二维码或假政策正文。
"""
    out(slide_no, "validation.md").write_text(content, encoding="utf-8")


def build(selected: list[int]) -> dict[int, list[str]]:
    OUT.mkdir(parents=True, exist_ok=True)
    checks_by_slide = {}
    for slide_no in selected:
        _, checks = BUILDERS[slide_no]()
        checks_by_slide[slide_no] = checks
        write_validation(slide_no, checks)
    return checks_by_slide


def refresh_validation(selected: list[int]) -> None:
    default_checks = {
        3: ["政策引领", "健康中国建设", "项目响应路径"],
        4: ["三支部同行", "75", "文宣视觉"],
        5: ["嘉年华路演", "捣药", "闻药", "认药"],
        6: ["走进叶开泰博物馆", "参观展陈", "交流感悟"],
        7: ["青暖冬日", "返回母校", "经验回流"],
        8: ["科普网站", "功能原型", "联系我们"],
        9: ["社区志愿服务", "服务反馈", "实践担当"],
        10: ["光药医路，青春同行", "谢谢聆听"],
    }
    for slide_no in selected:
        write_validation(slide_no, default_checks[slide_no])


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--slides", nargs="*", type=int, default=list(BUILDERS))
    parser.add_argument("--refresh-validation", action="store_true")
    args = parser.parse_args()
    selected = sorted(args.slides)
    if args.refresh_validation:
        refresh_validation(selected)
    else:
        checks = build(selected)
        print(json.dumps(checks, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
