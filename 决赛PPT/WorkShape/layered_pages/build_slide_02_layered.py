from __future__ import annotations

import json
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
SRC = ROOT / "决赛PPT" / "输出终稿" / "image2_premium_ppt_v2" / "slides" / "slide_02.png"
OUT = ROOT / "决赛PPT" / "WorkShape" / "layered_pages"

CANVAS_W = 1920
CANVAS_H = 1080
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5

PPTX_OUT = OUT / "slide_02_layered.pptx"
BG_OUT = OUT / "slide_02_background_base.png"
PREVIEW_OUT = OUT / "slide_02_layered_preview.png"
POWERPOINT_PREVIEW = OUT / "slide_02_layered_powerpoint_export.png"
MANIFEST_OUT = OUT / "slide_02_layer_manifest.json"
STRATEGY_OUT = OUT / "slide_02_split_strategy.md"
VALIDATION_OUT = OUT / "slide_02_validation.md"

SURVEY_LAYER = OUT / "slide_02_survey_screenshots_layer.png"
RESULT_LAYER = OUT / "slide_02_results_charts_layer.png"
PHOTO_01 = OUT / "slide_02_card_photo_01.png"
PHOTO_02 = OUT / "slide_02_card_photo_02.png"
PHOTO_03 = OUT / "slide_02_card_photo_03.png"

FONT_YAHEI = "Microsoft YaHei"
FONT_KAI = "KaiTi"

FONT_REG = Path("C:/Windows/Fonts/msyh.ttc")
FONT_BOLD = Path("C:/Windows/Fonts/msyhbd.ttc")
FONT_KAI_FILE = Path("C:/Windows/Fonts/simkai.ttf")

GREEN = (16, 82, 54)
GREEN_2 = (24, 97, 64)
GOLD = (196, 155, 76)
GOLD_LIGHT = (242, 220, 150)
IVORY = (249, 246, 232)
IVORY_2 = (255, 252, 241)
DARK = (45, 54, 39)
MUTED = (89, 101, 79)


def pil_font(size: int, *, bold: bool = False, kai: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_KAI_FILE if kai else (FONT_BOLD if bold else FONT_REG)
    return ImageFont.truetype(str(path), size)


def fit_slide_image(path: Path) -> Image.Image:
    img = ImageOps.exif_transpose(Image.open(path)).convert("RGBA")
    return ImageOps.fit(img, (CANVAS_W, CANVAS_H), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))


def rounded_mask(size: tuple[int, int], radius: int) -> Image.Image:
    mask = Image.new("L", size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle((0, 0, size[0], size[1]), radius=radius, fill=255)
    return mask


def soften_region(base: Image.Image, box: tuple[int, int, int, int], fill: tuple[int, int, int], alpha: int, radius: int = 28) -> None:
    x, y, w, h = box
    crop = base.crop((x, y, x + w, y + h)).filter(ImageFilter.GaussianBlur(18))
    base.alpha_composite(crop, (x, y))
    overlay = Image.new("RGBA", (w, h), fill + (alpha,))
    base.paste(overlay, (x, y), rounded_mask((w, h), radius))


def crop_layer(base: Image.Image, box: tuple[int, int, int, int], out: Path) -> None:
    x, y, w, h = box
    base.crop((x, y, x + w, y + h)).save(out)


def draw_center(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], text: str, font: ImageFont.FreeTypeFont, fill: tuple[int, int, int], spacing: int = 8) -> None:
    x, y, w, h = box
    bbox = draw.multiline_textbbox((0, 0), text, font=font, spacing=spacing, align="center")
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    draw.multiline_text((x + (w - tw) / 2, y + (h - th) / 2 - bbox[1]), text, font=font, fill=fill, spacing=spacing, align="center")


def draw_layer_image(canvas: Image.Image, img_path: Path, box: tuple[int, int, int, int], radius: int = 16) -> None:
    x, y, w, h = box
    img = Image.open(img_path).convert("RGBA").resize((w, h), Image.Resampling.LANCZOS)
    canvas.paste(img, (x, y), rounded_mask((w, h), radius))


def build_background_and_layers() -> Image.Image:
    source = fit_slide_image(SRC)
    crop_layer(source, (24, 186, 920, 575), SURVEY_LAYER)
    crop_layer(source, (963, 179, 890, 572), RESULT_LAYER)
    crop_layer(source, (376, 815, 178, 122), PHOTO_01)
    crop_layer(source, (955, 809, 140, 126), PHOTO_02)
    crop_layer(source, (1455, 808, 178, 126), PHOTO_03)

    bg = source.copy()
    soften_region(bg, (225, 20, 1410, 125), IVORY, 224, 26)
    soften_region(bg, (20, 180, 930, 590), IVORY, 172, 20)
    soften_region(bg, (958, 174, 900, 585), IVORY, 172, 20)
    soften_region(bg, (20, 770, 1650, 178), IVORY, 198, 20)
    soften_region(bg, (85, 965, 1750, 92), GREEN, 120, 34)
    bg.save(BG_OUT)
    return bg


def make_preview(bg: Image.Image) -> None:
    img = bg.copy()
    d = ImageDraw.Draw(img)

    d.rounded_rectangle((225, 22, 1665, 122), radius=24, fill=IVORY_2 + (240,))
    draw_center(d, (250, 34, 1340, 62), "趣味调研：青年愿意了解中药，但更需要体验式入口", pil_font(43, bold=True, kai=True), GREEN)
    d.line((180, 138, 715, 138), fill=GOLD, width=3)
    d.line((1205, 138, 1740, 138), fill=GOLD, width=3)
    d.rounded_rectangle((705, 111, 1215, 163), radius=18, fill=IVORY_2 + (235,))
    draw_center(d, (720, 115, 480, 50), "问卷截图与结果图均来自项目材料", pil_font(27), DARK)

    for box, layer in [((24, 186, 920, 575), SURVEY_LAYER), ((963, 179, 890, 572), RESULT_LAYER)]:
        x, y, w, h = box
        d.rounded_rectangle((x - 4, y - 4, x + w + 4, y + h + 4), radius=24, fill=(233, 222, 188), outline=GOLD, width=2)
        draw_layer_image(img, layer, box, 18)

    cards = [
        ((30, 790, 548, 150), "青年样本", ["18–25岁青年为主力样本", "多数对中医药持积极态度", "认知基础有待提升"], PHOTO_01, "样"),
        ((610, 790, 505, 150), "认知痛点", ["概念混淆，专业门槛高", "辨识困难，实物接触少", "缺乏有趣、可参与的体验"], PHOTO_02, "痛"),
        ((1140, 790, 505, 150), "活动启发", ["以体验破解认知壁垒", "以互动提升参与热情", "以传播扩大文化影响"], PHOTO_03, "启"),
    ]
    for box, title, bullets, photo, icon in cards:
        x, y, w, h = box
        d.rounded_rectangle((x, y, x + w, y + h), radius=20, fill=IVORY_2, outline=(218, 202, 160), width=2)
        d.ellipse((x + 18, y + 34, x + 84, y + 100), fill=GREEN, outline=GOLD_LIGHT, width=3)
        draw_center(d, (x + 18, y + 34, 66, 66), icon, pil_font(28, bold=True), GOLD_LIGHT)
        d.rounded_rectangle((x + 122, y - 10, x + 282, y + 28), radius=18, fill=GREEN, outline=GOLD_LIGHT, width=2)
        draw_center(d, (x + 122, y - 9, 160, 36), title, pil_font(25, bold=True), IVORY_2)
        for i, b in enumerate(bullets):
            d.text((x + 108, y + 45 + i * 28), "• " + b, font=pil_font(17), fill=DARK)
        draw_layer_image(img, photo, (x + w - 148, y + 38, 120, 82), 10)

    d.rounded_rectangle((90, 970, 1830, 1046), radius=34, fill=GREEN, outline=GOLD_LIGHT, width=3)
    d.ellipse((130, 988, 178, 1036), fill=GOLD_LIGHT)
    draw_center(d, (190, 978, 1540, 58), "调研结论直接指向活动设计：用可参与、可体验、可传播的方式讲好本草文化", pil_font(34, bold=True), GOLD_LIGHT)
    img.convert("RGB").save(PREVIEW_OUT, quality=95)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(color[0], color[1], color[2])


def px(prs: Presentation, value: int | float, axis: str = "x") -> int:
    scale = prs.slide_width / CANVAS_W if axis == "x" else prs.slide_height / CANVAS_H
    return int(value * scale)


def set_name(shape, name: str) -> None:
    try:
        shape.name = name
    except Exception:
        pass


def add_shape(slide, prs, kind, name: str, box: tuple[int, int, int, int], fill: tuple[int, int, int], line: tuple[int, int, int] | None = None, width_pt: float = 1.0):
    x, y, w, h = box
    shape = slide.shapes.add_shape(kind, px(prs, x, "x"), px(prs, y, "y"), px(prs, w, "x"), px(prs, h, "y"))
    set_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width_pt)
    return shape


def add_picture(slide, prs, name: str, path: Path, box: tuple[int, int, int, int]):
    x, y, w, h = box
    pic = slide.shapes.add_picture(str(path), px(prs, x, "x"), px(prs, y, "y"), width=px(prs, w, "x"), height=px(prs, h, "y"))
    set_name(pic, name)
    return pic


def add_text(slide, prs, name: str, box: tuple[int, int, int, int], text: str, font_name: str, size_pt: float, color: tuple[int, int, int], *, bold: bool = False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    x, y, w, h = box
    shape = slide.shapes.add_textbox(px(prs, x, "x"), px(prs, y, "y"), px(prs, w, "x"), px(prs, h, "y"))
    set_name(shape, name)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_bullet(slide, prs, name: str, box: tuple[int, int, int, int], bullets: list[str], size_pt: float = 11.2):
    text = "\n".join("• " + b for b in bullets)
    return add_text(slide, prs, name, box, text, FONT_YAHEI, size_pt, DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_picture(slide, prs, "background-base-image2-muted", BG_OUT, (0, 0, CANVAS_W, CANVAS_H))

    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "top-title-backplate", (225, 22, 1440, 100), IVORY_2, None)
    add_text(slide, prs, "slide-title", (250, 34, 1340, 62), "趣味调研：青年愿意了解中药，但更需要体验式入口", FONT_KAI, 22.8, GREEN, bold=True)
    add_shape(slide, prs, MSO_SHAPE.RECTANGLE, "subtitle-line-left", (180, 137, 535, 3), GOLD)
    add_shape(slide, prs, MSO_SHAPE.RECTANGLE, "subtitle-line-right", (1205, 137, 535, 3), GOLD)
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "subtitle-backplate", (705, 111, 510, 52), IVORY_2, None)
    add_text(slide, prs, "subtitle", (720, 115, 480, 50), "问卷截图与结果图均来自项目材料", FONT_YAHEI, 14.2, DARK)

    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "survey-panel-frame", (20, 182, 928, 583), IVORY_2, GOLD, 1.0)
    add_picture(slide, prs, "survey-screenshots-layer", SURVEY_LAYER, (24, 186, 920, 575))
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "results-panel-frame", (959, 175, 898, 580), IVORY_2, GOLD, 1.0)
    add_picture(slide, prs, "results-charts-layer", RESULT_LAYER, (963, 179, 890, 572))

    card_specs = [
        ("sample", (30, 790, 548, 150), "青年样本", ["18–25岁青年为主力样本", "多数对中医药持积极态度", "认知基础有待提升"], PHOTO_01, "样"),
        ("pain", (610, 790, 505, 150), "认知痛点", ["概念混淆，专业门槛高", "辨识困难，实物接触少", "缺乏有趣、可参与的体验"], PHOTO_02, "痛"),
        ("activity", (1140, 790, 505, 150), "活动启发", ["以体验破解认知壁垒", "以互动提升参与热情", "以传播扩大文化影响"], PHOTO_03, "启"),
    ]
    for key, box, title, bullets, photo, icon in card_specs:
        x, y, w, h = box
        add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, f"{key}-card", box, IVORY_2, (218, 202, 160), 1.0)
        add_shape(slide, prs, MSO_SHAPE.OVAL, f"{key}-icon-circle", (x + 18, y + 34, 66, 66), GREEN, GOLD_LIGHT, 1.4)
        add_text(slide, prs, f"{key}-icon-text", (x + 18, y + 34, 66, 66), icon, FONT_YAHEI, 17, GOLD_LIGHT, bold=True)
        add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, f"{key}-label", (x + 122, y - 10, 160, 38), GREEN, GOLD_LIGHT, 1.0)
        add_text(slide, prs, f"{key}-label-text", (x + 122, y - 8, 160, 34), title, FONT_YAHEI, 13.2, IVORY_2, bold=True)
        add_bullet(slide, prs, f"{key}-bullets", (x + 104, y + 34, w - 252, 95), bullets, size_pt=8.8)
        add_picture(slide, prs, f"{key}-photo", photo, (x + w - 148, y + 38, 120, 82))

    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "conclusion-strip", (90, 970, 1740, 76), GREEN, GOLD_LIGHT, 1.6)
    add_shape(slide, prs, MSO_SHAPE.CLOUD, "conclusion-cloud-left", (126, 988, 58, 42), GOLD_LIGHT, GOLD, 0.8)
    add_shape(slide, prs, MSO_SHAPE.CLOUD, "conclusion-cloud-right", (1740, 988, 58, 42), GOLD_LIGHT, GOLD, 0.8)
    add_text(slide, prs, "conclusion-text", (190, 978, 1540, 58), "调研结论直接指向活动设计：用可参与、可体验、可传播的方式讲好本草文化", FONT_YAHEI, 19.5, GOLD_LIGHT, bold=True)

    prs.save(PPTX_OUT)


def rel(path: Path) -> str:
    return str(path.relative_to(ROOT)).replace("\\", "/")


def build_manifest() -> None:
    elements = [
        {"id": "background-base-image2-muted", "type": "image", "source": "generated_background_from_image2_reference", "file": rel(BG_OUT), "bbox": {"x": 0, "y": 0, "w": CANVAS_W, "h": CANVAS_H}, "z": 0, "editable": False, "animation_group": "background"},
        {"id": "slide-title", "type": "text", "source": "native", "text": "趣味调研：青年愿意了解中药，但更需要体验式入口", "font": FONT_KAI, "font_size": 22.8, "color": "RGB(16,82,54)", "align": "center", "bbox": {"x": 250, "y": 34, "w": 1340, "h": 62}, "z": 10, "editable": True},
        {"id": "subtitle", "type": "text", "source": "native", "text": "问卷截图与结果图均来自项目材料", "font": FONT_YAHEI, "font_size": 14.2, "color": "RGB(45,54,39)", "align": "center", "bbox": {"x": 720, "y": 115, "w": 480, "h": 50}, "z": 20, "editable": True},
        {"id": "survey-screenshots-layer", "type": "image", "source": "visual_reference_crop", "file": rel(SURVEY_LAYER), "bbox": {"x": 24, "y": 186, "w": 920, "h": 575}, "z": 30, "editable": False, "animation_group": "evidence"},
        {"id": "results-charts-layer", "type": "image", "source": "visual_reference_crop", "file": rel(RESULT_LAYER), "bbox": {"x": 963, "y": 179, "w": 890, "h": 572}, "z": 40, "editable": False, "animation_group": "evidence"},
        {"id": "sample-card", "type": "shape", "source": "native", "bbox": {"x": 30, "y": 790, "w": 548, "h": 150}, "z": 50, "editable": True},
        {"id": "sample-card-text", "type": "text", "source": "native", "text": "青年样本\n18–25岁青年为主力样本\n多数对中医药持积极态度\n认知基础有待提升", "font": FONT_YAHEI, "font_size": 8.8, "color": "RGB(45,54,39)", "align": "left", "bbox": {"x": 134, "y": 782, "w": 296, "h": 128}, "z": 60, "editable": True},
        {"id": "pain-card-text", "type": "text", "source": "native", "text": "认知痛点\n概念混淆，专业门槛高\n辨识困难，实物接触少\n缺乏有趣、可参与的体验", "font": FONT_YAHEI, "font_size": 8.8, "color": "RGB(45,54,39)", "align": "left", "bbox": {"x": 714, "y": 782, "w": 253, "h": 128}, "z": 70, "editable": True},
        {"id": "activity-card-text", "type": "text", "source": "native", "text": "活动启发\n以体验破解认知壁垒\n以互动提升参与热情\n以传播扩大文化影响", "font": FONT_YAHEI, "font_size": 8.8, "color": "RGB(45,54,39)", "align": "left", "bbox": {"x": 1244, "y": 782, "w": 253, "h": 128}, "z": 80, "editable": True},
        {"id": "conclusion-text", "type": "text", "source": "native", "text": "调研结论直接指向活动设计：用可参与、可体验、可传播的方式讲好本草文化", "font": FONT_YAHEI, "font_size": 19.5, "color": "RGB(242,220,150)", "align": "center", "bbox": {"x": 190, "y": 978, "w": 1540, "h": 58}, "z": 90, "editable": True},
    ]
    manifest = {
        "slide": 2,
        "canvas": {"w": CANVAS_W, "h": CANVAS_H},
        "mode": "premium image-first semantic layer merge",
        "reference_image": rel(SRC),
        "output_pptx": rel(PPTX_OUT),
        "preview": rel(PREVIEW_OUT),
        "powerpoint_preview": rel(POWERPOINT_PREVIEW),
        "elements": elements,
    }
    MANIFEST_OUT.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def build_strategy_doc() -> None:
    content = """# slide_02 分层策略

## 页面判断

- 页面类型：调研发现页。
- 处理目标：保留 image-2 形成的问卷/图表视觉密度，同时把主叙事文字和结论卡片改成可编辑对象。
- 拆分方式：语义拆分，不做 PSD 级自动抠图，不重造问卷数据或图表数据。

## 保留为背景/图片

- 原始页面图 `slide_02.png` 作为视觉参考。
- 由原图生成 `slide_02_background_base.png`：保留草药枝叶、山水纸纹、右下药材器具和整体氛围；对标题、主体面板、底部卡片、结论条区域做柔化压淡，避免与可编辑层重影。
- `slide_02_survey_screenshots_layer.png`：左侧问卷截图组合，作为独立图片对象保留，避免重造问卷题目与选项。
- `slide_02_results_charts_layer.png`：右侧饼图与条形图组合，作为独立图片对象保留，避免生成假比例或假数据。
- 三张底部小照片由原图裁切为独立图片对象，放入对应原生卡片。

## 原生 PPT 文字

- 页面标题：`趣味调研：青年愿意了解中药，但更需要体验式入口`。
- 副标题：`问卷截图与结果图均来自项目材料`。
- 底部三张结论卡：`青年样本`、`认知痛点`、`活动启发` 及各自三条要点。
- 底部结论条：`调研结论直接指向活动设计：用可参与、可体验、可传播的方式讲好本草文化`。

## 原生 PPT shape

- 标题下方分隔线。
- 左右证据面板边框。
- 底部三张圆角卡片、深绿色标签、圆形编号/主题标识。
- 底部深绿色结论条与金色装饰云纹占位。

## 不做的事

- 不生成假问卷、假比例、假数据、假 logo 或假二维码。
- 不把第 2 页做成仅一张整页图片。
- 不改动第 1 页或第 3-10 页。
"""
    STRATEGY_OUT.write_text(content, encoding="utf-8")


def validate_package_smoke() -> dict:
    result = {"pptx_exists": PPTX_OUT.exists(), "zip_ok": False, "slide_xml_count": 0, "media_count": 0, "contains_chinese": False, "shape_tags": 0, "picture_tags": 0}
    with zipfile.ZipFile(PPTX_OUT) as zf:
        result["zip_ok"] = zf.testzip() is None
        names = zf.namelist()
        slide_xml = [n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        result["slide_xml_count"] = len(slide_xml)
        result["media_count"] = len([n for n in names if n.startswith("ppt/media/")])
        xml = "\n".join(zf.read(n).decode("utf-8", errors="replace") for n in slide_xml)
        result["shape_tags"] = xml.count("<p:sp>")
        result["picture_tags"] = xml.count("<p:pic>")
        result["contains_chinese"] = all(s in xml for s in ["趣味调研", "体验式入口", "问卷截图与结果图", "青年样本", "调研结论直接指向活动设计"])
    return result


def build_validation_doc(smoke: dict) -> None:
    content = f"""# slide_02 分层 PPTX 验证记录

## 产物

- 单页 PPTX：`slide_02_layered.pptx`
- 背景母图：`slide_02_background_base.png`
- 问卷截图层：`slide_02_survey_screenshots_layer.png`
- 调研结果图表层：`slide_02_results_charts_layer.png`
- 脚本预览：`slide_02_layered_preview.png`
- PowerPoint 导出预览：`slide_02_layered_powerpoint_export.png`
- 分层 manifest：`slide_02_layer_manifest.json`
- 拆分策略：`slide_02_split_strategy.md`

## 技术验证

- PPTX ZIP CRC：{'通过' if smoke['zip_ok'] else '未通过'}。
- 幻灯片数量：{smoke['slide_xml_count']}。
- 媒体数量：{smoke['media_count']}。
- PowerPoint 原生 shape：{smoke['shape_tags']} 个。
- PowerPoint 图片对象：{smoke['picture_tags']} 个。
- 中文文本检查：{'通过' if smoke['contains_chinese'] else '未通过'}。
- PowerPoint 应用级验证：待运行导出后补充。
- UTF-8/mojibake 扫描：待统一扫描。
- `git diff --check`：待统一检查。

## 视觉复核

- 问卷截图和调研图表作为独立图片层保留，未重造数据。
- 标题、副标题、三张底部结论卡和底部结论条为原生 PPT 文本。
- 右上真实/生成 logo 区域不重绘，不生成假 logo。

## 已知取舍

- 问卷题目、图表分类和图形比例保留为图片层，不逐项转成原生图表，避免误造或误读数据。
- 复杂枝叶、山水、药材器具和纸纹氛围统一保留在背景母图。
"""
    VALIDATION_OUT.write_text(content, encoding="utf-8")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    bg = build_background_and_layers()
    make_preview(bg)
    build_pptx()
    build_manifest()
    build_strategy_doc()
    smoke = validate_package_smoke()
    build_validation_doc(smoke)
    print(json.dumps(smoke, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
