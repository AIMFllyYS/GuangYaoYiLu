from __future__ import annotations

import json
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
SRC = ROOT / "决赛PPT" / "输出终稿" / "image2_premium_ppt_v2" / "slides" / "slide_01.png"
OUT = ROOT / "决赛PPT" / "WorkShape" / "layered_pages"

CANVAS_W = 1920
CANVAS_H = 1080
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5

PPTX_OUT = OUT / "slide_01_layered.pptx"
BG_OUT = OUT / "slide_01_background_base.png"
PREVIEW_OUT = OUT / "slide_01_layered_preview.png"
MANIFEST_OUT = OUT / "slide_01_layer_manifest.json"
STRATEGY_OUT = OUT / "slide_01_split_strategy.md"

FONT_YAHEI = "Microsoft YaHei"
FONT_KAI = "KaiTi"

FONT_REG = Path("C:/Windows/Fonts/msyh.ttc")
FONT_BOLD = Path("C:/Windows/Fonts/msyhbd.ttc")
FONT_KAI_FILE = Path("C:/Windows/Fonts/simkai.ttf")

GREEN = (12, 71, 46)
GREEN_2 = (20, 91, 60)
GOLD = (211, 171, 84)
GOLD_LIGHT = (246, 222, 143)
IVORY = (246, 242, 225)
IVORY_2 = (255, 251, 235)
INK = (31, 45, 34)
SHADOW = (36, 32, 22)


def pil_font(size: int, *, bold: bool = False, kai: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_KAI_FILE if kai else (FONT_BOLD if bold else FONT_REG)
    return ImageFont.truetype(str(path), size)


def fit_slide_image(path: Path) -> Image.Image:
    img = ImageOps.exif_transpose(Image.open(path)).convert("RGBA")
    return ImageOps.fit(img, (CANVAS_W, CANVAS_H), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))


def rounded_mask(size: tuple[int, int], radius: int) -> Image.Image:
    mask = Image.new("L", size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle((0, 0, size[0], size[1]), radius=radius, fill=255)
    return mask


def soften_region(base: Image.Image, box: tuple[int, int, int, int], fill: tuple[int, int, int], alpha: int, radius: int = 32) -> None:
    x, y, w, h = box
    crop = base.crop((x, y, x + w, y + h)).filter(ImageFilter.GaussianBlur(18))
    base.alpha_composite(crop, (x, y))
    overlay = Image.new("RGBA", (w, h), fill + (alpha,))
    base.paste(overlay, (x, y), rounded_mask((w, h), radius))


def build_background() -> Image.Image:
    bg = fit_slide_image(SRC)
    soften_region(bg, (145, 135, 1630, 248), IVORY, 190, 38)
    soften_region(bg, (1515, 20, 370, 165), IVORY, 105, 70)
    soften_region(bg, (500, 410, 930, 245), GREEN, 112, 48)
    soften_region(bg, (365, 672, 1190, 132), IVORY, 170, 36)
    soften_region(bg, (650, 800, 650, 108), IVORY, 170, 34)
    bg.save(BG_OUT)
    return bg


def draw_center(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], text: str, font: ImageFont.FreeTypeFont, fill: tuple[int, int, int], spacing: int = 8) -> None:
    x, y, w, h = box
    bbox = draw.multiline_textbbox((0, 0), text, font=font, spacing=spacing, align="center")
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    draw.multiline_text((x + (w - tw) / 2, y + (h - th) / 2 - bbox[1]), text, font=font, fill=fill, spacing=spacing, align="center")


def make_preview(bg: Image.Image) -> None:
    img = bg.copy()
    d = ImageDraw.Draw(img)

    d.rounded_rectangle((164, 157, 1760, 315), radius=38, outline=GOLD + (210,), width=2)
    slogan_font = pil_font(68, bold=True, kai=True)
    draw_center(d, (205, 175, 1510, 130), "金光青黛同济兴，百草杏林华夏清", slogan_font, SHADOW)
    draw_center(d, (198, 168, 1510, 130), "金光青黛同济兴，百草杏林华夏清", slogan_font, GREEN)

    d.line((510, 390, 875, 390), fill=GOLD, width=3)
    d.line((1045, 390, 1410, 390), fill=GOLD, width=3)
    d.polygon((956, 378, 973, 390, 956, 402, 939, 390), fill=GOLD)
    d.ellipse((903, 382, 919, 398), outline=GOLD, width=3)
    d.ellipse((1001, 382, 1017, 398), outline=GOLD, width=3)

    d.rounded_rectangle((500, 430, 1420, 626), radius=50, fill=(190, 151, 70), outline=(104, 81, 36), width=3)
    d.rounded_rectangle((516, 446, 1404, 608), radius=42, fill=GREEN_2, outline=GOLD_LIGHT, width=3)
    title_font = pil_font(114, bold=True, kai=True)
    draw_center(d, (520, 430, 880, 184), "光药医路", title_font, SHADOW)
    draw_center(d, (512, 420, 880, 184), "光药医路", title_font, GOLD_LIGHT)

    d.rounded_rectangle((395, 705, 1530, 783), radius=32, fill=GREEN_2, outline=GOLD_LIGHT, width=3)
    d.ellipse((418, 690, 510, 782), fill=GREEN, outline=GOLD_LIGHT, width=3)
    draw_center(d, (418, 690, 92, 92), "药", pil_font(42, bold=True), GOLD_LIGHT)
    team = "药学（中外）2503、光电2506、基础医学（强基）2501班团支部"
    draw_center(d, (525, 708, 970, 70), team, pil_font(32, bold=True), IVORY_2)

    d.rounded_rectangle((700, 820, 1220, 892), radius=32, fill=GREEN_2, outline=GOLD_LIGHT, width=3)
    d.ellipse((688, 802, 782, 896), fill=GREEN, outline=GOLD_LIGHT, width=3)
    draw_center(d, (688, 802, 94, 94), "日", pil_font(42, bold=True), GOLD_LIGHT)
    draw_center(d, (788, 824, 400, 60), "答辩时间：2026.5.24", pil_font(35, bold=True), IVORY_2)

    for x in (1572, 1736):
        d.ellipse((x, 34, x + 126, 160), fill=GREEN, outline=GOLD_LIGHT, width=5)
        d.ellipse((x + 8, 42, x + 118, 152), outline=(29, 38, 25), width=2)
        draw_center(d, (x + 18, 55, 90, 80), "校徽\n放置处", pil_font(31, bold=True), IVORY_2, spacing=2)

    img.convert("RGB").save(PREVIEW_OUT, quality=95)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(color[0], color[1], color[2])


def px(prs: Presentation, value: int | float, axis: str = "x") -> int:
    scale = prs.slide_width / CANVAS_W if axis == "x" else prs.slide_height / CANVAS_H
    return int(value * scale)


def set_name(shape, name: str) -> None:
    try:
        shape.name = name
    except Exception:
        pass


def add_shape(slide, prs, kind, name: str, box: tuple[int, int, int, int], fill: tuple[int, int, int] | None, line: tuple[int, int, int] | None = None, width_pt: float = 1.0):
    x, y, w, h = box
    shape = slide.shapes.add_shape(kind, px(prs, x, "x"), px(prs, y, "y"), px(prs, w, "x"), px(prs, h, "y"))
    set_name(shape, name)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width_pt)
    return shape


def add_line(slide, prs, name: str, start: tuple[int, int], end: tuple[int, int], color: tuple[int, int, int], width_pt: float = 1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px(prs, start[0], "x"),
        px(prs, start[1], "y"),
        px(prs, end[0], "x"),
        px(prs, end[1], "y"),
    )
    set_name(line, name)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width_pt)
    return line


def add_text(slide, prs, name: str, box: tuple[int, int, int, int], text: str, font_name: str, size_pt: float, color: tuple[int, int, int], *, bold: bool = False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    x, y, w, h = box
    shape = slide.shapes.add_textbox(px(prs, x, "x"), px(prs, y, "y"), px(prs, w, "x"), px(prs, h, "y"))
    set_name(shape, name)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    lines = text.split("\n")
    for i, line_text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line_text
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_text_shadow(slide, prs, base_name: str, box: tuple[int, int, int, int], text: str, font_name: str, size_pt: float, front: tuple[int, int, int], shadow: tuple[int, int, int], *, bold: bool = True, dx: int = 5, dy: int = 6):
    x, y, w, h = box
    add_text(slide, prs, base_name + "-shadow", (x + dx, y + dy, w, h), text, font_name, size_pt, shadow, bold=bold)
    return add_text(slide, prs, base_name, box, text, font_name, size_pt, front, bold=bold)


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_picture(str(BG_OUT), 0, 0, width=prs.slide_width, height=prs.slide_height)
    set_name(bg, "background-base-image2-muted")

    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "top-slogan-backplate", (164, 157, 1596, 158), None, GOLD, 1.4)
    add_text_shadow(slide, prs, "top-slogan", (198, 168, 1510, 130), "金光青黛同济兴，百草杏林华夏清", FONT_KAI, 38, GREEN, SHADOW, dx=3, dy=4)

    add_line(slide, prs, "divider-left", (510, 390), (875, 390), GOLD, 1.5)
    add_line(slide, prs, "divider-right", (1045, 390), (1410, 390), GOLD, 1.5)
    add_shape(slide, prs, MSO_SHAPE.DIAMOND, "divider-center-diamond", (939, 378, 34, 24), GOLD, None)
    add_shape(slide, prs, MSO_SHAPE.OVAL, "divider-dot-left", (903, 382, 16, 16), IVORY, GOLD, 1.4)
    add_shape(slide, prs, MSO_SHAPE.OVAL, "divider-dot-right", (1001, 382, 16, 16), IVORY, GOLD, 1.4)

    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "main-title-plaque-outer", (500, 430, 920, 196), (190, 151, 70), (104, 81, 36), 1.7)
    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "main-title-plaque-inner", (516, 446, 888, 162), GREEN_2, GOLD_LIGHT, 1.8)
    add_text_shadow(slide, prs, "slide-title", (520, 432, 880, 174), "光药医路", FONT_KAI, 62, GOLD_LIGHT, SHADOW, dx=3, dy=4)

    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "team-strip", (395, 705, 1135, 78), GREEN_2, GOLD_LIGHT, 1.7)
    add_shape(slide, prs, MSO_SHAPE.OVAL, "team-icon-circle", (418, 690, 92, 92), GREEN, GOLD_LIGHT, 1.8)
    add_text(slide, prs, "team-icon-text", (418, 690, 92, 92), "药", FONT_YAHEI, 24, GOLD_LIGHT, bold=True)
    add_text(slide, prs, "team-text", (525, 708, 970, 70), "药学（中外）2503、光电2506、基础医学（强基）2501班团支部", FONT_YAHEI, 17, IVORY_2, bold=True)

    add_shape(slide, prs, MSO_SHAPE.ROUNDED_RECTANGLE, "date-strip", (700, 820, 520, 72), GREEN_2, GOLD_LIGHT, 1.7)
    add_shape(slide, prs, MSO_SHAPE.OVAL, "date-icon-circle", (688, 802, 94, 94), GREEN, GOLD_LIGHT, 1.8)
    add_text(slide, prs, "date-icon-text", (688, 802, 94, 94), "日", FONT_YAHEI, 24, GOLD_LIGHT, bold=True)
    add_text(slide, prs, "date-text", (788, 824, 400, 60), "答辩时间：2026.5.24", FONT_YAHEI, 20, IVORY_2, bold=True)

    for i, x in enumerate((1572, 1736), start=1):
        add_shape(slide, prs, MSO_SHAPE.OVAL, f"logo-placeholder-{i}-outer", (x, 34, 126, 126), GREEN, GOLD_LIGHT, 2.3)
        add_shape(slide, prs, MSO_SHAPE.OVAL, f"logo-placeholder-{i}-inner-line", (x + 8, 42, 110, 110), GREEN, INK, 0.8)
        add_text(slide, prs, f"logo-placeholder-{i}-text", (x + 18, 55, 90, 80), "校徽\n放置处", FONT_YAHEI, 15, IVORY_2, bold=True)

    prs.save(PPTX_OUT)


def build_manifest() -> None:
    elements = [
        {"id": "background-base-image2-muted", "type": "image", "source": "generated_background_from_image2_reference", "file": str(BG_OUT.relative_to(ROOT)).replace("\\", "/"), "bbox": {"x": 0, "y": 0, "w": CANVAS_W, "h": CANVAS_H}, "z": 0, "editable": False, "animation_group": "background"},
        {"id": "top-slogan-backplate", "type": "shape", "source": "native", "bbox": {"x": 164, "y": 157, "w": 1596, "h": 158}, "z": 10, "editable": True},
        {"id": "top-slogan", "type": "text", "source": "native", "text": "金光青黛同济兴，百草杏林华夏清", "font": FONT_KAI, "font_size": 38, "color": "RGB(12,71,46)", "align": "center", "bbox": {"x": 198, "y": 168, "w": 1510, "h": 130}, "z": 20, "editable": True, "animation_group": "text"},
        {"id": "divider", "type": "shape", "source": "native", "bbox": {"x": 510, "y": 378, "w": 900, "h": 40}, "z": 30, "editable": True},
        {"id": "main-title-plaque", "type": "shape", "source": "native", "bbox": {"x": 500, "y": 430, "w": 920, "h": 196}, "z": 40, "editable": True},
        {"id": "slide-title", "type": "text", "source": "native", "text": "光药医路", "font": FONT_KAI, "font_size": 62, "color": "RGB(246,222,143)", "align": "center", "bbox": {"x": 520, "y": 432, "w": 880, "h": 174}, "z": 50, "editable": True, "animation_group": "title"},
        {"id": "team-strip", "type": "shape", "source": "native", "bbox": {"x": 395, "y": 705, "w": 1135, "h": 78}, "z": 60, "editable": True},
        {"id": "team-text", "type": "text", "source": "native", "text": "药学（中外）2503、光电2506、基础医学（强基）2501班团支部", "font": FONT_YAHEI, "font_size": 17, "color": "RGB(255,251,235)", "align": "center", "bbox": {"x": 525, "y": 708, "w": 970, "h": 70}, "z": 70, "editable": True, "animation_group": "text"},
        {"id": "date-strip", "type": "shape", "source": "native", "bbox": {"x": 700, "y": 820, "w": 520, "h": 72}, "z": 80, "editable": True},
        {"id": "date-text", "type": "text", "source": "native", "text": "答辩时间：2026.5.24", "font": FONT_YAHEI, "font_size": 20, "color": "RGB(255,251,235)", "align": "center", "bbox": {"x": 788, "y": 824, "w": 400, "h": 60}, "z": 90, "editable": True, "animation_group": "text"},
        {"id": "logo-placeholder-1", "type": "shape", "source": "native", "bbox": {"x": 1572, "y": 34, "w": 126, "h": 126}, "z": 100, "editable": True},
        {"id": "logo-placeholder-2", "type": "shape", "source": "native", "bbox": {"x": 1736, "y": 34, "w": 126, "h": 126}, "z": 110, "editable": True},
    ]
    manifest = {
        "slide": 1,
        "canvas": {"w": CANVAS_W, "h": CANVAS_H},
        "mode": "premium image-first semantic layer merge",
        "reference_image": str(SRC.relative_to(ROOT)).replace("\\", "/"),
        "output_pptx": str(PPTX_OUT.relative_to(ROOT)).replace("\\", "/"),
        "preview": str(PREVIEW_OUT.relative_to(ROOT)).replace("\\", "/"),
        "elements": elements,
    }
    MANIFEST_OUT.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def build_strategy_doc() -> None:
    content = """# slide_01 分层策略

## 页面判断

- 页面类型：封面页。
- 处理目标：保留 image-2 的完整高级氛围，但将可编辑信息层从底图中分离出来。
- 拆分方式：语义拆分，不做 PSD 级自动抠图。

## 保留为背景/图片

- 原始页面图 `slide_01.png` 作为视觉参考。
- 由原图生成 `slide_01_background_base.png`：保留山水、校园/活动照片拼贴、中药器具、光效、草药枝叶、卷轴与桌面氛围；对标语、主标题、信息条、日期条和校徽占位区域做柔化压淡，避免原图烧录文字与原生文字严重重影。
- 左下药柜、研钵、药材盘、右下卷轴与笔记本区域保留为背景氛围，不强行拆成可编辑形状。

## 原生 PPT 文字

- 顶部标语：`金光青黛同济兴，百草杏林华夏清`。
- 主标题：`光药医路`。
- 团队信息：`药学（中外）2503、光电2506、基础医学（强基）2501班团支部`。
- 答辩时间：`答辩时间：2026.5.24`。
- 双校徽占位：`校徽 / 放置处`，仅作为占位，不生成假校徽。

## 原生 PPT shape

- 顶部标语浅金背板。
- 金色分隔线、中心菱形与小圆点。
- 深绿色主标题牌及金色边框。
- 团队信息条、日期信息条。
- 圆形药字图标、日期图标、两个校徽占位圆章。

## 不做的事

- 不生成假 logo、假二维码、假政策文件正文。
- 不把第 1 页做成仅一张整页图片。
- 不改动后续第 2-10 页。
"""
    STRATEGY_OUT.write_text(content, encoding="utf-8")


def validate_package_smoke() -> dict:
    result = {"pptx_exists": PPTX_OUT.exists(), "zip_ok": False, "slide_xml_count": 0, "media_count": 0, "contains_chinese": False}
    with zipfile.ZipFile(PPTX_OUT) as zf:
        bad = zf.testzip()
        result["zip_ok"] = bad is None
        names = zf.namelist()
        result["slide_xml_count"] = len([n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")])
        result["media_count"] = len([n for n in names if n.startswith("ppt/media/")])
        xml = "\n".join(zf.read(n).decode("utf-8", errors="replace") for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
        result["contains_chinese"] = all(s in xml for s in ["光药医路", "金光青黛同济兴", "答辩时间"])
    return result


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    bg = build_background()
    make_preview(bg)
    build_pptx()
    build_manifest()
    build_strategy_doc()
    smoke = validate_package_smoke()
    print(json.dumps(smoke, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
