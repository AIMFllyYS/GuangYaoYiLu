#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt


DEFAULT_CANVAS_W = 1920
DEFAULT_CANVAS_H = 1080
DEFAULT_PPT_W_EMU = 9144000
DEFAULT_PPT_H_EMU = 5143500
FONT_DIR = Path("C:/Windows/Fonts")
FONT_FALLBACKS = [
    "msyh.ttc",
    "simhei.ttf",
    "simsun.ttc",
    "arial.ttf",
]


def write_json(data: dict[str, Any]) -> None:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
    print(json.dumps(data, ensure_ascii=False, indent=2))


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def slide_sort_key(path: Path) -> tuple[int, str]:
    try:
        data = load_json(path)
        return int(data.get("slide", 0)), path.name
    except Exception:
        match = re.search(r"(\d+)", path.stem)
        return int(match.group(1)) if match else 0, path.name


def parse_rgb(value: str | None, default: tuple[int, int, int] = (0, 0, 0)) -> tuple[int, int, int]:
    if not value:
        return default
    raw = value.strip()
    if raw.startswith("#") and len(raw) == 7:
        return int(raw[1:3], 16), int(raw[3:5], 16), int(raw[5:7], 16)
    match = re.fullmatch(r"RGB\((\d+),\s*(\d+),\s*(\d+)\)", raw)
    if match:
        return tuple(int(match.group(i)) for i in range(1, 4))  # type: ignore[return-value]
    raise ValueError(f"Unsupported color value: {value}")


def align_value(value: str | None) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get((value or "left").lower(), PP_ALIGN.LEFT)


def resolve_path(workspace: Path, manifest_path: Path, value: str) -> Path:
    path = Path(value)
    if path.is_absolute():
        return path
    workspace_candidate = (workspace / path).resolve()
    if workspace_candidate.exists():
        return workspace_candidate
    manifest_candidate = (manifest_path.parent / path).resolve()
    return manifest_candidate


def px_to_emu_x(px: int | float, canvas_w: int, ppt_w_emu: int) -> int:
    return int(float(px) * ppt_w_emu / canvas_w)


def px_to_emu_y(px: int | float, canvas_h: int, ppt_h_emu: int) -> int:
    return int(float(px) * ppt_h_emu / canvas_h)


def local_font_path(font_name: str | None = None) -> str | None:
    candidates: list[str] = []
    if font_name:
        candidates.extend([font_name, f"{font_name}.ttf", f"{font_name}.ttc"])
    candidates.extend(FONT_FALLBACKS)
    for candidate in candidates:
        path = Path(candidate)
        if path.is_absolute() and path.exists():
            return str(path)
        font_path = FONT_DIR / candidate
        if font_path.exists():
            return str(font_path)
    return None


def preview_font(font_name: str | None, size_px: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    path = local_font_path(font_name)
    if path:
        return ImageFont.truetype(path, size_px)
    return ImageFont.load_default()


def set_east_asian_font(run: Any, font_name: str) -> None:
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()  # noqa: SLF001 - python-pptx has no public East Asian font setter.
    r_pr.set(qn("a:ea"), font_name)


def add_text(slide: Any, element: dict[str, Any], canvas: dict[str, int], ppt_size: tuple[int, int]) -> None:
    bbox = element["bbox"]
    ppt_w, ppt_h = ppt_size
    box = slide.shapes.add_textbox(
        px_to_emu_x(bbox["x"], canvas["w"], ppt_w),
        px_to_emu_y(bbox["y"], canvas["h"], ppt_h),
        px_to_emu_x(bbox["w"], canvas["w"], ppt_w),
        px_to_emu_y(bbox["h"], canvas["h"], ppt_h),
    )
    box.name = element.get("id", "text")
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align_value(element.get("align"))
    run = paragraph.add_run()
    run.text = str(element.get("text", ""))
    font_name = str(element.get("font", "Microsoft YaHei"))
    set_east_asian_font(run, font_name)
    run.font.size = Pt(float(element.get("font_size", 24)))
    run.font.color.rgb = RGBColor(*parse_rgb(element.get("color"), (0, 0, 0)))


def add_image(
    slide: Any,
    workspace: Path,
    manifest_path: Path,
    element: dict[str, Any],
    canvas: dict[str, int],
    ppt_size: tuple[int, int],
) -> Path:
    image_path = resolve_path(workspace, manifest_path, element["file"])
    if not image_path.exists():
        raise FileNotFoundError(image_path)
    bbox = element["bbox"]
    ppt_w, ppt_h = ppt_size
    pic = slide.shapes.add_picture(
        str(image_path),
        px_to_emu_x(bbox["x"], canvas["w"], ppt_w),
        px_to_emu_y(bbox["y"], canvas["h"], ppt_h),
        px_to_emu_x(bbox["w"], canvas["w"], ppt_w),
        px_to_emu_y(bbox["h"], canvas["h"], ppt_h),
    )
    pic.name = element.get("id", image_path.stem)
    return image_path


def add_shape(slide: Any, element: dict[str, Any], canvas: dict[str, int], ppt_size: tuple[int, int]) -> None:
    bbox = element["bbox"]
    ppt_w, ppt_h = ppt_size
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        px_to_emu_x(bbox["x"], canvas["w"], ppt_w),
        px_to_emu_y(bbox["y"], canvas["h"], ppt_h),
        px_to_emu_x(bbox["w"], canvas["w"], ppt_w),
        px_to_emu_y(bbox["h"], canvas["h"], ppt_h),
    )
    shape.name = element.get("id", "shape")
    fill = element.get("fill")
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*parse_rgb(fill, (255, 255, 255)))
    else:
        shape.fill.background()
    stroke = element.get("stroke")
    if stroke:
        shape.line.color.rgb = RGBColor(*parse_rgb(stroke, (0, 0, 0)))
    else:
        shape.line.fill.background()


def draw_preview_text(draw: ImageDraw.ImageDraw, element: dict[str, Any]) -> None:
    bbox = element["bbox"]
    size_px = max(8, int(float(element.get("font_size", 24)) * 1.33))
    font = preview_font(element.get("font"), size_px)
    text = str(element.get("text", ""))
    color = parse_rgb(element.get("color"), (0, 0, 0))
    text_box = draw.textbbox((0, 0), text, font=font)
    text_w = text_box[2] - text_box[0]
    text_h = text_box[3] - text_box[1]
    align = str(element.get("align", "left")).lower()
    if align == "center":
        x = bbox["x"] + (bbox["w"] - text_w) / 2
    elif align == "right":
        x = bbox["x"] + bbox["w"] - text_w
    else:
        x = bbox["x"]
    y = bbox["y"] + (bbox["h"] - text_h) / 2
    draw.text((x, y), text, fill=color, font=font)


def build_preview(
    workspace: Path,
    manifest_path: Path,
    manifest: dict[str, Any],
    preview_path: Path,
) -> None:
    canvas = normalize_canvas(manifest)
    preview = Image.new("RGBA", (canvas["w"], canvas["h"]), (255, 255, 255, 255))
    draw = ImageDraw.Draw(preview)
    for element in sorted(manifest["elements"], key=lambda item: item.get("z", 0)):
        element_type = element["type"]
        bbox = element["bbox"]
        if element_type == "image":
            image_path = resolve_path(workspace, manifest_path, element["file"])
            image = Image.open(image_path).convert("RGBA").resize((bbox["w"], bbox["h"]), Image.Resampling.LANCZOS)
            preview.alpha_composite(image, (bbox["x"], bbox["y"]))
        elif element_type == "text":
            draw_preview_text(draw, element)
        elif element_type == "shape":
            fill = parse_rgb(element.get("fill"), (255, 255, 255))
            stroke = parse_rgb(element.get("stroke"), fill)
            draw.rectangle((bbox["x"], bbox["y"], bbox["x"] + bbox["w"], bbox["y"] + bbox["h"]), fill=fill, outline=stroke)
    preview_path.parent.mkdir(parents=True, exist_ok=True)
    preview.convert("RGB").save(preview_path)


def normalize_canvas(manifest: dict[str, Any]) -> dict[str, int]:
    canvas = manifest.get("canvas") or {}
    return {
        "w": int(canvas.get("w", DEFAULT_CANVAS_W)),
        "h": int(canvas.get("h", DEFAULT_CANVAS_H)),
    }


def validate_manifest(manifest: dict[str, Any], manifest_path: Path) -> None:
    for key in ["slide", "elements"]:
        if key not in manifest:
            raise ValueError(f"{manifest_path}: missing required top-level field {key}")
    canvas = normalize_canvas(manifest)
    if canvas["w"] <= 0 or canvas["h"] <= 0:
        raise ValueError(f"{manifest_path}: invalid canvas")
    for element in manifest["elements"]:
        for key in ["id", "type", "bbox", "z"]:
            if key not in element:
                raise ValueError(f"{manifest_path}: element missing required field {key}")
        bbox = element["bbox"]
        for key in ["x", "y", "w", "h"]:
            if key not in bbox:
                raise ValueError(f"{manifest_path}: element {element['id']} bbox missing {key}")
        if element["type"] == "image" and "file" not in element:
            raise ValueError(f"{manifest_path}: image element {element['id']} missing file")
        if element["type"] == "text":
            for key in ["text", "font", "font_size", "color", "align"]:
                if key not in element:
                    raise ValueError(f"{manifest_path}: text element {element['id']} missing {key}")


def build_deck(args: argparse.Namespace) -> dict[str, Any]:
    workspace = Path(args.workspace).resolve()
    manifest_dir = Path(args.manifest_dir)
    if not manifest_dir.is_absolute():
        manifest_dir = workspace / manifest_dir
    manifest_paths = sorted(manifest_dir.glob("*.json"), key=slide_sort_key)
    if args.expected_slides is not None and len(manifest_paths) != args.expected_slides:
        raise ValueError(f"Expected {args.expected_slides} manifests, found {len(manifest_paths)}")
    if not manifest_paths:
        raise ValueError(f"No manifest JSON files found in {manifest_dir}")

    prs = Presentation()
    prs.slide_width = Emu(DEFAULT_PPT_W_EMU)
    prs.slide_height = Emu(DEFAULT_PPT_H_EMU)
    preview_dir = Path(args.preview_dir)
    if not preview_dir.is_absolute():
        preview_dir = workspace / preview_dir
    output_path = Path(args.out)
    if not output_path.is_absolute():
        output_path = workspace / output_path
    output_path.parent.mkdir(parents=True, exist_ok=True)

    built_slides = []
    for manifest_path in manifest_paths:
        manifest = load_json(manifest_path)
        validate_manifest(manifest, manifest_path)
        canvas = normalize_canvas(manifest)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for element in sorted(manifest["elements"], key=lambda item: item.get("z", 0)):
            if element["type"] == "image":
                add_image(slide, workspace, manifest_path, element, canvas, (DEFAULT_PPT_W_EMU, DEFAULT_PPT_H_EMU))
            elif element["type"] == "text":
                add_text(slide, element, canvas, (DEFAULT_PPT_W_EMU, DEFAULT_PPT_H_EMU))
            elif element["type"] == "shape":
                add_shape(slide, element, canvas, (DEFAULT_PPT_W_EMU, DEFAULT_PPT_H_EMU))
            else:
                raise ValueError(f"{manifest_path}: unsupported element type {element['type']}")

        preview_path = preview_dir / f"slide_{int(manifest['slide']):02d}_preview.png"
        build_preview(workspace, manifest_path, manifest, preview_path)
        built_slides.append({"slide": manifest["slide"], "manifest": str(manifest_path), "preview": str(preview_path)})

    prs.save(output_path)
    return {"pptx": str(output_path), "slides": built_slides, "slide_count": len(built_slides)}


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a PPTX deck and previews from slide manifest JSON files.")
    parser.add_argument("workspace", help="PPTX production workspace root")
    parser.add_argument("--manifest-dir", default="03_assembly/manifests", help="Manifest directory, absolute or relative to workspace")
    parser.add_argument("--out", default="04_final/pptx/deck.pptx", help="Output PPTX path, absolute or relative to workspace")
    parser.add_argument("--preview-dir", default="03_assembly/previews", help="Preview output directory, absolute or relative to workspace")
    parser.add_argument("--expected-slides", type=int, help="Expected manifest/slide count")
    args = parser.parse_args()
    try:
        write_json(build_deck(args))
    except Exception as exc:
        write_json({"ok": False, "error": str(exc)})
        raise SystemExit(1)


if __name__ == "__main__":
    main()
